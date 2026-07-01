#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""lint_deck.py — deterministic LAYOUT lint for a built .pptx: catches the overlaps a visual
critic can miss (a callout tucked a few px under a panel/image, a block hanging off the slide,
content colliding with the footer).

    python scripts/lint_deck.py deck.pptx

Run it right after building, BEFORE/with the critic loop — it's a cheap, deterministic safety net
for the "no block/text/image overlap" rule, not a replacement for the visual critic (which judges
crop, balance, legibility, fidelity). Exits non-zero if it finds anything.

Checks (tuned for low false-positives):
  1. OVERFLOW   — a shape extends beyond the slide edge.
  2. OVERLAP    — two SOLID shapes (filled auto-shape / picture / table / freeform) partially
                  overlap with NEITHER contained in the other. Intentional layering — text on a
                  card, an image inside its frame, a header band on a card, a full-bleed
                  background — is *containment* and is NOT flagged; only "two separate blocks
                  collide" partial overlaps are.
  3. FOOTER     — the bottom footer band is reserved deck chrome: NO content block (solid or text)
                  may cover the footer/page-number text or dip into its zone. Unlike OVERLAP this has
                  NO containment escape — a low content band that swallows the footer text into its
                  bbox is still a collision (that "text-on-a-card" exclusion was the blind spot that
                  let a band overlap the footer).
  Plus, measured against each text's RENDERED extent (alignment- + anchor-aware, so centred / middle
  text isn't mis-flagged) — the recurring "looks fine to a box-overlap check but cramped to the eye" bugs:
  4. TEXT PADDING   — text inside a filled card crammed against (or past) the card's BOTTOM edge with
                      less than the minimum interior margin (not just hard overflow).
  5. CHIP/LABEL TOO SMALL — a text box coincident with a small filled pill whose label is taller/wider
                      than the pill (a row of chips where the text overruns the chip — size the chip to its text).
  6. TEXT COLLISION — the rendered bottom of an upper text box overruns the rendered top of a lower one
                      in the same column (a wrap-collision two *declared* boxes never reveal).
  Plus: off-slide overflow, TEXT-OVERFLOWS-CARD, UNEVEN CARD HEIGHTS, ORPHANED PUNCTUATION / WIDOW
  (a wrapped box whose last line is a lone 。/，or a single CJK glyph — 避头尾), CJK-TEXT-WITHOUT-EA-FONT
  (the root cause of orphaned punctuation: no <a:ea> → PowerPoint applies no kinsoku, plus tofu risk —
  checked across ANY text: text boxes, table cells, and grouped shapes), whole-page-image/editability,
  and orphan/empty slides.
  Plus soft [warn]s (advisory — printed but do NOT fail the exit code): MISSING ALT-TEXT on an
  informative image (accessibility; invisible to every other check), and MATH-FONT TOFU (an
  equation_native math font not installed on the render host → equations render as boxes).
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

EMU = 914400.0
TOL = 0.05        # inches — ignore hairline/touching overlaps
CONTAIN = 0.90    # A is "inside" B when >=90% of A's area lies within B
SOLID = {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.FREEFORM}
MATH_FONTS = {"STIX Two Math", "Cambria Math", "Latin Modern Math", "XITS Math", "Asana Math", "TeX Gyre Termes Math"}
try:                                                  # reuse deckkit's font-availability probe for the tofu warn
    import os as _os
    sys.path.insert(0, _os.path.dirname(_os.path.abspath(__file__)))
    from deckkit import _font_substituted as _fsub
except Exception:
    def _fsub(_name):                                 # can't check → never warn (no false positive)
        return False

def _no_real_alt(descr):
    # python-pptx auto-sets a picture's descr to its FILE NAME; treat that (or None) as no real
    # alt-text. A deliberate alt="" (decorative) or a real sentence is fine and NOT flagged.
    if descr is None:
        return True
    d = descr.strip().lower()
    return d.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"))


def _boxes(slide, sw, sh):
    out = []
    for s in slide.shapes:
        try:
            l, t, w, h = s.left / EMU, s.top / EMU, s.width / EMU, s.height / EMU
        except (TypeError, AttributeError):
            continue
        if not w or not h or w <= 0 or h <= 0:
            continue
        full = s.text_frame.text.strip() if s.has_text_frame else ""
        txt = full.replace("\n", " ")[:26]
        paras, size, align, anchor, mathfont = [], 0.0, None, None, None
        if s.has_text_frame:
            size = max((r.font.size.pt for p in s.text_frame.paragraphs for r in p.runs if r.font.size),
                       default=12.0)
            for p in s.text_frame.paragraphs:
                pr = []
                for r in p.runs:
                    pr.append((r.text, (r.font.size.pt if r.font.size else size)))
                    if mathfont is None and r.font.name in MATH_FONTS:
                        mathfont = r.font.name           # an equation_native run in a math font
                if pr:
                    paras.append(pr)
            try: anchor = s.text_frame.vertical_anchor
            except Exception: anchor = None
            try: align = s.text_frame.paragraphs[0].alignment if s.text_frame.paragraphs else None
            except Exception: align = None
        descr = None                                     # accessibility alt-text (cNvPr 'descr')
        try:
            cn = s._element.find(".//" + qn("p:cNvPr"))
            if cn is not None:
                descr = cn.get("descr")
        except Exception:
            descr = None
        out.append({"l": l, "t": t, "w": w, "h": h, "r": l + w, "b": t + h,
                    "st": str(s.shape_type).split()[0], "txt": txt, "full": full, "size": size or 12.0,
                    "paras": paras, "solid": s.shape_type in SOLID, "align": align, "anchor": anchor,
                    "text": bool(s.has_text_frame and txt), "descr": descr, "mathfont": mathfont,
                    "bg": (w * h) >= 0.95 * (sw * sh)})
    return out


def _est_lines(paras, width_in):
    """CJK-aware estimate of total wrapped line count across paragraphs, using each RUN's own font
    size (so a mixed-size 'small label + big value' stat line counts as one line, not two). CJK glyph
    ≈ 1 em, Latin ≈ 0.52 em."""
    if width_in <= 0 or not paras:
        return 1
    total = 0
    for pr in paras:
        lines, w = 1, 0.0
        for text, size_pt in pr:
            em = max(0.01, size_pt / 72.0)
            for ch in text:
                if ch == "\n":
                    lines += 1; w = 0.0; continue
                cw = em * (1.0 if ord(ch) > 0x2E80 else 0.52)
                if w + cw > width_in and w > 0:
                    lines += 1; w = cw
                else:
                    w += cw
        total += lines
    return total or 1


def _last_line(paras, width_in):
    """Return the text on the LAST wrapped visual line (same CJK-aware wrap estimate as _est_lines) —
    so we can catch a lone trailing punctuation mark or a single orphaned glyph (避头尾 / widow)."""
    if width_in <= 0 or not paras:
        return ""
    line = ""
    for pr in paras:
        line, w = "", 0.0
        for text, size_pt in pr:
            em = max(0.01, size_pt / 72.0)
            for ch in text:
                if ch == "\n":
                    line, w = "", 0.0; continue
                cw = em * (1.0 if ord(ch) > 0x2E80 else 0.52)
                if w + cw > width_in and w > 0:
                    line, w = ch, cw          # this char starts a new line
                else:
                    line += ch; w += cw
    return line                                # last paragraph's last line = the box's last line


# punctuation that must never stand alone at the start of a line (closing marks)
_CLOSERS = set("。．，、！？：；）》】」』〕〗｝….,!?:;)]}、。")


def _walk_runs(shapes):
    """Yield every text run on the slide — recursing into GROUPS and descending into TABLE cells —
    so checks cover ANY text scenario, not just top-level text boxes."""
    for s in shapes:
        try:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk_runs(s.shapes); continue
            if getattr(s, "has_table", False):
                for row in s.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            yield from p.runs
                continue
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    yield from p.runs
        except Exception:
            continue


def _run_cjk_no_ea(run):
    """True if the run has CJK glyphs but no <a:ea> font (→ no kinsoku + tofu/uncontrolled-font risk)."""
    try:
        if not any(ord(ch) > 0x2E80 for ch in run.text):
            return False
        rPr = run._r.find(qn("a:rPr"))
        return rPr is None or rPr.find(qn("a:ea")) is None
    except Exception:
        return False


def _inter(a, b):
    return (max(0.0, min(a["r"], b["r"]) - max(a["l"], b["l"])),
            max(0.0, min(a["b"], b["b"]) - max(a["t"], b["t"])))


def _frac_inside(a, b):
    ix, iy = _inter(a, b)
    return (ix * iy) / (a["w"] * a["h"] + 1e-9)


def lint(path):
    try:
        prs = Presentation(path)
    except Exception:
        print(f"lint_deck: cannot open '{path}' (missing or not a .pptx)", file=sys.stderr)
        sys.exit(2)
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    total = 0
    warn_total = 0
    for si, slide in enumerate(prs.slides):
        bx = _boxes(slide, sw, sh)
        finds = []
        # 1) overflow
        for s in bx:
            if s["bg"]:
                continue
            ov = []
            if s["l"] < -TOL: ov.append("left")
            if s["t"] < -TOL: ov.append("top")
            if s["r"] > sw + TOL: ov.append(f"right+{round(s['r']-sw,2)}")
            if s["b"] > sh + TOL: ov.append(f"bottom+{round(s['b']-sh,2)}")
            if ov:
                finds.append(f"OVERFLOW [{','.join(ov)}] {s['st']} '{s['txt']}'")
        # 2) solid vs solid partial overlap (neither contained)
        sol = [s for s in bx if s["solid"] and not s["bg"]]
        for i in range(len(sol)):
            for j in range(i + 1, len(sol)):
                a, b = sol[i], sol[j]
                ix, iy = _inter(a, b)
                # A clean overlap is one where BOTH dims clear the hairline tolerance, OR a
                # "bar over a card" signature: one dim is broad (>= 0.5in) while the other is
                # only a sliver (a wide bottom bar dipping a few pt into the card above it, or a
                # tall sidebar grazing content). That sliver still CLIPS rounded corners / sits
                # on top, so flag it even though it's under TOL in the short dim. The 0.012in
                # floor still ignores true sub-pixel coincidence (flush edges).
                clean = ix > TOL and iy > TOL
                bar_on_block = (max(ix, iy) >= 0.5 and min(ix, iy) > 0.012)
                if (clean or bar_on_block) and _frac_inside(a, b) < CONTAIN and _frac_inside(b, a) < CONTAIN:
                    # skip a hard offset-shadow / sticker pair: same size, offset by a few pt (intentional)
                    same_size = abs(a["w"] - b["w"]) < 0.06 and abs(a["h"] - b["h"]) < 0.06
                    tiny_offset = abs(a["l"] - b["l"]) < 0.18 and abs(a["t"] - b["t"]) < 0.18
                    if same_size and tiny_offset:
                        continue
                    finds.append(f"OVERLAP {round(ix,2)}x{round(iy,2)}in  {a['st']}'{a['txt']}' x {b['st']}'{b['txt']}'")
        # 3) footer-zone reservation: the bottom footer band is deck chrome — NO content block (solid
        #    or text) may cover the footer text or dip into its reserved zone. (No containment escape:
        #    a content band whose bottom collides with the footer is a bug even when the footer text
        #    falls inside the band's bbox — that "text on a card" exclusion is the blind spot that let
        #    a low band overlap the footer.)
        footers = [s for s in bx if s["text"] and s["t"] > sh - 0.6 and not s["bg"]]
        if footers:
            fid = {id(f) for f in footers}
            fb_top = min(f["t"] for f in footers) - 0.08      # keep content above this line
            for s in bx:
                if s["bg"] or id(s) in fid or not (s["solid"] or s["text"]) or s["h"] <= 0.2:
                    continue
                if s["t"] < fb_top - 0.04 and s["b"] > fb_top:    # starts in the content area, dips into the footer zone
                    covers = [f for f in footers if _inter(s, f)[0] > TOL and _inter(s, f)[1] > TOL]
                    if covers:
                        finds.append(f"FOOTER collision  {s['st']}'{s['txt']}' overlaps footer '{covers[0]['txt']}' "
                                     f"— keep content above {round(fb_top,2)}in (content_band()/bottom_callout())")
                    else:
                        finds.append(f"FOOTER-ZONE intrusion  {s['st']}'{s['txt']}' (bottom {round(s['b'],2)}in) "
                                     f"dips into the reserved footer band — keep content above {round(fb_top,2)}in")
        # 4) editability: a slide that is ~one whole-page image with no native text/objects
        bigpic = next((s for s in bx if s["st"] == "PICTURE" and s["w"] * s["h"] >= 0.85 * sw * sh), None)
        if bigpic is not None:
            others = [s for s in bx if s is not bigpic and not s["bg"] and (s["text"] or s["w"] * s["h"] > 0.5)]
            if not others:
                finds.append("EDITABILITY: slide is ~one whole-page image with no native text/objects "
                             "(build content as native shapes; images are plates/figures, not the whole slide)")
        # 5) orphan / blank slide
        if not [s for s in bx if not s["bg"] and (s["text"] or s["solid"])]:
            finds.append("EMPTY/ORPHAN slide: no native content (blank or background only)")
        # 6) INTERIOR PADDING / OVERFLOW: text inside a filled card must keep a minimum margin on every
        #    side, measured against its RENDERED extent — not merely "not overflow the card". Catches both
        #    text running PAST the card and text CRAMMED against an edge (the recurring "too close to the
        #    boundary" bug that a pure-overflow check misses).
        PAD = 0.09                                       # minimum interior padding (in)
        def _cjk(t): return any(ord(ch) > 0x2E80 for ch in t["full"])
        def _txt_h(t):
            return _est_lines(t["paras"], t["w"]) * (t["size"] / 72.0) * (1.4 if _cjk(t) else 1.25)
        def _nat_width(t):                               # natural one-line width of the widest paragraph
            return max((sum((sz / 72.0) * (1.0 if ord(ch) > 0x2E80 else 0.52) for s_, sz in pr for ch in s_)
                        for pr in t["paras"]), default=0.0)
        def _rbox(t):
            """The RENDERED bounding box of the text — alignment- and anchor-aware (so a centred /
            middle-anchored label isn't mistaken for cramped-against-the-edge)."""
            tw = min(_nat_width(t), t["w"]) or t["w"]; eh = _txt_h(t)
            al = t["align"]
            if al == PP_ALIGN.CENTER:   rl = t["l"] + (t["w"] - tw) / 2; rr = rl + tw
            elif al == PP_ALIGN.RIGHT:  rr = t["r"]; rl = rr - tw
            else:                       rl = t["l"]; rr = rl + tw
            an = t["anchor"]
            if an == MSO_ANCHOR.MIDDLE: rt = t["t"] + (t["h"] - eh) / 2; rb = rt + eh
            elif an == MSO_ANCHOR.BOTTOM: rb = t["b"]; rt = rb - eh
            else:                       rt = t["t"] + 0.04; rb = rt + eh
            return rl, rt, rr, rb
        # a "card" for padding is a SMALL filled block (not a full/half-slide scrim or background plate)
        cards = [s for s in bx if s["solid"] and not s["bg"] and not s["text"] and s["h"] > 0.35
                 and s["w"] * s["h"] < 0.45 * sw * sh]
        for t in [s for s in bx if s["text"]]:
            host = None
            for c in cards:
                if c["l"] - 0.12 <= t["l"] and t["r"] <= c["r"] + 0.12 and c["t"] - 0.12 <= t["t"] <= c["b"] - 0.05:
                    if host is None or c["b"] < host["b"]:
                        host = c
            if not host:
                continue
            rl, rt, rr, rb = _rbox(t)
            nlines = _est_lines(t["paras"], t["w"])
            if rb > host["b"] - PAD:                          # rendered text crammed against / past the card bottom
                kind = "runs PAST" if rb > host["b"] + 0.03 else "is cramped against (< pad)"
                finds.append(f"TEXT PADDING: '{t['txt']}' (~{nlines} lines) {kind} the card bottom "
                             f"({round(host['b'],2)}in) — size the card to the text, shorten, or add bottom padding")
        # 6d) CHIP / LABEL too small: a text box coincident with a small filled pill whose text overruns it
        #     (a row of chips where the label is wider/taller than the pill — text cramped to the pill edges).
        fills = [s for s in bx if s["solid"] and not s["text"] and not s["bg"]]
        for t in [s for s in bx if s["text"] and s["h"] < 0.55]:
            pill = next((c for c in fills if abs(c["l"] - t["l"]) < 0.06 and abs(c["t"] - t["t"]) < 0.06
                         and abs(c["w"] - t["w"]) < 0.16 and abs(c["h"] - t["h"]) < 0.16), None)
            if pill:
                nl = _est_lines(t["paras"], t["w"] - 0.10)         # inner pad
                if nl * (t["size"] / 72.0) * (1.4 if _cjk(t) else 1.25) > t["h"] - 0.02:
                    finds.append(f"CHIP/LABEL TOO SMALL: '{t['txt']}' (~{nl} lines) overruns its "
                                 f"{round(t['w'],2)}×{round(t['h'],2)}in pill — size the chip to its text (or shorten)")
        # 6e) TEXT-vs-TEXT collision: the RENDERED bottom of an upper text box overruns the rendered TOP of a
        #     lower one in the same column (the wrap-collision a 'declared-box' overlap check never sees).
        txts = [s for s in bx if s["text"] and not s["bg"] and s["t"] < sh - 0.55 and len(s["full"].strip()) > 1]
        for a in txts:
            arl, art, arr, arb = _rbox(a)
            for b in txts:
                if b is a:
                    continue
                brl, brt, brr, brb = _rbox(b)
                if brt <= art + 0.05:                                  # b must sit clearly below a
                    continue
                if min(arr, brr) - max(arl, brl) < 0.30:               # must share an x-column
                    continue
                if arb > brt + 0.06:
                    finds.append(f"TEXT COLLISION: '{a['txt']}' (~{_est_lines(a['paras'], a['w'])} lines) overruns "
                                 f"into the text below '{b['txt']}' — add vertical gap or size the box")
                    break
        # 6b) orphaned punctuation / widow: a wrapped box whose LAST line is just a punctuation mark
        #     (the 避头尾 bug — a lone 。/，pushed to its own row) or a single orphaned CJK glyph
        for t in [s for s in bx if s["text"] and s["w"] > 0]:
            if _est_lines(t["paras"], t["w"]) < 2:
                continue
            ll = _last_line(t["paras"], t["w"]).strip()
            if ll and all(c in _CLOSERS for c in ll):
                finds.append(f"ORPHANED PUNCTUATION: the last line of '{t['txt']}' is just '{ll}' — "
                             f"widen the box / lower the size / reword so the mark stays attached (避头尾)")
            elif len(ll) == 1 and ord(ll) > 0x2E80:
                finds.append(f"WIDOW: a lone glyph '{ll}' stranded on the last line of '{t['txt']}' — "
                             f"widen the box or reword so the last line has company")
        # 6c) CJK text with NO East-Asian font — the ROOT cause of orphaned punctuation (no <a:ea> →
        #     PowerPoint applies no kinsoku (避头尾), so a 。/，can start a line; also tofu/uncontrolled
        #     font). Checked across ANY text scenario — text boxes, TABLE cells, and grouped shapes —
        #     not just top-level boxes. Reliable, render-independent. Fix: set deckkit.EAFONT (+ EADISPLAY).
        bad_ea = [r for r in _walk_runs(slide.shapes) if _run_cjk_no_ea(r)]
        if bad_ea:
            sample = next((r.text.strip() for r in bad_ea if r.text.strip()), "")[:18]
            finds.append(f"CJK TEXT without an EA font: {len(bad_ea)} run(s) (e.g. '{sample}', incl. any "
                         f"table/grouped text) have CJK with no East-Asian font — set deckkit.EAFONT "
                         f"(no kinsoku → orphaned punctuation; uncontrolled font → tofu)")
        # 7) uneven card heights in a row (sibling cards must share ONE height)
        cset = [s for s in bx if s["solid"] and not s["bg"] and not s["text"] and s["h"] > 0.5 and s["w"] < 0.6 * sw]
        bycol = {}                                  # dedupe layered shapes per (top,left): keep the tallest (the card, not its header band)
        for c in cset:
            key = (round(c["t"] * 10), round(c["l"] * 10))
            if key not in bycol or c["h"] > bycol[key]["h"]:
                bycol[key] = c
        rows = {}
        for c in bycol.values():
            rows.setdefault(round(c["t"] * 10), []).append(c)
        for _top, row in rows.items():
            if len(row) >= 2 and (max(r["w"] for r in row) - min(r["w"] for r in row)) < 0.4:  # a row of similar-width siblings
                hs = [r["h"] for r in row]
                if max(hs) - min(hs) > 0.12:
                    finds.append(f"UNEVEN CARD HEIGHTS: a row of {len(row)} cards has heights "
                                 f"{sorted(round(h,2) for h in hs)} — sibling cards in a row must share ONE "
                                 f"height (size the row to the tallest card's content)")
        # --- soft WARNs: advisory (do NOT fail the build) — accessibility + font-portability nudges
        #     that are otherwise invisible to the critic and the geometry checks ---
        warns = []
        for s in bx:
            if s["st"] == "PICTURE" and not s["bg"] and _no_real_alt(s["descr"]):
                warns.append(f"MISSING ALT-TEXT: an informative image ({round(s['w'],1)}×{round(s['h'],1)}in) has no "
                             f"alt-text — deckkit.alt_text(shape, '…'), or pass alt='' if purely decorative")
            if s["mathfont"] and _fsub(s["mathfont"]):
                warns.append(f"MATH-FONT TOFU RISK: '{s['txt']}' uses {s['mathfont']}, which isn't installed on this "
                             f"render host — install it, or set deckkit.EQ_MATHFONT to a math font present here")
        for m in finds:
            print(f"  slide {si+1}: {m}")
        for m in warns:
            print(f"  slide {si+1}: [warn] {m}")
        total += len(finds)
        warn_total += len(warns)
    tail = ("" if total else "  ✓ clean (no hard findings)") + (f"  ·  {warn_total} warning(s)" if warn_total else "")
    print(f"\n{path}: {total} layout finding(s){tail}")
    return total


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: python lint_deck.py <deck.pptx>"); sys.exit(2)
    sys.exit(1 if lint(sys.argv[1]) > 0 else 0)
