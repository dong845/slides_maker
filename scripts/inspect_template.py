#!/usr/bin/env python3
"""Inspect a .pptx template so you can build new slides on it correctly.

Prints: slide size, all layouts (index + name), the layouts used by the existing
title and content slides, their placeholder ids/types/geometry, and which media
(logos/bands) the master and layouts reference.

Usage:  python3 inspect_template.py /path/to/template.pptx
"""
import sys, zipfile, re
from pptx import Presentation
from pptx.util import Emu

path = sys.argv[1]
prs = Presentation(path)

print(f"slide size: {Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in")

print("\n=== layouts (index : name) ===")
for i, l in enumerate(prs.slide_layouts):
    print(f"  {i:2d} : {l.name}")

def ph_dump(layout):
    for ph in layout.placeholders:
        f = ph.placeholder_format
        print(f"     idx={f.idx} type={f.type} '{ph.name}'  "
              f"x={Emu(ph.left).inches:.2f} y={Emu(ph.top).inches:.2f} "
              f"w={Emu(ph.width).inches:.2f} h={Emu(ph.height).inches:.2f}")

slides = list(prs.slides)   # python-pptx Slides has no slice support; materialize first
if slides:
    t = slides[0].slide_layout
    print(f"\n=== existing TITLE slide uses layout {list(prs.slide_layouts).index(t)} '{t.name}' ===")
    ph_dump(t)
    for s in slides[1:]:
        c = s.slide_layout
        if c.name != t.name:
            print(f"\n=== existing CONTENT slide uses layout {list(prs.slide_layouts).index(c)} '{c.name}' ===")
            ph_dump(c)
            break

print("\n=== media referenced by master / layouts (logos & bands live here) ===")
with zipfile.ZipFile(path) as z:
    for rels in sorted(n for n in z.namelist()
                       if re.search(r'slide(Master|Layout)\d+\.xml\.rels$', n)):
        data = z.read(rels).decode('utf-8', 'ignore')
        imgs = sorted(set(re.findall(r'media/(image[0-9]+\.[a-z]+)', data)))
        if imgs:
            print(f"  {rels.split('/')[-1]}: {', '.join(imgs)}")
    import tempfile
    print("\n  tip: .emf files are usually vector LOGOS; .png on the master is often"
          "\n       the decorative header band. Convert an .emf to view it with:"
          f"\n       soffice --headless --convert-to png --outdir {tempfile.gettempdir()} ppt/media/imageN.emf"
          "\n       (on Windows, soffice may be soffice.exe / soffice.com)")
