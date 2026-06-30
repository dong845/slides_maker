#!/usr/bin/env python3
"""deckkit — reusable helpers for building clean slides with python-pptx.

General-purpose: works whether you build ON a user's template (open_template) or
from scratch when they have none (blank_deck + title_bar/footer). Import this from
a small per-deck build script rather than copy-pasting primitives every time.
A brand-free worked example lives at references/examples/build_example_generic.py.

Design intent (see references/design-principles.md — these are the PRESENTED-deck defaults; a
read-alone / reference / poster deck legitimately runs denser, so flex by delivery mode):
  - terse, few-word points for a SPOKEN deck (a visual aid, not a document); a read-alone deck
    carries the fuller sentence a speaker would otherwise say
  - native diagrams (boxes + arrows) over walls of text
  - a clear title per slide; results figures always get a legend + a takeaway

The palette/fonts below are sensible DEFAULTS. When building on a template, pull
the real brand colors from it (inspect_template.py / theme) and set FONT to match;
when building from scratch, pick colors from the user's brand or pick a clean set.

Equations: NEVER rely on Unicode modifier-letter super/subscripts (ᴴ ᵀ ᵣ) — many
display fonts lack those glyphs and render tofu/overlap. Use eq_par(), which draws
real baseline-shifted ASCII in a full-coverage font, so it is crisp anywhere.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import math

# ---- default professional palette (a neutral blue scheme). NOT tied to any brand —
# when building on a template, override these with the template's real theme colours.
DEEP    = RGBColor(0x00, 0x3C, 0x66)   # deep navy — strong text / dark panels
BLUE    = RGBColor(0x00, 0x7C, 0xC2)   # accent blue — bullets, primary boxes
TEAL    = RGBColor(0x00, 0x9F, 0xBD)   # secondary accent
MAGENTA = RGBColor(0xE3, 0x00, 0x4F)   # highlight / "after" / callout rule
SLATE   = RGBColor(0x3A, 0x4A, 0x55)   # body text
MUTE    = RGBColor(0x5A, 0x64, 0x72)   # captions / secondary (clears 4.5:1 on white; code
                                       # panels use text_c, so this is safe on dark too)
TINT    = RGBColor(0xEA, 0xF3, 0xFA)   # light callout fill
LIGHT   = RGBColor(0xF4, 0xF8, 0xFB)   # lighter fill
PALE    = RGBColor(0xBF, 0xDE, 0xF0)   # pale text on dark panels
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
# secondary accents — for VARIETY so diagrams don't read monotone. The first two
# are from the template theme (brand-coherent); rotate through ACCENTS for chips.
GOLD    = RGBColor(0xC0, 0x96, 0x5C)   # theme accent5
STEEL   = RGBColor(0x6E, 0x90, 0xA6)   # theme accent3
VIOLET  = RGBColor(0x6A, 0x4C, 0x93)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
ACCENTS = [BLUE, TEAL, GOLD, STEEL, VIOLET, GREEN]   # cycle for multi-item diagrams

# layout: keep a consistent gutter of whitespace between a figure and any adjacent
# text/callout/edge. Crowding elements together reads as amateur — give them room.
GUTTER  = 0.4   # inches
# the bottom band reserved for footer() chrome (tag + page number sit at h_in - 0.35).
# content_band()/bottom_callout() keep content out of this zone, so a bottom callout can
# never grow into the footer — the recurring "callout collided with the footer" failure.
FOOTER_BAND = 0.5   # inches, measured up from the slide's bottom edge


def contrast_ratio(c1, c2):
    """WCAG contrast ratio between two colours (RGBColor or 'RRGGBB' hex string).
    Returns a number from 1 (none) to 21 (black-on-white). Body text wants >= 4.5;
    large/bold text >= 3. Use it to sanity-check a text colour against its fill before
    you even render — e.g. a pale caption on a tint, or grey body on white.
        if contrast_ratio(MUTE, WHITE) < 4.5: ...  # too faint, darken it
    Note: this checks the *colour pair*; the render still tells you about size/overlap."""
    def lum(c):
        if isinstance(c, str):
            c = RGBColor.from_string(c)
        chans = []
        for v in (c[0], c[1], c[2]):
            s = v / 255.0
            chans.append(s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4)
        return 0.2126 * chans[0] + 0.7152 * chans[1] + 0.0722 * chans[2]
    l1, l2 = sorted((lum(c1), lum(c2)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


def _rgb_dist(a, b):
    """Euclidean RGB distance — a quick proxy for 'are these two fills visually distinct?'.
    (Use this, NOT contrast_ratio, for category colours: two different hues can share a
    luminance — cyan vs amber — so contrast_ratio reads ~1 while they're clearly distinct.)"""
    if isinstance(a, str): a = RGBColor.from_string(a)
    if isinstance(b, str): b = RGBColor.from_string(b)
    return sum((a[i] - b[i]) ** 2 for i in range(3)) ** 0.5


def palette(n, accents=None):
    """Return ``n`` DISTINCT categorical fills for a sequence of blocks/chips/cards.

    A sequence of blocks must read as a *thought-through* set of colours, not an accident:
    each block a deliberate, well-separated hue, **adjacent blocks visibly different**, and
    **no neutral gray dropped in as a category** (gray reads as disabled/secondary, not a
    category — reserve it for genuinely de-emphasised items). Pass the deck's own palette
    (e.g. your style's ``ACCENTS``); defaults to deckkit's ``ACCENTS``. Cycles with a warning
    if ``n`` exceeds the available hues (extend ``ACCENTS`` rather than repeat a category
    colour), and warns if any two ADJACENT fills are perceptually too close — so a same-colour
    or near-duplicate pair surfaces at build time instead of in the render.

        fills = dk.palette(4, ACCENTS)          # 4 distinct, well-separated hues
        for (x,y,w,h), label, fill in zip(cells, labels, fills):
            dk.chip(s, x, y, w, h, *label, fill)
    """
    pool = list(ACCENTS if accents is None else accents)
    if not pool:
        raise ValueError("palette needs a non-empty accents list")
    if n > len(pool):
        import warnings
        warnings.warn(f"palette({n}) exceeds {len(pool)} distinct hues — colours will repeat; "
                      f"add more entries to ACCENTS rather than reuse a category colour")
    fills = [pool[i % len(pool)] for i in range(n)]
    import warnings
    for a, b in zip(fills, fills[1:]):
        if _rgb_dist(a, b) < 60:                 # ~perceptually similar / identical
            warnings.warn("palette: two adjacent category fills are nearly identical — "
                          "reorder or extend ACCENTS so neighbouring blocks contrast")
    for c in fills:
        cc = c if not isinstance(c, str) else RGBColor.from_string(c)
        if max(cc[0], cc[1], cc[2]) - min(cc[0], cc[1], cc[2]) < 30:   # near-neutral / gray
            warnings.warn("palette: a near-gray fill is used as a category colour — gray reads "
                          "as disabled/secondary, not a category; use a saturated hue (reserve "
                          "gray for genuinely de-emphasised items)")
    return fills


def palette_from_image(path, n=5, *, keep_neutrals=False):
    """Extract ``n`` representative ACCENT colours from a generated template image, as
    ``RGBColor``, most-frequent first.

    This is the bridge that makes NATIVE content "fit a generated template": derive the
    template's palette from its hero/divider image, set your ``style.py`` colours to it, and
    every native card / chip / heading / motif comes out in the same hues — so the inserted
    blocks read as part of the generated look rather than pasted on top. By default it skips
    near-white and near-black (the background/ink, not accents) and the deck's *base* colour;
    pass ``keep_neutrals=True`` to keep them. Deduplicates perceptually-close hues.

        BASE = palette_from_image("assets/template_bg.png", 6)   # the template's own colours
        ACCENTS = BASE                                           # native content now matches
    """
    from PIL import Image
    from collections import Counter
    im = Image.open(path).convert("RGB")
    im.thumbnail((220, 220))
    q = im.quantize(colors=64, method=Image.FASTOCTREE).convert("RGB")
    out = []
    for (r, g, b), _ in Counter(q.getdata()).most_common():
        mx, mn = max(r, g, b), min(r, g, b)
        if not keep_neutrals:
            if mx > 236 and mn > 226:        # near-white background
                continue
            if mx < 30:                      # near-black ink
                continue
            if mx - mn < 22:                 # washed/near-gray, not an accent
                continue
        c = RGBColor(r, g, b)
        if all(_rgb_dist(c, o) > 48 for o in out):   # keep hues visibly distinct
            out.append(c)
        if len(out) >= n:
            break
    return out

FONT    = "Calibri"       # display font — cross-platform default (ships with MS Office,
                          # renders the same on Windows PowerPoint and macOS/Keynote).
MONO    = "Consolas"      # code / filenames — Calibri's cross-platform monospace pair.
EQFONT  = "Arial"         # equations — universal glyph + Greek coverage
EAFONT  = None            # East-Asian font for CJK text (e.g. "PingFang SC" / "Heiti SC"
                          # / "Microsoft YaHei" / "Noto Sans CJK SC"). When set, every run
                          # ALSO carries an <a:ea> typeface so PowerPoint/Keynote render
                          # Chinese/Japanese/Korean glyphs with THIS font (not an
                          # uncontrolled default) while Latin/numbers stay on FONT. Leave
                          # None for Latin decks. See references/multilingual.md.
DISPLAY = None            # optional DISPLAY/title font (Latin) — when set, title_bar uses it for
                          # the title so headings get their own face vs the FONT body. Falls back
                          # to FONT. Pairing roles (display / body / mono) beats one font for the
                          # whole deck — see references/font-guidance.md ("Type pairing").
EADISPLAY = None          # optional CJK DISPLAY/title font (e.g. "PingFang SC" titles over a
                          # "Hiragino Sans GB"/"Noto Sans CJK SC" EAFONT body). Falls back to EAFONT.
# To re-theme a whole deck (e.g. to match a style example), reassign these AND the
# palette constants above right after importing deckkit, before building — set_font
# resolves FONT at call time, so `deckkit.FONT = "Helvetica Neue"` takes effect.


# ====================================================================== text
def _apply_ea(run, typeface):
    """Set the East-Asian (<a:ea>) typeface so PowerPoint/Keynote render CJK glyphs with
    the chosen font (Latin chars keep the <a:latin> font). python-pptx only writes
    <a:latin>, so we add <a:ea> directly, in the correct schema position (after latin)."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        latin = rPr.find(qn('a:latin'))
        (latin.addnext(ea) if latin is not None else rPr.append(ea))
    ea.set('typeface', typeface)

def set_font(run, size, color, bold=False, italic=False, font=None, ea=None):
    run.font.name = font or FONT          # resolve FONT at call time so re-theming works
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    eaf = ea or EAFONT                     # also tag CJK font when set (mixed CN/EN stays correct)
    if eaf:
        _apply_ea(run, eaf)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs = list of paragraphs; each paragraph = list of run tuples
    (txt, size, color, bold, italic[, font]).

    Vertical centring: to centre text inside a filled box/card, pass anchor=MSO_ANCHOR.MIDDLE
    AND give this textbox the SAME (x, y, w, h) as the box. A y-offset (e.g. y+0.07) combined
    with the box's full height pushes the centre below the box's true middle — text then reads
    "a bit low". Want top padding? Use margin_top, not a y-offset with unchanged height."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, *rest) in para:
            r = p.add_run(); r.text = txt
            set_font(r, size, color, bold, italic, rest[0] if rest else None)
    return tb


# ===================================================================== shapes
def _grad_fill(shape, stops, angle=90.0, radial=False):
    """Apply a gradient fill WITH PER-STOP ALPHA to a shape (python-pptx solid fills can't do
    this). `stops` = list of (pos 0..1, colour as 'RRGGBB'/RGBColor, alpha 0..1). `angle` is the
    linear direction in degrees (0=→, 90=↓); `radial=True` makes a centre-out radial (for glows).
    This is the enabler for glass cards, soft glows, and graduated photo scrims."""
    sp = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        e = sp.find(qn(tag))
        if e is not None:
            sp.remove(e)
    gs = "".join(
        f'<a:gs pos="{int(round(max(0.0,min(1.0,p))*100000))}">'
        f'<a:srgbClr val="{_hex(c)}"><a:alpha val="{int(round(max(0.0,min(1.0,a))*100000))}"/></a:srgbClr></a:gs>'
        for (p, c, a) in stops)
    direction = ('<a:path path="circle"><a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path>'
                 if radial else f'<a:lin ang="{int(round((angle % 360) * 60000))}" scaled="1"/>')
    el = parse_xml(f'<a:gradFill {nsdecls("a")}><a:gsLst>{gs}</a:gsLst>{direction}</a:gradFill>')
    geom = sp.find(qn("a:prstGeom"))
    if geom is None:
        geom = sp.find(qn("a:custGeom"))
    ln = sp.find(qn("a:ln"))
    if geom is not None:
        geom.addnext(el)
    elif ln is not None:
        ln.addprevious(el)
    else:
        sp.append(el)
    return shape


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, round=False, corners="all", r=None,
        grad=None, grad_angle=90.0, grad_radial=False):
    """A rectangle. `round=True` rounds all four corners (radius = 8% of the shorter side,
    or `r` inches if given). For a colored HEADER BAND sitting on top of a rounded card,
    use `corners='top'` and pass `r=<the card's corner radius in inches>` so the band's
    curve MATCHES the card — a square band over a rounded card (corners poking out) is the
    tell to avoid. `corners='bottom'` rounds the bottom two. (A thin accent strip can
    instead be inset by the radius so its square ends fall on the card's straight edge.)

    `grad` gives a GRADIENT fill with per-stop alpha instead of a solid `fill`: a list of
    (pos 0..1, colour, alpha 0..1); `grad_angle` sets linear direction (deg), `grad_radial=True`
    a centre-out radial. Powers glass/glow/scrim — usually via the `glass_card`/`glow`/
    `scrim_overlay` helpers rather than called directly."""
    if not (round or r is not None or corners != "all"):
        t = MSO_SHAPE.RECTANGLE
    elif corners in ("top", "bottom"):
        t = MSO_SHAPE.ROUND_2_SAME_RECTANGLE   # rounds the two top corners (rotate for bottom)
    else:
        t = MSO_SHAPE.ROUNDED_RECTANGLE
    s = slide.shapes.add_shape(t, Inches(x), Inches(y), Inches(w), Inches(h))
    if grad is not None: _grad_fill(s, grad, angle=grad_angle, radial=grad_radial)
    elif fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = _as_rgb(fill)
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = _as_rgb(line); s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if t != MSO_SHAPE.RECTANGLE:
        adj = (r / min(w, h)) if r is not None else 0.08
        adj = max(0.0, min(0.5, adj))
        try: s.adjustments[0] = adj
        except Exception: pass
        if corners == "bottom":
            s.rotation = 180
    return s


def glow(slide, cx, cy, w, h, color, alpha=0.5):
    """A soft radial colour GLOW (centre-out) for depth/atmosphere on a DARK slide — place 1-2
    off-centre behind glass cards so a flat black slide gets dimensional lighting. Invisible on
    light decks; don't use there. (cx,cy) is the glow centre in inches."""
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h))
    o.line.fill.background(); o.shadow.inherit = False
    _grad_fill(o, [(0.0, color, alpha), (1.0, color, 0.0)], radial=True)
    return o


def scrim_overlay(slide, x, y, w, h, *, stops=((0.0, 0.0), (1.0, 0.75)), color="000000", angle=90.0):
    """A GRADUATED alpha scrim over a photo so overlaid text stays legible WHERE THE TEXT IS,
    while the image stays bright elsewhere — far better than a flat dark overlay. `stops` =
    (position, alpha) pairs; aim the gradient toward the text via `angle` (90 = darker at bottom,
    270 = darker at top). Size the rect to the text's zone, not always the whole slide."""
    return box(slide, x, y, w, h, grad=[(p, color, a) for (p, a) in stops], grad_angle=angle)


def glass_card(slide, x, y, w, h, tint, *, accent=None, r=0.14, rim=1.0):
    """A frosted-glass card (UI glassmorphism) rebuilt natively in three layers: (1) a low-alpha
    `tint` gradient body, (2) a white diagonal sheen, (3) a 1px white rim. It only reads as glass
    on a DARK / glowing / photographic background — pair with `glow()` on a dark base. Optional
    `accent` adds a colored top header band. Returns the body shape (place content on top)."""
    body = box(slide, x, y, w, h, round=True, r=r,
               grad=[(0.0, tint, 0.46), (1.0, tint, 0.20)], grad_angle=120)
    body.line.color.rgb = WHITE; body.line.width = Pt(rim)
    sheen = box(slide, x, y, w, h, round=True, r=r,
                grad=[(0.0, "FFFFFF", 0.20), (0.45, "FFFFFF", 0.05), (1.0, "FFFFFF", 0.0)], grad_angle=120)
    sheen.line.fill.background()
    if accent is not None:
        box(slide, x, y, w, 0.58, corners="top", r=r,
            grad=[(0.0, accent, 0.95), (1.0, accent, 0.70)], grad_angle=0)
    return body


def offset_shadow(slide, x, y, w, h, fill, *, dx=0.06, dy=0.06, shadow=None,
                  line=None, line_w=2.0, round=True, r=0.1):
    """A HARD offset 'sticker' / letterpress shadow (riso / print look): a crisp solid shadow
    copy behind the shape, offset by (dx,dy) — NOT a soft blur. Returns the top shape so you can
    place text on it. Use for bold/editorial/retro-print decks; skip in minimal/scientific ones."""
    sh = shadow if shadow is not None else RGBColor(0x1B, 0x1B, 0x1B)
    box(slide, x + dx, y + dy, w, h, fill=sh, round=round, r=r)
    return box(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, round=round, r=r)


_GOOD = RGBColor(0x1F, 0x9D, 0x55)   # positive delta (green)
_BAD = RGBColor(0xE0, 0x3A, 0x2E)    # negative delta (red)

def scorecard(slide, x, y, w, h, label, value, *, delta=None, caption=None, good_up=True,
              ink=DEEP, accent=BLUE, glass_tint=None):
    """A KPI scorecard tile: small-caps label · oversized value · colored ▲/▼ delta · tiny
    caption — the 'current state in numbers' building block. `value`/`label` may be numbers or
    strings (coerced). `delta` is a string ('+3.2pp' / '-18%'); an unsigned value reads as an
    increase, so **a sign is required to mark a decrease**. Its colour is auto-set GREEN/RED by
    direction vs `good_up` (a falling cost with good_up=False is green). `glass_tint` makes it a
    glass tile (use on dark decks). Lay out 3-6 with columns()."""
    if glass_tint is not None:
        glass_card(slide, x, y, w, h, glass_tint); lab_c, val_c, cap_c = WHITE, WHITE, RGBColor(0xCF, 0xD7, 0xE6)
    else:
        box(slide, x, y, w, h, fill=WHITE, line=RGBColor(0xE3, 0xE8, 0xEE), line_w=1.0, round=True, r=0.1)
        box(slide, x, y, 0.1, h, fill=accent, round=True, r=0.05)          # accent spine
        lab_c, val_c, cap_c = MUTE, ink, MUTE
    px = x + 0.28
    text(slide, px, y + 0.22, w - 0.5, 0.3, [[(str(label).upper(), 11, lab_c, True, False)]], space_after=0)
    text(slide, px, y + 0.5, w - 0.5, 0.8, [[(str(value), 33, val_c, True, False)]], space_after=0)
    cy = y + 1.32
    if delta:
        d = str(delta).strip()
        up = not d.startswith(("-", "▼"))            # unsigned / '+' = increase; sign required for a decrease
        dc = _GOOD if (up == good_up) else _BAD
        text(slide, px, cy, w - 0.5, 0.3, [[(("▲ " if up else "▼ ") + d.lstrip("+-▲▼ "), 12, dc, True, False)]], space_after=0)
        cy += 0.3
    if caption:
        text(slide, px, cy, w - 0.5, 0.4, [[(caption, 10.5, cap_c, False, False)]], space_after=0)


def leaderboard(slide, x, y, w, rows, *, row_h=0.5, gap=0.1, ink=DEEP):
    """Ranked / part-to-whole list keyed to a chart: each row = a colored left swatch (matching a
    chart wedge/series) + name + right-aligned value (+ optional sub). Pass the SAME colour list
    you built the chart with so legend and chart stay in sync. rows = [(color, name, value[, sub])]."""
    cy = y
    for row in rows:
        color = row[0]; name = str(row[1]); value = str(row[2]) if len(row) > 2 else ""
        sub = str(row[3]) if len(row) > 3 else None
        box(slide, x, cy, 0.12, row_h, fill=color, round=True, r=0.04)
        if sub:                                              # name + sub both stay INSIDE this row
            text(slide, x + 0.28, cy + 0.02, w - 2.0, row_h * 0.58, [[(name, 13, ink, True, False)]],
                 anchor=MSO_ANCHOR.BOTTOM, space_after=0)
            text(slide, x + 0.28, cy + row_h * 0.56, w - 2.0, row_h * 0.42, [[(sub, 9, MUTE, False, False)]], space_after=0)
        else:
            text(slide, x + 0.28, cy, w - 2.0, row_h, [[(name, 13, ink, True, False)]],
                 anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        text(slide, x + w - 1.9, cy, 1.9, row_h, [[(value, 14, ink, True, False)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        cy += row_h + gap
        hrule(slide, x, cy - gap / 2, w, color=RGBColor(0xE6, 0xE9, 0xEE), weight=0.01)
    return cy


def takeaway_rail(slide, x, y, w, label, hero, body, *, accent=MAGENTA, ink=DEEP, body_c=SLATE):
    """The narrative 'so-what' rail beside a chart (the ~35% right column): small caps accent
    label → one restated hero stat → a 2-3 line interpretation. Pair with a chart in the left
    ~65% (content_band). Every designed chart should carry one of these."""
    text(slide, x, y, w, 0.3, [[(label.upper(), 11, accent, True, False)]], space_after=0)
    text(slide, x, y + 0.34, w, 0.9, [[(hero, 34, ink, True, False)]], space_after=0)
    text(slide, x, y + 1.3, w, 2.0, [[(body, 14, body_c, False, False)]], space_after=0, line_spacing=1.2)


# ============================================ layout patterns (editorial / diagram / wayfinding)
def editorial_header(slide, eyebrow, title, *, x=0.6, y=0.55, w=None, accent=MAGENTA, ink=DEEP,
                     serif=None, size=28, rule_w=1.2):
    """Editorial header lockup: a caps eyebrow, a large title, and a short accent hairline beneath.
    The premium/showcase alternative to title_bar. The title uses the **DISPLAY** face by default
    (falls back to FONT), and EADISPLAY for CJK; pass `serif=` to override per call."""
    if w is None:
        sw, _ = _slide_size(slide); w = sw - 2 * x
    disp = serif if serif is not None else (DISPLAY or FONT)
    text(slide, x, y, w, 0.3, [[(eyebrow.upper(), 12, accent, True, False)]], space_after=0)
    tb = text(slide, x, y + 0.34, w, 0.8, [[(title, size, ink, True, False, disp)]], space_after=0)
    if EADISPLAY:
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                _apply_ea(r, EADISPLAY)
    box(slide, x + 0.02, y + 0.34 + size / 72.0 * 1.18, rule_w, 0.05, fill=accent)


def big_numeral(slide, x, y, n, *, mode="marker", color=MAGENTA, size=None, w=None,
                italic=True, serif="Georgia"):
    """An oversized index figure as wayfinding/rhythm. mode='marker' (solid accent, ~44pt) for a
    numbered item; 'ghost' (very large, near-bg) as a watermark behind a title. The box is sized
    GENEROUSLY WIDE so a short token like '01' / '04' never wraps to two stacked glyphs (the bug
    seen in the Swiss deck — LibreOffice ignores word-wrap=off, so we prevent it by width)."""
    s = size or (44 if mode == "marker" else 132)
    c = color if mode == "marker" else RGBColor(0xE8, 0xE8, 0xE8)
    if w is None:
        sw, _ = _slide_size(slide)
        w = min(len(str(n)) * s / 72.0 * 1.0 + 0.5, sw - x - 0.1)   # wide enough to stay one line, but on-canvas
    tb = text(slide, x, y, w, s / 72.0 * 1.35, [[(str(n), s, c, True, italic, serif)]], space_after=0)
    tb.text_frame.word_wrap = False
    return tb


def stat_row(slide, x, y, w, items, *, ink=DEEP, accent=MAGENTA, serif=None, dividers=True,
             fig_size=34, label_c=MUTE):
    """Editorial big-number row: items = [(figure, unit, label), ...] in 2-4 equal columns with
    optional vertical hairline dividers. For 2-4 standout numbers with no trend to plot."""
    if not items:
        return y
    n = len(items); gap = 0.4; cw = (w - (n - 1) * gap) / n
    for i, item in enumerate(items):
        fig, unit, label = item if len(item) == 3 else (item[0], "", item[1])  # unit is optional
        cx = x + i * (cw + gap)
        runs = [(str(fig), fig_size, ink, True, False, serif)]
        if unit:
            runs.append((" " + str(unit), fig_size * 0.42, accent, True, False, serif))
        text(slide, cx, y, cw, 0.7, [runs], space_after=0)
        text(slide, cx, y + 0.66, cw, 0.4, [[(str(label), 12, label_c, False, False)]], space_after=0)
        if dividers and i > 0:
            box(slide, cx - gap / 2, y + 0.06, 0.014, 0.9, fill=RGBColor(0xDD, 0xDD, 0xDD))
    return y + 1.1


def _set_baseline(run, pct):
    """Raise (or lower, if negative) a run by `pct` percent of its font size via the OOXML baseline
    shift — used to vertically centre a small run beside a much larger one WITHOUT splitting them
    into separate boxes (so natural, equal spacing around the operator is preserved)."""
    run._r.get_or_add_rPr().set("baseline", str(int(round(pct * 1000))))


def change_stat(slide, x, y, w, h, before, after, *, accent=MAGENTA, ink=DEEP, before_size=16,
                after_size=26, arrow="→", font=None, align=None):
    """A 'before → after' change stat with the AFTER value emphasized large. It is **one text box**
    — so the spaces around the arrow are natural and EQUAL on both sides, and it packs into a narrow
    column like a normal line — and the small `before + arrow` run is **baseline-shifted UP** to
    vertically centre on the big AFTER value (mixing sizes in a line otherwise baseline-aligns them,
    sinking the small prefix/arrow below the big number). Returns the textbox."""
    al = align if align is not None else PP_ALIGN.LEFT
    tb = text(slide, x, y, w, h,
              [[(f"{before} {arrow} ", before_size, ink, False, False, font),
                (str(after), after_size, accent, True, False, font)]],
              align=al, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    pct = max(0.0, 0.36 * (after_size - before_size) / max(1, before_size) * 100)  # centre small on big
    runs = tb.text_frame.paragraphs[0].runs
    if runs:
        _set_baseline(runs[0], pct)
    return tb


def quadrant(slide, x, y, w, h, *, x_labels=("", ""), y_labels=("", ""), gap=0.35, axis_c=MUTE):
    """A 2×2 matrix whose AXES carry meaning (e.g. frequency × severity). Draws edge axis captions
    and returns the four cell rects (TL, TR, BL, BR) to fill with cards/scorecards. Use only when
    items truly classify on two independent dimensions; else use a plain grid. Leave ≈1.4in of left
    margin (place at x≈1.5) when using `y_labels`, so they sit in the margin without clipping."""
    cw = (w - gap) / 2; ch = (h - gap) / 2
    if x_labels[0] or x_labels[1]:
        text(slide, x, y - 0.32, cw, 0.28, [[(x_labels[0].upper(), 10.5, axis_c, True, False)]], space_after=0)
        text(slide, x + cw + gap, y - 0.32, cw, 0.28, [[(x_labels[1].upper(), 10.5, axis_c, True, False)]], space_after=0)
    if y_labels[0] or y_labels[1]:
        lx = max(0.05, x - 1.35); lw = max(0.5, x - lx - 0.1)
        text(slide, lx, y + ch * 0.42, lw, 0.3, [[(y_labels[0].upper(), 10.5, axis_c, True, False)]], align=PP_ALIGN.RIGHT, space_after=0)
        text(slide, lx, y + ch + gap + ch * 0.42, lw, 0.3, [[(y_labels[1].upper(), 10.5, axis_c, True, False)]], align=PP_ALIGN.RIGHT, space_after=0)
    return [(x, y, cw, ch), (x + cw + gap, y, cw, ch), (x, y + ch + gap, cw, ch), (x + cw + gap, y + ch + gap, cw, ch)]


def _connector(slide, x0, y0, x1, y1, color, w=1.5, dash=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    if dash:
        ln = c.line._get_or_add_ln()
        ln.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="dash"/>'))
    return c


def hub_spoke(slide, cx, cy, radius, center, spokes, *, hub=(1.7, 1.0), node=(1.8, 0.78),
              accent=BLUE, ink=DEEP, fill=None):
    """Radial 'one core, many peers' diagram: a hub at (cx,cy) with N spoke cards evenly angled on
    a circle of `radius`, connectors drawn behind. `center`='label', spokes=['a','b',...] or
    [(title,sub)]. Use for a platform+modules / metric+drivers; not for sequences (use a pipeline)."""
    n = len(spokes); pts = []
    for i in range(n):
        ang = math.radians(-90 + i * 360.0 / n)
        pts.append((cx + radius * math.cos(ang), cy + radius * math.sin(ang)))
    for (sx, sy) in pts:                       # connectors first (behind the nodes)
        _connector(slide, cx, cy, sx, sy, RGBColor(0xC8, 0xCE, 0xDA), w=1.4, dash=True)
    for (sx, sy), sp in zip(pts, spokes):
        title, sub = (sp if isinstance(sp, (tuple, list)) else (sp, ""))
        chip(slide, sx - node[0] / 2, sy - node[1] / 2, node[0], node[1], title, sub,
             fill if fill is not None else WHITE, tcolor=ink)
    chip(slide, cx - hub[0] / 2, cy - hub[1] / 2, hub[0], hub[1], center, "", accent)
    return pts


def spaced_centers(x, w, n, *, label_w=2.0, total_w=10.0, margin=0.05):
    """Center x's for n evenly-spaced markers across [x, x+w], **inset at the ends** so a
    `label_w`-wide caption CENTERED on each end marker still fits the canvas [margin, total_w-margin].
    Returns (centers, axis_x0, axis_w). Use for ANY row of markers carrying centered captions —
    a timeline, tick row, numbered steps, a stat row under dots. The end inset is what keeps each
    caption **co-centered with its marker**: the failure to avoid is clamping a caption to the margin
    *independently* of its marker, which desyncs them near a slide edge (the classic "first/last label
    sits off to the side of its dot" bug). With this, place each caption at `center - label_w/2` and it
    is guaranteed to land on-canvas AND centered under the marker — no per-caption clamp needed."""
    if n <= 1:
        return ([x + w / 2.0], x, w)
    pad = max(0.0, label_w / 2.0 - (x - margin), label_w / 2.0 - ((total_w - margin) - (x + w)))
    pad = min(pad, w / 2.0 - 0.3)
    x0 = x + pad
    aw = w - 2 * pad
    step = aw / (n - 1)
    return ([x0 + i * step for i in range(n)], x0, aw)


def mid(*vals):
    """Midpoint of the given coordinates — e.g. centre a connector endpoint on a block:
    `connector(s, (ax, mid(by, by+bh)), ...)`, or a hub between two block centres."""
    return sum(vals) / len(vals)


def span_center(boxes, size):
    """Top-left coord that centres a shape of length `size` on the COMBINED SPAN of `boxes`
    (each box = (start, length) on the SAME axis). The one rule for a **converge / fan-out / hub**
    node: a many→one, one→many, or hub-and-spoke node must sit on the geometric centre of the nodes
    it links — never eyeballed to one member's level. Compute it:
        hub_y = span_center([(y_top,h_top), (y_bot,h_bot), ...], hub_h)   # then place the hub at hub_y
    so the hub's centre = (topmost member's top + bottommost member's bottom) / 2. Anchor every
    connector at each member's centre via `mid(y, y+h)`. (Mirrors `spaced_centers` for the across-axis
    case — both exist so diagram nodes are placed by computation, not by eye.)"""
    starts = [b[0] for b in boxes]
    ends = [b[0] + b[1] for b in boxes]
    return (min(starts) + max(ends)) / 2.0 - size / 2.0


def timeline(slide, x, y, w, events, *, orientation="h", highlight=None, accent=MAGENTA,
             ink=DEEP, axis_c=RGBColor(0x9A, 0xA0, 0xAE), h=1.4):
    """Native timeline. events = [(when, title[, caption]), ...]. orientation='h' (axis L→R, 3-6
    evenly-weighted events) or 'v' (top→bottom spine, when each event needs 2+ lines). One node is
    recolored `accent` via `highlight` index. For chronology/roadmaps/evolution — not comparisons.
    End nodes are inset (via `spaced_centers`) so the first/last captions stay centered on their dots."""
    n = len(events)
    if orientation == "h":
        ay = y + 0.2
        sw, _sh = _slide_size(slide)
        def _lx(cx, lw):                                  # keep a centered label box on-canvas
            return max(0.05, min(cx - lw / 2, sw - lw - 0.05))
        centers, ax0, axw = spaced_centers(x, w, n, label_w=2.0, total_w=sw)
        box(slide, ax0, ay - 0.012, axw, 0.024, fill=axis_c)
        for i, ev in enumerate(events):
            ex = centers[i] if n > 1 else (x + w / 2)
            when, title = ev[0], ev[1]; cap = ev[2] if len(ev) > 2 else ""
            em = (highlight is None or i == highlight)
            dc = accent if em else axis_c
            box(slide, ex - 0.09, ay - 0.09, 0.18, 0.18, fill=dc, round=True, r=0.09)
            text(slide, _lx(ex, 2.0), ay + 0.18, 2.0, 0.3, [[(str(when), 13, dc, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
            text(slide, _lx(ex, 2.0), ay + 0.5, 2.0, 0.3, [[(title, 12, ink, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
            if cap:
                text(slide, _lx(ex, 2.2), ay + 0.78, 2.2, 0.5, [[(cap, 10.5, MUTE, False, False)]], align=PP_ALIGN.CENTER, space_after=0)
    else:
        ax = x + 0.12
        box(slide, ax - 0.012, y, 0.024, h, fill=axis_c)
        step = h / max(1, n)
        for i, ev in enumerate(events):
            ey = y + i * step + 0.1
            when, title = ev[0], ev[1]; cap = ev[2] if len(ev) > 2 else ""
            em = (highlight is None or i == highlight); dc = accent if em else axis_c
            box(slide, ax - 0.09, ey - 0.09, 0.18, 0.18, fill=dc, round=True, r=0.09)
            text(slide, ax + 0.35, ey - 0.16, w - 0.5, 0.3, [[(str(when) + "  ", 13, dc, True, False), (title, 13, ink, True, False)]], space_after=0)
            if cap:
                text(slide, ax + 0.35, ey + 0.16, w - 0.5, 0.4, [[(cap, 10.5, MUTE, False, False)]], space_after=0)
    return y + (h if orientation == "v" else 1.4)


def image_tab(slide, x, y, text_str, *, color=DEEP, tcolor=WHITE, size=10.5):
    """A small solid corner label-tab that belongs to a photo (EXTERIOR / BEFORE / 第一). Place at
    the image's corner so the label reads as part of it, not floating beside it."""
    wpt = 0.16 + 0.085 * len(text_str)
    box(slide, x, y, wpt, 0.32, fill=color)
    text(slide, x, y, wpt, 0.32, [[(text_str.upper(), size, tcolor, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def before_after(slide, x, y, w, h, img_a, img_b, label_a="BEFORE", label_b="AFTER", *,
                 accent=MAGENTA, gap=0.7):
    """Two equal image cards with corner label-tabs and a connecting accent arrow — for old↔new /
    v1↔v2 / renovation comparisons (exactly two)."""
    cw = (w - gap) / 2
    picture(slide, img_a, x, y, cw, h, fit="cover"); image_tab(slide, x + 0.12, y + 0.12, label_a, color=accent)
    picture(slide, img_b, x + cw + gap, y, cw, h, fit="cover"); image_tab(slide, x + cw + gap + 0.12, y + 0.12, label_b, color=accent)
    arrow(slide, x + cw + 0.08, y + h / 2 - 0.13, gap - 0.16, 0.26, color=accent)


def photo_triptych(slide, imgs, *, x=0.0, y=1.1, w=None, h=4.4, scrim=False):
    """Three full-height image columns as one hero band (zero gutter). For one subject with several
    strong views / a 'scale & grandeur' message. Optional bottom scrim for an overlaid caption."""
    if w is None:
        sw, _ = _slide_size(slide); w = sw - 2 * x if x else sw
    cw = w / 3
    for i, im in enumerate(imgs[:3]):
        picture(slide, im, x + i * cw, y, cw, h, fit="cover")
    if scrim:
        scrim_overlay(slide, x, y + h - 1.4, w, 1.4, stops=((0.0, 0.0), (1.0, 0.7)), angle=90)


def corner_frame(slide, *, corners=("tl", "br"), color=MAGENTA, length=0.9, weight=0.04, inset=0.5):
    """Decorative L-corner brackets to frame a sparse closing/quote slide so it reads as composed,
    not unfinished. Use on one slide, not throughout."""
    sw, sh = _slide_size(slide)
    for c in corners:
        if c == "tl":   px, py, hx, vy = inset, inset, inset, inset
        elif c == "tr": px, py, hx, vy = sw - inset, inset, sw - inset - length, inset
        elif c == "bl": px, py, hx, vy = inset, sh - inset, inset, sh - inset - length
        else:           px, py, hx, vy = sw - inset, sh - inset, sw - inset - length, sh - inset - length
        box(slide, hx, py - weight / 2, length, weight, fill=color)   # horizontal arm
        box(slide, px - weight / 2, vy, weight, length, fill=color)   # vertical arm


def accent_one(items, featured_idx, accent, neutral=RGBColor(0xC2, 0xC6, 0xD2)):
    """One-accent discipline as a list: return a colour per item with ONLY `featured_idx` in the
    accent and every other item on `neutral`. Pass to chips/cards/series so the eye goes to the
    one thing that matters (the Swiss/data decks' restraint). Featured=None → all neutral."""
    return [accent if (featured_idx is not None and i == featured_idx) else neutral for i in range(len(items))]


# ================================== publication templates · editorial chrome · self-demo · texture
def _set_spc(shape, pts):
    """Letter-spacing (tracking) in points on every run of a textbox — for tracked caps eyebrows."""
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r._r.get_or_add_rPr().set("spc", str(int(pts * 100)))
    return shape


def part_eyebrow(slide, x, y, text_str, *, w=6.0, color=MUTE, font=MONO, size=11, track=1.5):
    """A small TRACKED caps eyebrow in the chrome (usually mono) font — the editorial/technical
    'part label'. Route every kicker/eyebrow through one chrome font for a quiet signature."""
    tb = text(slide, x, y, w, 0.3, [[(text_str.upper(), size, color, True, False, font)]], space_after=0)
    return _set_spc(tb, track)


def page_marker(slide, page, total=None, *, font=MONO, color=MUTE, size=9):
    """A tiny mono page marker at bottom-right ('03 / 14') — chrome, not content."""
    sw, sh = _slide_size(slide)
    label = f"{int(page):02d} / {int(total):02d}" if total else f"{int(page):02d}"
    text(slide, sw - 1.6, sh - 0.42, 1.2, 0.3, [[(label, size, color, True, False, font)]],
         align=PP_ALIGN.RIGHT, space_after=0)


def cover(slide, title, *, issue_label=None, subtitle=None, mode_caption=None, x=0.7, y=None,
          accent=MAGENTA, ink=DEEP, bg=None, display=None, chrome=MONO):
    """A publication-style COVER (issue label + big display title + accent rule + subtitle + a
    date/mode caption) designed to be mirrored by colophon() as a bookend. Stronger than a plain
    title slide for editorial/report/zine decks."""
    sw, sh = _slide_size(slide)
    if bg is not None:
        box(slide, 0, 0, sw, sh, fill=bg)
    yy = (sh / 2 - 1.1) if y is None else y
    if issue_label:
        part_eyebrow(slide, x + 0.02, yy - 0.5, issue_label, color=accent, font=chrome)
    text(slide, x, yy, sw - 2 * x, 1.4, [[(title, 46, ink, True, False, display)]], space_after=2, line_spacing=1.0)
    box(slide, x + 0.02, yy + 1.5, 1.4, 0.06, fill=accent)
    if subtitle:
        text(slide, x, yy + 1.7, sw - 2 * x, 0.5, [[(subtitle, 16, ink, False, False)]], space_after=0)
    if mode_caption:
        part_eyebrow(slide, x + 0.02, sh - 0.7, mode_caption, color=MUTE, font=chrome)


def colophon(slide, tagline, *, credits=None, tooling=None, x=0.7, accent=MAGENTA, ink=DEEP,
             bg=None, display=None, chrome=MONO):
    """A closing COLOPHON mirroring the cover: a payoff tagline + small mono credits/tooling. A
    stronger close than 'Thanks'; the credits slot doubles as a research deck's sources note.
    `credits`/`tooling` may be a string or a list (joined with ' · ')."""
    if isinstance(credits, (list, tuple)):
        credits = " · ".join(map(str, credits))
    if isinstance(tooling, (list, tuple)):
        tooling = " · ".join(map(str, tooling))
    sw, sh = _slide_size(slide)
    if bg is not None:
        box(slide, 0, 0, sw, sh, fill=bg)
    yy = sh / 2 - 0.9
    text(slide, x, yy, sw - 2 * x, 1.4, [[(tagline, 40, ink, True, False, display)]], space_after=2, line_spacing=1.0)
    box(slide, x + 0.02, yy + 1.4, 1.4, 0.06, fill=accent)
    cy = yy + 1.72
    if credits:
        part_eyebrow(slide, x + 0.02, cy, "credits", color=accent, font=chrome)
        text(slide, x, cy + 0.28, sw - 2 * x, 0.6, [[(credits, 12, ink, False, False, chrome)]], space_after=0)
        cy += 0.92
    if tooling:
        part_eyebrow(slide, x + 0.02, cy, "made with", color=accent, font=chrome)
        text(slide, x, cy + 0.28, sw - 2 * x, 0.4, [[(tooling, 11, MUTE, False, False, chrome)]], space_after=0)


def sources_page(slide, sources, *, title="Sources", cols=2, x=0.7, y=1.4, accent=MAGENTA, ink=DEEP, chrome=MONO):
    """Render references as mono numbered columns under an accent header — a research deck's
    colophon / a credible 'where this came from' close."""
    sw, sh = _slide_size(slide)
    part_eyebrow(slide, x, 0.6, title, color=accent, font=chrome, size=13)
    box(slide, x, 1.02, 1.2, 0.05, fill=accent)
    w = (sw - 2 * x - 0.4 * (cols - 1)) / cols
    per = (len(sources) + cols - 1) // cols
    for ci in range(cols):
        cx = x + ci * (w + 0.4); cy = y
        for i in range(ci * per, min((ci + 1) * per, len(sources))):
            text(slide, cx, cy, w, 0.5, [[(f"{i+1:02d}  ", 9, accent, True, False, chrome),
                                          (sources[i], 9, ink, False, False, chrome)]], space_after=0, line_spacing=1.1)
            cy += 0.46


def specimen_card(slide, x, y, w, h, specimen, label, *, accent=MAGENTA, ink=DEEP, featured=False, serif=None):
    """A rule-on-top SPEC CARD with a giant specimen (a glyph 'Aa', a monogram, a number) as the
    hero — for comparing fonts / brands / metrics. The featured card's rule + specimen recolor to
    the accent. A lighter, more Swiss alternative to a boxed icon-card."""
    rc = accent if featured else RGBColor(0xDD, 0xDD, 0xDD)
    sc = accent if featured else ink
    box(slide, x, y, w, 0.06, fill=rc)                                  # top rule
    text(slide, x, y + 0.14, w, 0.3, [[(label.upper(), 11, MUTE, True, False)]], space_after=0)
    text(slide, x, y + 0.5, w, max(0.3, h - 0.6), [[(specimen, 64, sc, True, False, serif)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def wireframe_grid(slide, x, y, w, h, cells, *, cols=4, rows=3, highlight=None, line=None, ink=DEEP, accent=MAGENTA):
    """A SELF-DEMONSTRATING annotated wireframe — labeled outline cells over a cols×rows grid, for
    decks ABOUT layout/design/systems (show the scaffolding). cells = [(label, c0, cspan, r0, rspan)];
    `highlight` recolors one cell. Pair with spec_list() for the 'derived = base × n' math."""
    if cols < 1 or rows < 1:
        raise ValueError("wireframe_grid needs cols >= 1 and rows >= 1")
    ln = line or RGBColor(0xCC, 0xCC, 0xCC)
    cw = w / cols; ch = h / rows
    for i in range(cols + 1): box(slide, x + i * cw - 0.005, y, 0.01, h, fill=ln)
    for j in range(rows + 1): box(slide, x, y + j * ch - 0.005, w, 0.01, fill=ln)
    for k, (label, c0, cspan, r0, rspan) in enumerate(cells):
        cx = x + c0 * cw; cy = y + r0 * ch
        em = (highlight == k)
        box(slide, cx + 0.04, cy + 0.04, cspan * cw - 0.08, rspan * ch - 0.08,
            fill=None, line=accent if em else ink, line_w=2 if em else 1)
        text(slide, cx + 0.12, cy + 0.1, cspan * cw - 0.24, 0.3,
             [[(label.upper(), 10, accent if em else ink, True, False, MONO)]], space_after=0)


def spec_list(slide, x, y, lines, *, font=MONO, color=DEEP, size=12, gap=0.32):
    """Monospace 'derived = base × n' spec lines — pairs with wireframe_grid for a systems deck."""
    cy = y
    for ln in lines:
        text(slide, x, cy, 6.0, 0.3, [[(ln, size, color, False, False, font)]], space_after=0); cy += gap
    return cy


def photo_card(slide, x, y, w, h, *, role="info", accent=MAGENTA, r=0.1, alpha=0.92):
    """A translucent tinted card to hold text ON a photo (keeps the image visible behind). `role`:
    'info' (light), 'primary' (dark), 'accent' (accent tint). Returns the text colour to use on it."""
    fc, tc = {"info": ("FFFFFF", DEEP), "primary": ("141414", WHITE),
              "accent": (_hex(accent), WHITE)}.get(role, ("FFFFFF", DEEP))
    box(slide, x, y, w, h, round=True, r=r, grad=[(0.0, fc, alpha), (1.0, fc, alpha)])
    return tc


def backdrop_motif(slide, *, kind="grid", color=None, spacing=0.6, accent_disc=None,
                   disc_at=None, disc_r=0.7):
    """A FAINT full-bleed texture (grid / graph-paper) + optional accent disc, to bookend a deck on
    its cover and closer as one object. Keep it faint (≈#EEE) so it never fights body content."""
    sw, sh = _slide_size(slide)
    c = color or RGBColor(0xEE, 0xEE, 0xEE)
    n = int(sw / spacing) + 1
    for i in range(n): box(slide, i * spacing - 0.004, 0, 0.008, sh, fill=c)
    for j in range(int(sh / spacing) + 1): box(slide, 0, j * spacing - 0.004, sw, 0.008, fill=c)
    if accent_disc is not None:
        cx, cy = disc_at or (sw - 1.6, 1.4)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - disc_r), Inches(cy - disc_r), Inches(2 * disc_r), Inches(2 * disc_r))
        o.fill.solid(); o.fill.fore_color.rgb = accent_disc; o.line.fill.background(); o.shadow.inherit = False


_ARROW_SHAPE = {"right": MSO_SHAPE.RIGHT_ARROW, "left": MSO_SHAPE.LEFT_ARROW,
                "up": MSO_SHAPE.UP_ARROW, "down": MSO_SHAPE.DOWN_ARROW}

def arrow(slide, x, y, w, h, color=BLUE, direction="right"):
    """A solid block arrow. **Point it in the direction the flow actually moves.**
    `direction` is "right" (default), "left", "up", or "down". For a vertical flow —
    e.g. two stacked boxes (problem on top, solution below) — use "down" (or "up"), NOT a
    sideways arrow: a left/right arrow between vertically-stacked blocks reads as wrong.
    For an up/down arrow give it a tall, narrow box (small w, larger h); for left/right a
    wide, short box. The arrow fills the (w, h) box you pass."""
    shape = _ARROW_SHAPE.get(direction.lower(), MSO_SHAPE.RIGHT_ARROW)
    a = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    try: a.adjustments[0] = 0.55; a.adjustments[1] = 0.55
    except Exception: pass
    return a


def arrow_label(slide, x, y, w, h, label, color=BLUE, size=9, lab_c=None, box_w=1.3, bold=True,
                direction="right"):
    """An arrow with its label snug against it — so connector labels (a verb, a transform
    name, a step — e.g. 'encode', 'train', 'merge', 'step 2') stay tight to the arrow rather
    than drifting. **Pass `direction` to match the flow** (right/left/up/down, same as
    `arrow`): for a horizontal arrow the label sits centred just above it; for a **vertical**
    (up/down) arrow it sits centred just to the **right** of it (placing it above a vertical
    arrow would read as belonging to the box above). `box_w` is the (transparent) label-box
    width; shrink it when the arrow sits in a narrow gap. Returns the arrow."""
    lab_c = lab_c or color
    th = size / 72.0 * 1.3
    if direction.lower() in ("up", "down"):
        text(slide, x + w + 0.06, y + h / 2.0 - th / 2.0, box_w, th + 0.04,
             [[(label, size, lab_c, bold, False)]], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    else:
        text(slide, x + w / 2.0 - box_w / 2.0, y - th - 0.04, box_w, th + 0.04,
             [[(label, size, lab_c, bold, False)]], align=PP_ALIGN.CENTER, space_after=0)
    return arrow(slide, x, y, w, h, color=color, direction=direction)


# ====================================================================== images
def _round_pic_geom(pic, radius_in, w_in, h_in):
    """Give a placed picture rounded corners (so an image matches a rounded card/frame around it).
    `radius_in` is the corner radius in inches, applied to the picture's ACTUAL placed size."""
    ss = max(0.01, min(w_in, h_in))
    adj = int(max(0.0, min(0.5, radius_in / ss)) * 100000)
    g = pic._element.spPr.find(qn('a:prstGeom'))
    if g is None:
        return
    g.set('prst', 'roundRect')
    av = g.find(qn('a:avLst'))
    if av is None:
        av = g.makeelement(qn('a:avLst'), {}); g.append(av)
    for gd in list(av):
        av.remove(gd)
    av.append(av.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': f'val {adj}'}))


def picture(slide, path, x, y, w, h, fit="contain", alt=None, round=False, r=None):
    """Place an image in a frame without distorting it.

    `fit="contain"` shows the whole image inside the frame, letterboxed by whitespace.
    Use it for source figures, charts, screenshots, and anything whose edges/labels matter.
    `fit="cover"` fills the frame and crops evenly from the long dimension. Use it for
    decorative plates, photo panels, and generated atmosphere where edge crop is acceptable.

    `round=True` (or `r=<inches>`) gives the image ROUNDED corners — match this to the deck's
    cards/panels so a square photo doesn't sit among rounded blocks (a consistency tell). For an
    image inside a rounded frame, use a radius ≈ the frame's radius minus the border so the curves
    stay concentric. Default radius is 8% of the image's shorter side.

    Pass `alt` for informative images; pass `alt=""` for decorative plates. Returns the
    picture shape. Requires Pillow for reliable aspect-ratio reads, matching the rest of
    deckkit's image/equation helpers.
    """
    import os
    from PIL import Image
    if not os.path.exists(path):
        raise FileNotFoundError(f"picture(): image not found: {path} — generate it first, or fix the path")
    with Image.open(path) as im:
        iw, ih = im.size
    if iw <= 0 or ih <= 0:
        raise ValueError(f"cannot read image dimensions for {path}")

    img_ar = iw / ih
    frame_ar = w / h
    fit = fit.lower()
    if fit == "contain":
        if frame_ar > img_ar:
            ph = h
            pw = h * img_ar
            px = x + (w - pw) / 2
            py = y
        else:
            pw = w
            ph = w / img_ar
            px = x
            py = y + (h - ph) / 2
        pic = slide.shapes.add_picture(path, Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))
    elif fit == "cover":
        pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        if img_ar > frame_ar:
            crop = (1.0 - frame_ar / img_ar) / 2.0
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ar < frame_ar:
            crop = (1.0 - img_ar / frame_ar) / 2.0
            pic.crop_top = crop
            pic.crop_bottom = crop
    else:
        raise ValueError("fit must be 'contain' or 'cover'")

    if round or r is not None:
        pwp, php = (pw, ph) if fit == "contain" else (w, h)
        _round_pic_geom(pic, r if r is not None else 0.08 * min(pwp, php), pwp, php)
    if alt is not None:
        alt_text(pic, alt)
    return pic


def icon(slide, png, x, y, size, *, alt=None, disc=None, disc_pad=None):
    """Place a square icon PNG (use `scripts/icons.py`'s `icon_png()` to fetch+recolor+rasterize
    an open-licensed SVG first). `size` is the icon's edge in inches — keep it small and consistent
    (≈0.32–0.5 in; an icon should read at ~heading size, never bigger than the title).

    `disc=<hex>` draws a soft rounded **tinted tile** behind the icon (a common, tidy treatment —
    an accent-tinted square with the icon centred); `disc_pad` is the inset of the icon inside the
    tile. Pass `alt` for a screen-reader label (icons are decorative-ish, but if it carries meaning
    give it alt text AND a visible text label — never rely on the icon alone). Returns the picture
    shape. The icon must sit in ONE consistent family/size/colour across the deck (see icons.md)."""
    if disc:
        if isinstance(disc, str):
            disc = disc.lstrip("#")
        pad = disc_pad if disc_pad is not None else 0.22 * size   # icon ≈56% of the tile
        box(slide, x, y, size, size, fill=disc, round=True, r=0.11 * size)
        return picture(slide, png, x + pad, y + pad, size - 2 * pad,
                       size - 2 * pad, fit="contain", alt=alt)
    return picture(slide, png, x, y, size, size, fit="contain", alt=alt)


def icon_card(slide, x, y, w, h, png, title, body="", *, fill=None, line=None,
              icon_size=0.42, accent=BLUE, ink=DEEP, body_c=SLATE, disc=None, pad=0.26):
    """A card with the **icon in the UPPER-LEFT corner**, then a title and optional body below —
    the clean "feature card" pattern. Use a row of these (built from one `columns(n)` grid) to label
    categories/features/sections; keep the SAME icon family, size, and colour across every card so
    they read as a system (CRAP Repetition).

    `png` is a placed icon PNG (recolour it to `accent` via icons.py first); `disc=<hex>` puts the
    icon in a tinted tile. `fill`/`line` style the card (defaults to a hairline-bordered light card).
    Title sits at ~heading size, body smaller — never let the icon exceed the title. Returns the
    card's bottom y."""
    box(slide, x, y, w, h, fill=fill if fill is not None else WHITE,
        line=line if line is not None else RGBColor(0xE3, 0xE8, 0xEE), line_w=1.0, round=True)
    icon(slide, png, x + pad, y + pad, icon_size, disc=disc)
    ty = y + pad + icon_size + 0.16
    text(slide, x + pad, ty, w - 2 * pad, 0.34,
         [[(title, 15, ink, True, False, DISPLAY or FONT)]], space_after=0)
    if body:
        text(slide, x + pad, ty + 0.34, w - 2 * pad, h - (ty + 0.34 - y) - pad * 0.5,
             [[(body, 12, body_c, False, False)]], space_after=0, line_spacing=1.02)
    return y + h


def icon_tile(slide, x, y, size, png, *, shape="circle", fill=None, grad=None, grad_angle=120.0,
              ring=None, ring_w=1.6, glass=False, sheen=False, pad=None, alt=None):
    """Place an icon inside a STYLED TILE — the versatile alternative to a bare monochrome drop.
    The same recolored icon reads very differently by container, so vary the treatment to fit the
    deck instead of always using a flat glyph (see references/icons.md "treatments"):

      shape : "circle" (disc), "squircle" (rounded square), or "square".
      fill  : a solid tile colour (hex/RGBColor). Omit for no fill (pair with `ring`/`glass`).
      grad  : a two-stop gradient tile — `(c0, c1)` or a full `[(pos,colour,alpha),…]` list —
              overriding `fill`. Gives depth a flat fill lacks (the glassy header-disc look);
              `grad_angle` is the linear direction in degrees (0=→, 90=↓).
      ring  : an accent OUTLINE colour → a badge (icon inside a thin coloured ring).
      glass : True → a frosted translucent tile (low-alpha tint of `fill` + white rim) for
              dark / glowing / photographic backgrounds (pair with `glow()`).
      sheen : True → a soft top highlight (the glassy edge in modern decks).

    `size` is the tile edge in inches; the icon is inset by `pad` (default ≈26% so the glyph sits
    at ~50-55% of the tile — the tidy proportion). For a ROW of icons keep size/shape/treatment
    IDENTICAL across siblings (CRAP Repetition); colour-code per category by varying only the hue
    (its tile + its label), as polished decks do. Returns the icon picture shape."""
    if isinstance(grad, str):
        grad = None
    sh_kind = {"circle": MSO_SHAPE.OVAL, "squircle": MSO_SHAPE.ROUNDED_RECTANGLE,
               "square": MSO_SHAPE.RECTANGLE}.get(shape, MSO_SHAPE.OVAL)
    t = slide.shapes.add_shape(sh_kind, Inches(x), Inches(y), Inches(size), Inches(size))
    if glass:
        tint = _as_rgb(fill) if fill is not None else WHITE
        _grad_fill(t, [(0.0, tint, 0.20), (1.0, tint, 0.06)], angle=grad_angle)
    elif grad is not None:
        stops = grad if isinstance(grad[0], (tuple, list)) \
            else [(0.0, _as_rgb(grad[0]), 1.0), (1.0, _as_rgb(grad[1]), 1.0)]
        _grad_fill(t, stops, angle=grad_angle)
    elif fill is not None:
        t.fill.solid(); t.fill.fore_color.rgb = _as_rgb(fill)
    else:
        t.fill.background()
    if ring is not None:
        t.line.color.rgb = _as_rgb(ring); t.line.width = Pt(ring_w)
    elif glass:
        t.line.color.rgb = WHITE; t.line.width = Pt(1.0)
    else:
        t.line.fill.background()
    t.shadow.inherit = False
    if shape == "squircle":
        try: t.adjustments[0] = 0.24
        except Exception: pass
    if sheen:
        # a faint white highlight across the top — sells the glassy edge
        hh = size * 0.5
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL if shape == "circle" else MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + size * 0.12), Inches(y + size * 0.06), Inches(size * 0.76), Inches(hh))
        _grad_fill(sh, [(0.0, WHITE, 0.40), (1.0, WHITE, 0.0)], angle=90.0)
        sh.line.fill.background(); sh.shadow.inherit = False
    pad = pad if pad is not None else 0.26 * size
    return picture(slide, png, x + pad, y + pad, size - 2 * pad, size - 2 * pad,
                   fit="contain", alt=alt)


def icon_badge(slide, x, y, size, png, *, ring=MAGENTA, ring_w=1.8, fill=None, alt=None):
    """An icon inside a thin accent RING — a light, outlined treatment that reads as a badge
    (good on a light deck where a solid tile would feel heavy). Thin wrapper over `icon_tile`."""
    return icon_tile(slide, x, y, size, png, shape="circle", fill=fill, ring=ring,
                     ring_w=ring_w, alt=alt)


def icon_ghost(slide, png, x, y, size, *, alt=""):
    """An OVERSIZED, FAINT icon used as a watermark behind a card's content (the ghost-glyph
    treatment that adds texture without clutter). Place it FIRST, then draw text/blocks on top.
    Recolor the PNG to a very light tint (icons.py `color=<pale hue>`) so it never competes with
    the text — this helper only sizes/places it big; the faintness comes from the recolor. Pass a
    big `size` (it may bleed past the card edge for effect). Decorative → `alt=""` by default."""
    return picture(slide, png, x, y, size, size, fit="contain", alt=alt)


def cjk_numeral(n, style="formal"):
    """CJK numeral string for n (1–99) — for East-Asian section markers (壹·贰·叁 …).
    style='formal' (壹贰叁肆伍陆柒捌玖拾, 大写) or 'simple' (一二三…). Use as a numeral marker on an
    ink/traditional CJK deck instead of Latin "01/02" (see references/east-asian-aesthetic.md)."""
    d = "〇一二三四五六七八九" if style == "simple" else "〇壹贰叁肆伍陆柒捌玖"
    ten = "十" if style == "simple" else "拾"
    n = int(n)
    if n < 0 or n > 99:
        return str(n)
    if n < 10:
        return d[n]
    if n < 20:
        return ten + (d[n - 10] if n > 10 else "")
    return d[n // 10] + ten + (d[n % 10] if n % 10 else "")


def seal(slide, x, y, size, char, *, fill=None, tcolor=None, shape="square",
         rounded=True, border=True):
    """A red SEAL / chop stamp with a 1–2 char CJK mark — the signature East-Asian accent
    (印章). Filled vermilion square with a light char by default (阴文); on a warm-paper ink
    deck it's the single spot of red. `shape="circle"` for a round seal; `border=True` adds the
    classic thin inner rule. Pass `fill` to override the seal colour, `tcolor` the char colour.
    Keep it SMALL (≈0.45–0.8 in) and use ONE per slide — it's a signature, not a sticker.
    Set deckkit.EADISPLAY (e.g. 'KaiTi') so the mark renders in a brush face, not tofu."""
    f = fill if fill is not None else RGBColor(0xA5, 0x2A, 0x2A)
    tc = tcolor if tcolor is not None else RGBColor(0xF5, 0xF1, 0xE8)
    if shape == "circle":
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background()
    else:
        box(slide, x, y, size, size, fill=f, round=rounded, r=0.1 * size)
    if border:
        inset = 0.1 * size
        if shape == "circle":
            b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + inset), Inches(y + inset),
                                       Inches(size - 2 * inset), Inches(size - 2 * inset))
            b.fill.background(); b.line.color.rgb = tc; b.line.width = Pt(1.0)
        else:
            box(slide, x + inset, y + inset, size - 2 * inset, size - 2 * inset,
                fill=None, line=tc, line_w=1.0, round=rounded, r=0.08 * size)
    fs = max(10, int(size * 72 * (0.42 if len(char) > 1 else 0.5)))
    fam = EADISPLAY or EAFONT or DISPLAY or FONT
    text(slide, x, y, size, size, [[(char, fs, tc, True, False, fam)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.9)
    return x + size


def gif_poster(path, out_png=None, frame="first"):
    """Extract ONE representative frame of an animated GIF to a PNG — for the render/critic
    (which see only a static image) and as a check on what the deck shows when NOT in slideshow.

    Why this matters: a `.gif` placed on a slide loops only in **slideshow**. In the editor,
    in a PDF/print export, in the LibreOffice render, and to the static critic, the GIF shows
    its **first frame**. A cine / 4D / training-run GIF often starts on a blank, black, or
    "loading" frame — so the slide looks broken everywhere except live playback. Use this to
    SEE that frame and verify it reads:

      `frame="first"` (default) — frame 0, i.e. exactly what the static views show. View it; if
        it's blank/unrepresentative, get a source GIF that *starts* on a meaningful frame (there
        is no separate poster-frame for GIFs in pptx — the first frame IS the poster).
      `frame="auto"` — the frame with the most visual content (highest pixel variance), i.e. a
        good representative still to hand the critic for a legibility judgement.
      `frame="middle"` or an int index — a specific frame.

    Returns the output PNG path (defaults to `<gif>.<frame>.png`). Pillow-based; raises if the
    file isn't a readable (animated) image.
    """
    import os
    from PIL import Image, ImageSequence
    if not os.path.exists(path):
        raise FileNotFoundError(f"gif_poster(): file not found: {path}")
    frames = []
    with Image.open(path) as im:
        for fr in ImageSequence.Iterator(im):
            frames.append(fr.convert("RGB").copy())
    if not frames:
        raise ValueError(f"gif_poster(): no frames read from {path}")
    if frame == "first":
        idx = 0
    elif frame == "middle":
        idx = len(frames) // 2
    elif frame == "auto":
        # most visual content = largest per-channel variance (skips blank/loading frames)
        def score(fr):
            stat = fr.resize((64, 64))
            px = list(stat.getdata())
            n = len(px)
            mean = [sum(c) / n for c in zip(*px)]
            return sum((c - mean[i]) ** 2 for p in px for i, c in enumerate(p))
        idx = max(range(len(frames)), key=lambda i: score(frames[i]))
    else:
        idx = int(frame)
    idx = max(0, min(idx, len(frames) - 1))
    if out_png is None:
        base, _ = os.path.splitext(path)
        out_png = f"{base}.{frame if isinstance(frame, str) else 'f%d' % idx}.png"
    frames[idx].save(out_png)
    return out_png


def gif(slide, path, x, y, w, h, *, fit="contain", alt=None, max_mb=8.0, warn=True):
    """Place an **animated GIF** as live content — it stays animated and loops in PowerPoint /
    Keynote slideshow. Use for any result whose *motion is the point*: a cine / 4D / time-resolved
    sequence, a segmentation-over-time, an optimisation/training trajectory, a rotating 3D model,
    a UI/interaction loop, a physics sim.

    This is a thin, GIF-aware wrapper over `picture()`: it places the GIF **whole and undistorted**
    (`fit="contain"` — preserves the real aspect, never stretches a square cine clip to 16:9; the
    original file bytes are embedded so EVERY frame survives), sets alt-text, and adds two checks
    that catch the real failure modes:
      • **size/perf** — warns if the file exceeds `max_mb` (a heavy cine GIF bloats the .pptx and
        can stutter live; palette-optimise / downsample fps or dimensions if so);
      • **animation** — warns if the file is a *single-frame* GIF (then it's just a still — use
        `picture()`), and reports the frame count.
    It does NOT rasterise the GIF (that would freeze it). Pair with `gif_poster()` to verify the
    first frame (what the render / a PDF export shows) is representative, and label it with a deck-
    font caption + a one-line "what to watch" like any figure. Returns the picture shape.
    """
    import os
    from PIL import Image, ImageSequence
    if not os.path.exists(path):
        raise FileNotFoundError(f"gif(): file not found: {path} — fix the path or generate it first")
    n_frames = 1
    try:
        with Image.open(path) as im:
            n_frames = getattr(im, "n_frames", 1)
    except Exception:
        pass
    if warn:
        mb = os.path.getsize(path) / 1e6
        if mb > max_mb:
            print(f"gif(): WARNING {os.path.basename(path)} is {mb:.1f} MB (> {max_mb} MB) — "
                  f"palette-optimise or downsample fps/size so the .pptx stays light and plays smoothly.")
        if n_frames <= 1:
            print(f"gif(): WARNING {os.path.basename(path)} has 1 frame — it's a still, not an "
                  f"animation; use picture() instead.")
    if fit == "cover":
        # cover crops each frame's edges — rarely right for a scientific GIF; allow but flag.
        if warn:
            print("gif(): note fit='cover' crops the GIF's edges every frame — use 'contain' unless "
                  "the GIF is edge-tolerant texture.")
    return picture(slide, path, x, y, w, h, fit=fit, alt=alt)


def columns(n=2, *, slide=None, w_in=None, h_in=None, top=1.15, bottom=0.55, margin=None, gap=None):
    """Return ``n`` equal-width content-column rects ``(x, y, w, h)`` with **symmetric**
    outer margins and equal gutters between columns.

    Use this for any split slide — text+figure, two-up, three-up, image+caption — so the
    left and right regions (and the white margins flanking them) come out the SAME width.
    The most common lopsided-slide tell is a left panel and right panel of different
    widths, or a wider white margin on one side than the other; that happens when x/w are
    eyeballed per panel. Deriving every panel from one grid makes the layout balanced by
    construction:

        L, R = dk.columns(2, slide=s)         # two equal halves, equal flanking margins
        dk.bullet(s, *L[:3], items)           # text in the left half
        dk.picture(s, fig, *R, fit="contain") # figure in the right half (same width)

    **Pass ``slide=`` so the grid matches the deck's REAL size** — the width/height are then
    read from the presentation, so it stays symmetric on a 16:9, a widescreen 13.333×7.5, or
    a poster. Only the no-arg call falls back to the standard 10×5.625 deck; if you call it
    without ``slide=`` on a differently-sized deck, pass the same ``w_in``/``h_in`` you gave
    ``blank_deck`` (otherwise the right margin will be wrong — the very lopsidedness this
    helper exists to prevent).

    ``margin`` (outer left == outer right) and ``gap`` (between columns) default to
    ``GUTTER``. ``top``/``bottom`` reserve room for the title bar and footer. Returns a
    list of ``(x, y, w, h)`` tuples in inches, left to right.
    """
    if slide is not None:
        prs = slide.part.package.presentation_part.presentation
        if w_in is None:
            w_in = prs.slide_width / 914400
        if h_in is None:
            h_in = prs.slide_height / 914400
    w_in = 10.0 if w_in is None else w_in
    h_in = 5.625 if h_in is None else h_in
    margin = GUTTER if margin is None else margin
    gap = GUTTER if gap is None else gap
    if n < 1:
        raise ValueError("columns(n) needs n >= 1")
    usable = w_in - 2 * margin - (n - 1) * gap
    if usable <= 0:
        raise ValueError("margins/gap leave no room for columns; reduce margin or gap")
    cw = usable / n
    y = top
    h = h_in - top - bottom
    return [(margin + i * (cw + gap), y, cw, h) for i in range(n)]


def rows(n=2, *, slide=None, w_in=None, h_in=None, x=None, y=None, w=None, top=1.15, bottom=0.55, gap=None):
    """Return ``n`` equal-**height** row rects ``(x, y, w, h)`` stacked top-to-bottom with
    equal gaps — the vertical counterpart of :func:`columns`.

    Use it for a **vertical stack** of boxes (problem→solution, before→after, a list of
    cards) so the gaps — and any ``arrow(..., direction="down")`` connectors you drop between
    them — are **equal by construction** rather than eyeballed. Pass ``x``/``w`` to confine the
    stack to one column (e.g. the left half from ``columns``); they default to a full-width
    band inset by ``GUTTER`` (derived from the deck's real width when ``slide=`` is given)::

        L, R = dk.columns(2, slide=s)
        top_box, bot_box = dk.rows(2, slide=s, x=L[0], w=L[2], top=1.5, bottom=1.3)
        dk.chip(s, *top_box, "abandon", "used once", STEEL)
        gap_mid = bot_box[1] - (top_box[1] + top_box[3])          # the equal gap between rows
        dk.arrow(s, top_box[0] + top_box[2] / 2 - 0.11, top_box[1] + top_box[3] + (gap_mid - 0.3) / 2,
                 0.22, 0.3, direction="down")                     # centred in that gap
        dk.chip(s, *bot_box, "reuse", "flies again", ORANGE)

    Pass ``slide=`` so it matches the deck's real size. Returns ``(x, y, w, h)`` tuples,
    top to bottom.
    """
    if slide is not None:
        prs = slide.part.package.presentation_part.presentation
        if w_in is None:
            w_in = prs.slide_width / 914400
        if h_in is None:
            h_in = prs.slide_height / 914400
    w_in = 10.0 if w_in is None else w_in
    h_in = 5.625 if h_in is None else h_in
    gap = GUTTER if gap is None else gap
    x = GUTTER if x is None else x
    w = (w_in - 2 * GUTTER) if w is None else w
    y = top
    if n < 1:
        raise ValueError("rows(n) needs n >= 1")
    usable = h_in - top - bottom - (n - 1) * gap
    if usable <= 0:
        raise ValueError("top/bottom/gap leave no room for rows; reduce them")
    rh = usable / n
    return [(x, y + i * (rh + gap), w, rh) for i in range(n)]


def content_band(slide, *, top=1.15, footer_gap=0.15):
    """The one authoritative SAFE content rect ``(x, y, w, h)`` — below the title bar and
    **above the footer band** — read from the deck's REAL size.

    Use it instead of hand-picking 'somewhere above the footer' y-coordinates (the magic
    numbers like 4.3 / 5.05 that drift and let auto-growing blocks collide with the footer).
    The bottom edge is ``h_in - FOOTER_BAND - footer_gap``, so anything placed within the
    returned rect clears ``footer()``. Pair with :func:`vstack` / :func:`bottom_callout`.
    """
    w_in, h_in = _slide_size(slide)
    y = top
    bottom_y = h_in - FOOTER_BAND - footer_gap
    return (GUTTER, y, w_in - 2 * GUTTER, bottom_y - y)


def vstack(slide, x, y, w, blocks, *, gap=GUTTER, bottom=None, anchor="top"):
    """Measured vertical packer — **equal gaps and no overlap, guaranteed by construction**.

    Unlike :func:`rows` (which returns equal *fixed*-height cells), ``vstack`` respects each
    block's CONTENT-driven height, so a callout/bullet-list/chip that auto-grows can't overflow
    its cell and collide with the block below (the recurring "raised callout overlapped the
    bullets above it" failure).

    ``blocks`` = list of ``(height_in, draw)`` where ``draw(x, y, w)`` renders the block at that
    top-left and width. Measure each height first with :func:`measure_callout` /
    :func:`measure_bullets` / :func:`measure_text` (or a figure's known aspect).

    With ``bottom`` given (e.g. from ``content_band``): raises a located error if the blocks
    can't fit (so a collision surfaces at BUILD time, not at render); ``anchor='center'`` centres
    the stack in the band and ``anchor='justify'`` spreads equal gaps to fill it — both kill the
    "too much dead-white on the bottom" tell. Returns the placed ``(x, y, w, h)`` rects.
    """
    n = len(blocks)
    if n == 0:
        return []
    heights = [h for h, _ in blocks]
    content = sum(heights)
    total = content + gap * (n - 1)
    if bottom is not None:
        avail = bottom - y
        if total > avail + 1e-6:
            raise ValueError(
                f"vstack overflow: blocks need {total:.2f}in but the band is only {avail:.2f}in "
                f"— shorten text, drop a block, or raise the top. (Surfaced at build time on purpose.)")
        if anchor == "center":
            y += (avail - total) / 2
        elif anchor == "justify" and n > 1:
            gap = (avail - content) / (n - 1)
    cy = y
    rects = []
    for h, draw in blocks:
        draw(x, cy, w)
        rects.append((x, cy, w, h))
        cy += h + gap
    return rects


# ================================================================= components
def _is_wide(o):
    """True for CJK / full-width code points — their glyph advance is ≈ one em."""
    return (0x1100 <= o <= 0x115F or 0x2E80 <= o <= 0xA4CF or 0xAC00 <= o <= 0xD7A3
            or 0xF900 <= o <= 0xFAFF or 0xFE30 <= o <= 0xFE4F or 0xFF00 <= o <= 0xFF60
            or 0xFFE0 <= o <= 0xFFE6 or 0x20000 <= o <= 0x3FFFD)


def _disp_len(s):
    """Display width in 'Latin-char' units: CJK / full-width glyphs count as 2. Used only by
    the heuristic FALLBACK in `_measure_lines`; the primary path measures real glyph
    advances. Counting wide glyphs as 2 keeps that fallback safe for CJK text."""
    return sum(2 if _is_wide(ord(ch)) else 1 for ch in s)


# Chars-per-line estimates over-count narrow glyphs (ASCII digits/spaces/dashes, CJK
# punctuation), so an item one char over the wrap threshold gets a phantom extra line. This
# slack absorbs that bias. It is now only used by the heuristic FALLBACK below — the primary
# path (`_measure_lines`) measures real glyph advances and needs no slack.
_WRAP_SLACK = 2

# ---- accurate text measurement (real glyph metrics, with a heuristic fallback) ----
_MEAS_PREC = 4                 # load fonts at size*PREC px for sub-point precision
_FONT_PATH_CACHE = {}
_PIL_FONT_CACHE = {}


def _font_file(name, bold=False):
    """Resolve a font family name to a file path (cached). Returns None if unresolvable."""
    key = (name, bold)
    if key not in _FONT_PATH_CACHE:
        try:
            from matplotlib import font_manager as _fm
            _FONT_PATH_CACHE[key] = _fm.findfont(
                _fm.FontProperties(family=name, weight=("bold" if bold else "normal")))
        except Exception:
            _FONT_PATH_CACHE[key] = None
    return _FONT_PATH_CACHE[key]


def _pil_font(name, size_pt, bold=False):
    """A cached Pillow font for `name` at `size_pt` (loaded at size*PREC px)."""
    key = (name, bold, round(size_pt, 2))
    f = _PIL_FONT_CACHE.get(key)
    if f is None:
        from PIL import ImageFont
        f = ImageFont.truetype(_font_file(name, bold), max(1, int(round(size_pt * _MEAS_PREC))))
        _PIL_FONT_CACHE[key] = f
    return f


def _lines_heuristic(text, size_pt, avail_in):
    """Chars-per-line line-count estimate — the fallback when measurement isn't available."""
    cpl = max(6, int(avail_in * 136.0 / size_pt))
    eff = max(1, _disp_len(text) - _WRAP_SLACK)
    return max(1, -(-eff // cpl))


def _measure_lines(runs, size_pt, avail_in, font=None):
    """How many lines styled text wraps to — MEASURED, not estimated.

    `runs` = [(text, bold), ...] set at `size_pt` within `avail_in` inches of usable width.
    Narrow runs are measured with the REAL Latin font's glyph advances (Pillow, the bold
    parts measured bold); CJK / full-width glyphs are one em (= size_pt) by definition.
    A greedy line-breaker then counts wraps, breaking at spaces, between CJK glyphs, and at
    CJK↔Latin boundaries (Latin words stay whole). Because it uses the same font metrics the
    renderer does, the count matches the rendered layout far more closely than a chars-per-
    line guess. Falls back to `_lines_heuristic` if Pillow or the font can't be loaded — so a
    build never breaks over measurement. Lazy-imports Pillow/matplotlib."""
    fontname = font or FONT
    flat = "".join(t for t, _ in runs)
    if not flat:
        return 1
    avail = max(1.0, avail_in * 72.0)                       # usable width, in points
    try:
        fonts = {b: _pil_font(fontname, size_pt, b) for b in (False, True)}
        getlen = lambda s, b: fonts[b].getlength(s) / _MEAS_PREC
        getlen("x", False)                                  # probe — raises if unusable
    except Exception:
        return _lines_heuristic(flat, size_pt, avail_in)

    items = []                                              # (width_pt, kind): 'w'ord 's'pace 'c'jk
    for text, bold in runs:
        word = []
        for ch in text:
            if ch == " ":
                if word:
                    items.append((getlen("".join(word), bold), "w")); word = []
                items.append((getlen(" ", bold), "s"))
            elif _is_wide(ord(ch)):
                if word:
                    items.append((getlen("".join(word), bold), "w")); word = []
                items.append((float(size_pt), "c"))
            else:
                word.append(ch)
        if word:
            items.append((getlen("".join(word), bold), "w"))

    x = 0.0
    lines = 1
    for w, kind in items:
        if kind == "s":                                     # a space never forces a wrap
            if x > 0:
                x += w
            continue
        if w > avail:                                       # token alone wider than the line
            if x > 0:
                lines += 1
            n = int(w // avail)
            lines += n
            x = w - n * avail
            continue
        if x + w > avail and x > 0:
            lines += 1
            x = 0.0
        x += w
    return max(1, lines)


# ---- public "measure before you place" helpers (so a build knows a block's true height
#      at a given width BEFORE choosing its y — measure-then-place, not place-and-pray) ----
def measure_lines(runs, size_pt, avail_in, font=None):
    """Public wrapper: how many lines ``runs`` = [(text, bold), ...] wrap to at ``size_pt``
    within ``avail_in`` inches. Same metric the renderer/`bullet`/`callout` use."""
    return _measure_lines(runs, size_pt, avail_in, font=font)


def measure_callout(label, body, w):
    """Height (inches) :func:`callout` will draw for this ``label``+``body`` at width ``w``.
    Measure it BEFORE placing so the box can be positioned to clear the footer / the block
    below — the single source of truth for the callout height formula."""
    nlines = _measure_lines([(label + "  ", True), (body, False)], 12.5, w - 0.44)
    return 0.30 + 0.245 * nlines   # 0.30 = top+bottom padding: snug to the text but not cramped


def measure_bullets(items, w, size=17, gap=0.26):
    """Height (inches) :func:`bullet` will occupy for ``items`` at width ``w`` (no trailing
    gap), using the same per-item line measurement — so a build can place the next block
    below the list (or hand the list to :func:`vstack`) without overlap."""
    line_h = size / 72.0 * 1.12
    total = 0.0
    for i, (lead, rest) in enumerate(items):
        nlines = _measure_lines([(lead, True), (rest, False)], size, w - 0.22)
        total += line_h * nlines
        if i < len(items) - 1:
            total += gap
    return total


def measure_text(runs, w, size, *, line_h_factor=1.12, pad=0.0):
    """Height (inches) a plain :func:`text` block of ``runs`` = [(text, bold), ...] needs at
    ``size`` within width ``w``. ``pad`` adds top+bottom slack."""
    nlines = _measure_lines(runs, size, w)
    return nlines * (size / 72.0 * line_h_factor) + pad


def bullet(slide, x, y, w, items, size=17, gap=0.26, marker=BLUE, lead_c=DEEP, body_c=SLATE):
    """Square-marker bullets. items = list of (lead, rest); keep both terse.
    Returns the bottom y, so a caller can place the next element (e.g. a callout)
    below the list without overlapping it.

    EVEN RHYTHM depends on every item occupying the SAME line count. Each marker is placed
    by advancing the cursor `line_h * nlines + gap`, where `nlines` is now MEASURED from the
    real font's glyph advances (`_measure_lines` — Latin runs via Pillow, the bold lead
    measured bold, CJK as one em), so it matches what the renderer actually lays out. This
    removes the old heuristic's two failure modes near the wrap boundary — undershoot
    (next bullet OVERLAPS the wrapped text) and overshoot (a PHANTOM blank line is reserved,
    so the gap before the next bullet looks uneven). If Pillow/the font can't load, it falls
    back to the chars-per-line heuristic. STILL TRUE regardless: an item that genuinely wraps
    to 2 lines takes more vertical space than its 1-line peers, so for an even-looking list
    keep items to a CONSISTENT line count (ideally each comfortably on one line). As always,
    verify the PNG. (At the default size=17 this matches the prior behaviour.)"""
    cy = y
    line_h = size / 72.0 * 1.12              # line height in inches, scaled to font size
    for lead, rest in items:
        box(slide, x, cy + 0.07, 0.09, 0.09, fill=marker)
        text(slide, x + 0.22, cy - 0.02, w - 0.22, 0.6,
             [[(lead, size, lead_c, True, False), (rest, size, body_c, False, False)]],
             space_after=0, line_spacing=1.02)
        # MEASURED line count (real glyph metrics; bold lead measured bold) so the marker
        # advance matches the renderer's layout — no phantom or missing lines.
        nlines = _measure_lines([(lead, True), (rest, False)], size, w - 0.22)
        cy += line_h * nlines + gap
    return cy


def callout(slide, x, y, w, h, label, body, label_c=MAGENTA, fill=TINT, body_c=DEEP):
    # auto-grow height so the body never spills outside the box. Height comes from
    # measure_callout (MEASURED glyph metrics, label measured bold) — the ONE place the
    # formula lives, so a build can call measure_callout()/bottom_callout() and get the
    # identical height it will render.
    h = max(h, measure_callout(label, body, w))
    box(slide, x, y, w, h, fill=fill, round=True)
    rad = 0.08 * min(w, h)                                   # inset the accent bar so its square
    box(slide, x, y + rad, 0.07, h - 2 * rad, fill=label_c)  # ends fall on the card's straight edge
    # text box spans the card's full height so MSO_ANCHOR.MIDDLE centres on the card's true
    # centre (y + h/2). A y-offset here with the same height would push the text below centre.
    text(slide, x + 0.24, y, w - 0.44, h,
         [[(label + "  ", 11, label_c, True, False), (body, 12.5, body_c, False, False)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.08)
    return y + h   # bottom edge, so callers can keep a margin below


def bottom_callout(slide, x, w, label, body, *, footer_gap=0.15, **kw):
    """A footer-SAFE bottom callout — **never collides with the footer**, whatever the body
    length. It MEASURES its own height (:func:`measure_callout`), anchors its BOTTOM just above
    the footer band, and grows UPWARD. Returns its TOP y, so the caller keeps content above it.

    This replaces the failure-prone pattern of hand-picking a low ``y`` and passing it to
    :func:`callout` (which grows DOWN into the footer when the text wraps). Use this for every
    bottom takeaway / WHY / NOTE bar::

        top = dk.bottom_callout(s, 0.6, W-1.2, "TAKEAWAY", "...")
        # ...place the slide's content within [title, top].
    """
    _, h_in = _slide_size(slide)
    ch = measure_callout(label, body, w)
    y = h_in - FOOTER_BAND - footer_gap - ch
    callout(slide, x, y, w, ch, label, body, **kw)
    return y


def chip(slide, x, y, w, h, title, sub, fill, tcolor=None):
    """A labelled box for pipeline stages. `sub` may contain '\\n' for line breaks.
    If `tcolor` is None the text colour is auto-picked (white or dark ink) for the better
    contrast against `fill` — so a chip on a light accent (gold/teal) gets dark text
    instead of unreadable white. Pass `tcolor` explicitly to override."""
    if tcolor is None:
        tcolor = WHITE if contrast_ratio(WHITE, fill) >= contrast_ratio(DEEP, fill) else DEEP
    box(slide, x, y, w, h, fill=fill, round=True)
    runs = [[(title, 14, tcolor, True, False)]]
    if sub:
        for line in sub.split("\n"):
            runs.append([(line, 10.5, tcolor, False, False)])
    # generous side padding so text never crowds the rounded edge
    text(slide, x + 0.13, y, w - 0.26, h, runs, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, space_after=1, line_spacing=0.98)


def repeat_row(slide, x, y, w, h, n, label_fmt="{i}", *, sub="", show=2, fill=BLUE,
               gap=0.3, ellipsis="…", badge=True, caption=None, caption_c=None,
               tcolor=None):
    """Show N identical-except-index parallel units as a **pattern**, NOT N duplicate cards.

    The anti-pattern this prevents: rendering many units that are identical except for an index as
    N full blocks (e.g. 8 "unit k / <same caption>" cards) — repeating the same content N× adds zero
    information, eats the whole canvas, and buries the actual message. Instead this draws ``show``
    representative chips (``label_fmt.format(i=1..show)``), an **ellipsis** cell, the **Nth** chip,
    and a ``× N`` badge — and states the shared detail common to every unit **once** (``caption``,
    defaulting to ``sub``) centered under the row.

    Domain-agnostic — use whenever units differ only by an index and N is large: parallel
    model/compute units (attention heads, stacked layers), service replicas / nodes / microservices,
    an M-model ensemble, N regional teams running one playbook, repeated pipeline stages, any long
    set of same-shaped items. ``label_fmt`` uses ``{i}`` for the index, e.g. ``"head {i}"`` →
    "head 1", "head 2", or ``"replica {i}"`` / ``"L{i}"``. When N is small (``n <= show + 2``) it
    just draws all N chips (no ellipsis/badge) — showing every one is fine there. Build the *flow*
    the units feed into (how they combine/aggregate) below the returned y; that structure, not the
    enumeration of clones, is the slide's real content.

    Returns the **bottom y** of the group (row + shared caption) in inches — anchor the next
    element (a down-arrow, the combine/aggregate block) there.
    """
    badge_w = 0.66 if (badge and n > show + 2) else 0.0
    if n <= show + 2:                                   # small N → just show them all
        labels = [label_fmt.format(i=i) for i in range(1, n + 1)]
        cells = [("chip", lbl) for lbl in labels]
    else:
        cells = [("chip", label_fmt.format(i=i)) for i in range(1, show + 1)]
        cells.append(("ellipsis", ellipsis))
        cells.append(("chip", label_fmt.format(i=n)))
    avail = w - (badge_w + gap if badge_w else 0.0)
    ncell = len(cells)
    cw = (avail - gap * (ncell - 1)) / ncell
    for k, (kind, val) in enumerate(cells):
        cx = x + k * (cw + gap)
        if kind == "chip":
            chip(slide, cx, y, cw, h, val, "", fill, tcolor=tcolor)
        else:                                            # ellipsis — a glyph, no box
            text(slide, cx, y, cw, h, [[(val, 26, MUTE, True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    if badge_w:                                          # the "× N" count badge at the right
        bh = min(h, 0.5)
        by = y + (h - bh) / 2
        bx = x + avail + gap
        box(slide, bx, by, badge_w, bh, fill=fill, round=True)
        bc = WHITE if contrast_ratio(WHITE, fill) >= contrast_ratio(DEEP, fill) else DEEP
        text(slide, bx, by, badge_w, bh, [[("× %d" % n, 16, bc, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    cap = caption if caption is not None else sub
    bottom = y + h
    if cap:                                              # shared detail — stated ONCE
        text(slide, x, bottom + 0.1, w, 0.34, [[(cap, 12.5, caption_c or MUTE, False, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, space_after=0)
        bottom += 0.44
    return bottom


def modbox(slide, x, y, w, h, role, fname, fill, tcolor=None, check=False):
    """Module box for a code/architecture diagram: big role word + mono filename.
    `tcolor=None` auto-picks white/dark ink for contrast against `fill` (see chip)."""
    if tcolor is None:
        tcolor = WHITE if contrast_ratio(WHITE, fill) >= contrast_ratio(DEEP, fill) else DEEP
    box(slide, x, y, w, h, fill=fill, round=True)
    role_runs = [[(line, 16, tcolor, True, False)] for line in role.split("\n")]
    text(slide, x + 0.05, y + 0.12, w - 0.1, 0.55, role_runs, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.92)
    text(slide, x + 0.05, y + h - 0.32, w - 0.1, 0.3, [[(fname, 9.5, tcolor, False, False, MONO)]],
         align=PP_ALIGN.CENTER, space_after=0)
    if check:
        box(slide, x + w - 0.34, y - 0.12, 0.24, 0.24, fill=TEAL)
        text(slide, x + w - 0.34, y - 0.16, 0.24, 0.28, [[("✓", 12, WHITE, True, False)]],
             align=PP_ALIGN.CENTER, space_after=0)


# ============================================================= equations
def equation_png(latex_lines, out_path, color="FFFFFF", fontsize=28, dpi=300, mathfont="cm"):
    """Render LaTeX-style math lines to a transparent PNG (proper italics, \\odot,
    real subscripts/superscripts, fractions, Greek...) for a genuinely FORMAL look —
    this is the PREFERRED way to put equations on a slide. The ASCII eq_par() below is
    only a quick fallback; baseline-shifted ASCII never looks as good as real typeset
    math, so reach for equation_png whenever the audience will read the formula.

    Pass mathtext strings, e.g. r"\\hat{x} = \\mathrm{arg\\,min}_{x}\\,\\|Ax-y\\|_2^2 +
    \\lambda R(x)". Returns (w_px, h_px); place with add_picture, then scale to a target
    HEIGHT in inches (height = target_h; width = target_h * w_px/h_px) so glyph size is
    consistent across slides. Lazy-imports matplotlib.

    `mathfont` picks the math typeface (matplotlib mathtext.fontset):
      'cm'        — Computer Modern, the classic LaTeX look (default; formal & elegant)
      'stixsans' / 'dejavusans' — upright SANS math, to sit better next to a sans deck
      'stix'      — Times-like serif math
    Pick 'cm' for a formal/classical feel (academic, defense, any serif deck); a sans set to match a crisp corporate deck.
    `color` is an RRGGBB hex string (a leading '#' is tolerated; e.g. '202A37' for dark
    text on a light deck, 'FFFFFF' for light text on a dark panel).

    mathtext quirks (it is NOT full LaTeX — some control words differ, and older matplotlib
    is stricter, so prefer the safe forms below and you won't hit version differences):
      • use \\leq \\geq \\neq \\times (NOT \\le \\ge \\ne — these reliably raise ParseException);
      • write upright text as \\mathrm{...} (don't rely on \\text{}); use \\, \\; \\ for
        spacing inside it (e.g. \\mathrm{arg\\,min});
      • \\| gives the norm bars; \\hat \\mathcal \\Psi \\lambda \\epsilon \\Delta all work;
      • keep prose annotations OUT of the math (render them as a separate deck-font label),
        so the PNG stays pure math and font-independent."""
    import os
    import tempfile
    if not os.environ.get("MPLCONFIGDIR"):
        default = os.path.join(os.path.expanduser("~"), ".matplotlib")
        if not (os.path.isdir(default) and os.access(default, os.W_OK)):
            path = os.path.join(tempfile.gettempdir(), "slide-maker-matplotlib")
            os.makedirs(path, exist_ok=True)
            os.environ["MPLCONFIGDIR"] = path
    import matplotlib; matplotlib.use("Agg")
    matplotlib.rcParams["mathtext.fontset"] = mathfont   # math typeface (see docstring)
    import matplotlib.pyplot as plt
    from PIL import Image
    n = len(latex_lines)
    fig = plt.figure(figsize=(8, 0.66 * n + 0.15)); fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
    col = "#" + color.lstrip("#")   # tolerate a leading '#' (a natural mistake)
    for i, ln in enumerate(latex_lines):
        m = ln if ln.strip().startswith("$") else f"${ln}$"   # math mode
        ax.text(0.01, 1 - (i + 0.5) / n, m, color=col, fontsize=fontsize, va="center")
    fig.savefig(out_path, dpi=dpi, transparent=True, bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return Image.open(out_path).size


def hrule(slide, x, y, w, color=MUTE, weight=0.012):
    """A thin horizontal rule — for real table lines / separators."""
    return box(slide, x, y, w, weight, fill=color)


# ================================================================= native (editable) charts
def native_chart(slide, x, y, w, h, categories, series, *, kind="line_markers",
                 palette=None, dark=False, font=None, highlight=None, legend=True,
                 value_fmt=None, smooth=True):
    """An **EDITABLE native PowerPoint chart** (a real chart object: click to edit data/labels in
    PowerPoint, and **any non-Latin labels — CJK · Cyrillic · Greek · …** — render via PowerPoint's own
    fonts, **no tofu**, unlike the rasterised designed_charts recipes). Prefer this whenever editability
    or non-Latin labels matter. Pass ``font=`` your deck's text font for the script (e.g. your EAFONT
    for CJK; a Cyrillic/Greek deck's FONT already covers those).

    `series` = [(name, [v, v, ...]), ...]; `categories` = the x labels. Themed to the deck (palette,
    dark). `kind`: 'line' | 'line_markers' | 'column' | 'bar'. `highlight` = index of the one series
    to keep in the accent (others dropped to grey). For a two-scale 'A↑ vs B↓' chart use
    `native_dual_axis` instead."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    KIND = {"line": XL_CHART_TYPE.LINE, "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED}
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    for name, vals in series:
        cd.add_series(str(name), tuple(vals))
    gf = slide.shapes.add_chart(KIND.get(kind, XL_CHART_TYPE.LINE_MARKERS),
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    _theme_chart(ch, series, palette=palette, dark=dark, font=font, highlight=highlight,
                 legend=legend, value_fmt=value_fmt, smooth=smooth, kind=kind)
    return ch


def _theme_chart(ch, series, *, palette, dark, font, highlight, legend, value_fmt, smooth, kind):
    from pptx.enum.chart import XL_LEGEND_POSITION
    ink = RGBColor(0xEA, 0xF2, 0xFF) if dark else RGBColor(0x22, 0x2A, 0x37)
    grid = RGBColor(0x2A, 0x35, 0x55) if dark else RGBColor(0xE7, 0xE9, 0xF0)
    muted = RGBColor(0x8A, 0x93, 0xA6) if dark else RGBColor(0xB8, 0xBE, 0xCC)
    pal = [_as_rgb(c) for c in (palette or ACCENTS)]
    fname = font or EAFONT or FONT          # the deck's script font → non-Latin labels render (no tofu)
    try:
        ch.font.name = fname; ch.font.size = Pt(11); ch.font.color.rgb = ink
    except Exception:
        pass
    ch.has_title = False
    multi = len(series) > 1
    ch.has_legend = bool(legend and multi)
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
        ch.legend.font.color.rgb = ink; ch.legend.font.name = fname
    for ax in (ch.category_axis, ch.value_axis):
        try:
            ax.tick_labels.font.color.rgb = ink; ax.tick_labels.font.name = fname; ax.tick_labels.font.size = Pt(10)
            ax.format.line.color.rgb = grid
        except Exception:
            pass
    try:
        ch.value_axis.major_gridlines.format.line.color.rgb = grid
        ch.value_axis.major_gridlines.format.line.width = Pt(0.5)
        ch.category_axis.has_major_gridlines = False
    except Exception:
        pass
    if value_fmt:
        try:
            ch.value_axis.tick_labels.number_format = value_fmt
            ch.value_axis.tick_labels.number_format_is_linked = False
        except Exception:
            pass
    for i, ser in enumerate(ch.series):
        col = pal[i % len(pal)]
        if highlight is not None:
            col = pal[0] if i == highlight else muted
        try:
            ser.format.line.color.rgb = col; ser.format.line.width = Pt(2.5); ser.smooth = smooth
        except Exception:
            pass
        try:
            ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = col
            ser.marker.format.line.color.rgb = col
        except Exception:
            pass
        if kind in ("column", "bar"):
            try:
                ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
            except Exception:
                pass
    return ch


def native_dual_axis(slide, x, y, w, h, categories, left, right, *, left_name="A", right_name="B",
                     palette=None, dark=False, font=None):
    """An **editable** two-scale line chart (real PowerPoint combo chart): `left` values on the left
    axis, `right` values on a SECONDARY right axis — the 'A↑ vs B↓' story (e.g. 占比% vs 成本指数).
    Click-to-edit and **any-language-safe**: non-Latin labels (CJK · Cyrillic · Greek · …) render via
    PowerPoint's fonts, no tofu (pass ``font=`` your deck's text font). The editable, non-rasterised
    replacement for ``designed_charts.dual_axis`` — especially for non-Latin labels."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    cd = CategoryChartData(); cd.categories = [str(c) for c in categories]
    cd.add_series(str(left_name), tuple(left)); cd.add_series(str(right_name), tuple(right))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    _theme_chart(ch, [(left_name, 1), (right_name, 1)], palette=palette, dark=dark, font=font,
                 highlight=None, legend=True, value_fmt=None, smooth=True, kind="line_markers")
    _chart_to_secondary(ch, dark=dark, font=font)
    return ch


def _chart_to_secondary(ch, *, dark, font):
    """Move the LAST series of a 2-series chart onto a SECONDARY right-hand value axis, drawn as a
    line (builds the combo-chart OOXML python-pptx has no public API for). Works whether the primary
    plot is a line or a bar chart — so it powers both dual-axis and Pareto (bars + cumulative line)."""
    ink = "EAF2FF" if dark else "222A37"
    fname = font or EAFONT or FONT
    pa = ch._chartSpace.xpath('.//c:plotArea')[0]
    prim = (pa.xpath('./c:lineChart') or pa.xpath('./c:barChart'))[0]
    axids = [int(e.get('val')) for e in prim.xpath('./c:axId')]
    cat2, val2 = max(axids) + 111, max(axids) + 222
    ser1 = prim.xpath('./c:ser')[-1]; prim.remove(ser1)
    txpr = (f'<c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="1000">'
            f'<a:solidFill><a:srgbClr val="{ink}"/></a:solidFill><a:latin typeface="{fname}"/>'
            f'</a:defRPr></a:pPr><a:endParaRPr lang="en-US"/></a:p></c:txPr>')
    lc2 = parse_xml(f'<c:lineChart {nsdecls("c")}><c:grouping val="standard"/><c:varyColors val="0"/>'
                    f'<c:marker val="1"/><c:axId val="{cat2}"/><c:axId val="{val2}"/></c:lineChart>')
    lc2.insert(2, ser1); prim.addnext(lc2)
    pa.append(parse_xml(f'<c:valAx {nsdecls("c", "a")}><c:axId val="{val2}"/><c:scaling>'
                        f'<c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="r"/>'
                        f'{txpr}<c:crossAx val="{cat2}"/><c:crosses val="max"/></c:valAx>'))
    pa.append(parse_xml(f'<c:catAx {nsdecls("c")}><c:axId val="{cat2}"/><c:scaling>'
                        f'<c:orientation val="minMax"/></c:scaling><c:delete val="1"/><c:axPos val="b"/>'
                        f'<c:crossAx val="{val2}"/></c:catAx>'))


def native_donut(slide, x, y, w, h, segments, center_value="", center_label="", *,
                 palette=None, dark=False, font=None):
    """Editable part-to-whole **DOUGHNUT** + an optional headline KPI in the hole (native chart;
    click-to-edit; any-language-safe). segments = [(label, value), ...]. Editable replacement for
    designed_charts.donut_kpi. (The KPI is a separate native textbox — nudge it onto the hole if needed.)"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    if not segments:
        raise ValueError("native_donut needs at least one segment")
    cd = CategoryChartData(); cd.categories = [str(s[0]) for s in segments]
    cd.add_series("", tuple(s[1] for s in segments))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ink = RGBColor(0xEA, 0xF2, 0xFF) if dark else RGBColor(0x22, 0x2A, 0x37)
    mute = RGBColor(0x8A, 0x93, 0xA6) if dark else RGBColor(0x9A, 0xA0, 0xAE)
    fname = font or EAFONT or FONT
    pal = [_as_rgb(c) for c in (palette or ACCENTS)]
    ch.has_title = False
    try:
        ch.font.name = fname; ch.font.color.rgb = ink; ch.font.size = Pt(11)
        ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False; ch.legend.font.color.rgb = ink; ch.legend.font.name = fname
    except Exception:
        pass
    for i, pt in enumerate(ch.series[0].points):
        try:
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = pal[i % len(pal)]
        except Exception:
            pass
    if center_value or center_label:
        runs = []
        if center_value:
            runs.append([(str(center_value), 28, ink, True, False, fname)])
        if center_label:
            runs.append([(str(center_label), 12, mute, False, False, fname)])
        text(slide, x + w * 0.18, y + h * 0.32, w * 0.64, h * 0.30, runs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return ch


def native_pareto(slide, x, y, w, h, items, *, palette=None, dark=False, font=None):
    """Editable **Pareto**: ranked columns + a cumulative-% line on a secondary axis (native combo
    chart; click-to-edit; any-language-safe). items = [(label, value), ...] (sorted desc by you).
    Editable replacement for designed_charts.pareto."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    if not items:
        raise ValueError("native_pareto needs at least one item")
    vals = [float(v) for _, v in items]
    tot = sum(vals) or 1.0
    cum, run = [], 0.0
    for v in vals:
        run += v; cum.append(round(100.0 * run / tot, 1))
    cd = CategoryChartData(); cd.categories = [str(k) for k, _ in items]
    cd.add_series("数量", tuple(vals)); cd.add_series("累计 %", tuple(cum))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    pal = [_as_rgb(c) for c in (palette or ACCENTS)]
    _theme_chart(ch, [("数量", 1), ("累计 %", 1)], palette=pal, dark=dark, font=font,
                 highlight=0, legend=True, value_fmt=None, smooth=False, kind="column")
    _chart_to_secondary(ch, dark=dark, font=font)
    return ch


def native_bubble(slide, x, y, w, h, points, *, palette=None, dark=False, font=None,
                  xlabel="", ylabel=""):
    """Editable **bubble** chart — x vs y with a third (size) dimension (native; click-to-edit;
    any-language-safe). points = [(x, y, size[, label]), ...]. Editable cousin of designed_charts.bubble_trend."""
    from pptx.chart.data import BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE
    if not points:
        raise ValueError("native_bubble needs at least one point")
    bcd = BubbleChartData()
    ser = bcd.add_series("")
    for p in points:
        ser.add_data_point(float(p[0]), float(p[1]), float(p[2]))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, Inches(x), Inches(y), Inches(w), Inches(h), bcd)
    ch = gf.chart
    ink = RGBColor(0xEA, 0xF2, 0xFF) if dark else RGBColor(0x22, 0x2A, 0x37)
    grid = RGBColor(0x2A, 0x35, 0x55) if dark else RGBColor(0xE7, 0xE9, 0xF0)
    fname = font or EAFONT or FONT
    pal = [_as_rgb(c) for c in (palette or ACCENTS)]
    ch.has_title = False; ch.has_legend = False
    try:
        ch.font.name = fname; ch.font.color.rgb = ink; ch.font.size = Pt(11)
        ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = pal[0]
    except Exception:
        pass
    for ax in (ch.category_axis, ch.value_axis):
        try:
            ax.tick_labels.font.color.rgb = ink; ax.tick_labels.font.name = fname
            ax.format.line.color.rgb = grid
            ax.major_gridlines.format.line.color.rgb = grid; ax.major_gridlines.format.line.width = Pt(0.5)
        except Exception:
            pass
    return ch


# ================================================================= tables & code
def _hex(c):
    return c if isinstance(c, str) else str(c)   # RGBColor.__str__ -> 'RRGGBB'

def _as_rgb(c):
    """Accept a colour as an RGBColor OR a hex string ('RRGGBB' or '#RRGGBB') — one convention
    everywhere, tolerant of a leading '#' so callers don't have to remember to strip it."""
    return RGBColor.from_string(c.lstrip("#")) if isinstance(c, str) else c

def _clear_table_style(tbl):
    """Strip PowerPoint's default banded-blue table theme so WE control every fill and
    rule (the default theme looks nothing like the deck). Sets the built-in 'No Style,
    No Grid' style and turns off first-row/banding emphasis."""
    NO_STYLE = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    el = tbl._tbl
    tblPr = el.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = el.makeelement(qn('a:tblPr'), {}); el.insert(0, tblPr)
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = tblPr.makeelement(qn('a:tableStyleId'), {}); tblPr.append(sid)
    sid.text = NO_STYLE

_FILLISH = {qn('a:noFill'), qn('a:solidFill'), qn('a:gradFill'), qn('a:blipFill'),
            qn('a:pattFill'), qn('a:grpFill'), qn('a:cell3D'), qn('a:headers')}

def _cell_rules(cell, top=None, bottom=None):
    """Add booktabs horizontal rules to a table cell. top/bottom = (color, weight_pt) or
    None. Inserts <a:lnT>/<a:lnB> in the correct schema position (before the fill)."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    anchor = next((ch for ch in tcPr if ch.tag in _FILLISH), None)
    def make(tag, spec):
        color, w = spec
        ln = tcPr.makeelement(qn(tag), {'w': str(int(w * 12700)), 'cap': 'flat',
                                        'cmpd': 'sng', 'algn': 'ctr'})
        sf = ln.makeelement(qn('a:solidFill'), {})
        sf.append(ln.makeelement(qn('a:srgbClr'), {'val': _hex(color)}))
        ln.append(sf)
        return ln
    for tag, spec in (('a:lnT', top), ('a:lnB', bottom)):   # lnT before lnB (schema order)
        if spec is None:
            continue
        ln = make(tag, spec)
        (anchor.addprevious(ln) if anchor is not None else tcPr.append(ln))


def table(slide, x, y, w, rows, col_w=None, header=True, highlight=None,
          numeric_cols=None, size=13, row_h=0.34, head_c=DEEP, body_c=SLATE,
          rule_c=MUTE, hi_fill=TINT, hi_c=MAGENTA, font=None):
    """A clean booktabs-style data table — for a dense comparison a chart can't carry
    (many methods × many metrics). `rows` = list of rows of cell strings; row 0 is the
    header when header=True.

    Foreground the comparison the AUTHORS make (see step 1): pass `highlight=<i>` (0-based
    over the BODY rows) to bold + tint the proposed method's row, so the eye lands on the
    one row that matters — a results table exists to make ONE comparison obvious, not to be
    read cell by cell. `numeric_cols` = column indices to right-align (numbers read better
    right-aligned). `col_w` = column widths in inches (defaults to an equal split of `w`).

    Styled with NO PowerPoint theme — no banded blue, no gridlines: just a top rule, a rule
    under the header, and a bottom rule, so it reads like a paper table. Returns the bottom
    y so the caller keeps a GUTTER below it.

    NOTE: on a *presented* slide keep cells terse and highlight one row — it's a slide, not a
    spreadsheet. A **read-alone reference / appendix** table can legitimately be denser (more rows,
    smaller `row_h`, per-column rules) since the reader studies it without a narrator. A cell long
    enough to wrap past `row_h` will grow the row, so verify the render. For a *trend*, prefer a chart
    (`equation_png`/matplotlib or a dedicated figure-making workflow, if available);
    a table is for exact values."""
    ncol = max(len(r) for r in rows)
    nrow = len(rows)
    col_w = col_w or [w / ncol] * ncol
    h = row_h * nrow
    tbl = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(y), Inches(w), Inches(h)).table
    _clear_table_style(tbl)
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    numeric = set(numeric_cols or [])
    hi_row = (highlight + (1 if header else 0)) if highlight is not None else None
    for i in range(nrow):
        tbl.rows[i].height = Inches(row_h)
        is_head = header and i == 0
        is_hi = (i == hi_row)
        for j in range(ncol):
            cell = tbl.cell(i, j)
            cell.text = rows[i][j] if j < len(rows[i]) else ""
            cell.margin_left = cell.margin_right = Pt(7)
            cell.margin_top = cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if is_hi:
                cell.fill.solid(); cell.fill.fore_color.rgb = hi_fill
            else:
                cell.fill.background()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j in numeric else PP_ALIGN.LEFT
            col = head_c if is_head else (hi_c if is_hi else body_c)
            for r in p.runs:
                set_font(r, size, col, bold=is_head or is_hi, font=font)
    # booktabs rules: \toprule, \midrule (under header), \bottomrule — nothing else
    spec = {}
    for j in range(ncol):
        spec.setdefault((0, j), {})['top'] = (head_c, 1.4)
        if header and nrow > 1:
            spec.setdefault((0, j), {})['bottom'] = (rule_c, 0.8)
        spec.setdefault((nrow - 1, j), {})['bottom'] = (head_c, 1.4)
    for (i, j), s in spec.items():
        _cell_rules(tbl.cell(i, j), top=s.get('top'), bottom=s.get('bottom'))
    return y + h


def code_block(slide, x, y, w, code, size=12, lang=None, highlight_lines=None,
               panel=DEEP, text_c=PALE, hi_c=WHITE, hi_fill=None, line_numbers=False,
               pad=0.16, line_h=None):
    """A monospace code panel — preserved indentation, optional per-line highlighting.
    `code` = a string with real newlines (or a list of lines); leading/trailing blank
    lines are trimmed. Renders on a dark rounded panel by default (for a LIGHT deck pass
    panel=LIGHT or a tint + text_c=DEEP). `highlight_lines` = 1-based line numbers to
    emphasise (brighter, bold; pass hi_fill=<color> for a highlight band too).

    Keep snippets SHORT — a slide shows the 5 lines that carry the idea, not a file; elide
    the rest with '# ...'. Uses MONO so indentation and columns line up. Returns bottom y."""
    lines = code.split("\n") if isinstance(code, str) else list(code)
    lines = [ln.expandtabs(4) for ln in lines]   # tabs -> spaces so indentation is real
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    if not lines:
        lines = [""]
    hl = set(highlight_lines or [])
    line_h = line_h or size / 72.0 * 1.42
    gutter_w = (size / 72.0 * 0.62 * (len(str(len(lines))) + 1)) if line_numbers else 0.0
    h = pad * 2 + line_h * len(lines)
    box(slide, x, y, w, h, fill=panel, round=True)
    if hi_fill is not None:
        for k in hl:
            if 1 <= k <= len(lines):
                box(slide, x + 0.05, y + pad + (k - 1) * line_h - line_h * 0.06,
                    w - 0.10, line_h, fill=hi_fill)
    runs = [[(ln if ln.strip() else " ", size, (hi_c if k in hl else text_c),
              k in hl, False, MONO)]
            for k, ln in enumerate(lines, start=1)]
    tb = text(slide, x + pad + gutter_w, y + pad, w - 2 * pad - gutter_w, h, runs,
              space_after=0, line_spacing=1.0)
    tb.text_frame.word_wrap = False   # a long line clips instead of wrapping (which would
    #                                   break indentation and the height estimate) — the
    #                                   docstring's "keep snippets short" is the real fix.
    if line_numbers:
        nruns = [[(str(k), size, text_c, False, False, MONO)] for k in range(1, len(lines) + 1)]
        text(slide, x + pad - 0.02, y + pad, gutter_w, h, nruns,
             align=PP_ALIGN.RIGHT, space_after=0, line_spacing=1.0)
    if lang:
        text(slide, x + w - 1.3, y + 0.06, 1.2, 0.22,
             [[(lang, 9, text_c, False, False, MONO)]], align=PP_ALIGN.RIGHT, space_after=0)
    return y + h


# ----- lightweight ASCII equation fallback (editable text; use when matplotlib
#       isn't available or the equation is trivial) -----
N   = lambda t: (t, 'n')      # normal run
SUP = lambda t: (t, 'sup')    # superscript run
SUB = lambda t: (t, 'sub')    # subscript run

def eq_par(tf, tokens, base, color, first=False, italic=False, font=None):
    """One equation line inside a text frame. tokens = list of N()/SUP()/SUB().
    Uses real baseline shifts + ASCII letters, so it renders crisply in any font.
    `font` resolves to EQFONT at call time (so reassigning deckkit.EQFONT re-themes math).
    Example:  eq_par(tf, [N('A'),SUP('H'),N(' y = '),N('D'),SUB('r'),SUP('T')], 14, WHITE)"""
    font = font or EQFONT          # resolve at call time so re-theming EQFONT takes effect
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(2); p.space_before = Pt(0); p.line_spacing = 1.12
    p.alignment = PP_ALIGN.LEFT
    for txt, kind in tokens:
        r = p.add_run(); r.text = txt
        size = base if kind == 'n' else base * 0.62
        set_font(r, size, color, italic=italic, font=font)
        if kind in ('sup', 'sub'):
            r._r.get_or_add_rPr().set('baseline', '30000' if kind == 'sup' else '-22000')
    return p


# ============================================================ editable native math
# A LaTeX-subset → real TEXT RUNS renderer: italic variables, upright operators, true
# sub/superscripts, math symbols, in a math font. The result is CLICK-EDITABLE native
# text that renders identically in PowerPoint / Keynote / LibreOffice / PDF — unlike
# equation_png (a flat raster) and unlike an OMML equation object (invisible in the
# LibreOffice render/PDF). Use this for LINEAR formulas (sums, norms, sub/superscripts,
# Greek, operators); 2-D math (fractions, matrices, stacked limits) raises — use
# equation_png there. See SKILL §4 "Equations".
_EQ_GREEK = {'alpha':'α','beta':'β','gamma':'γ','delta':'δ','epsilon':'ε','varepsilon':'ε',
'zeta':'ζ','eta':'η','theta':'θ','vartheta':'ϑ','iota':'ι','kappa':'κ','lambda':'λ','mu':'μ',
'nu':'ν','xi':'ξ','pi':'π','rho':'ρ','sigma':'σ','tau':'τ','upsilon':'υ','phi':'φ','varphi':'φ',
'chi':'χ','psi':'ψ','omega':'ω','Gamma':'Γ','Delta':'Δ','Theta':'Θ','Lambda':'Λ','Xi':'Ξ',
'Pi':'Π','Sigma':'Σ','Upsilon':'Υ','Phi':'Φ','Psi':'Ψ','Omega':'Ω'}
_EQ_SYM = {'sum':'Σ','prod':'Π','int':'∫','oint':'∮','partial':'∂','nabla':'∇','infty':'∞',
'cdot':'·','times':'×','div':'÷','pm':'±','mp':'∓','ast':'∗','star':'⋆','circ':'∘','bullet':'∙',
'leq':'≤','le':'≤','geq':'≥','ge':'≥','neq':'≠','ne':'≠','approx':'≈','sim':'∼','simeq':'≃',
'equiv':'≡','cong':'≅','propto':'∝','ll':'≪','gg':'≫','in':'∈','notin':'∉','subset':'⊂',
'subseteq':'⊆','supset':'⊃','cup':'∪','cap':'∩','forall':'∀','exists':'∃','rightarrow':'→',
'to':'→','Rightarrow':'⇒','leftarrow':'←','Leftarrow':'⇐','leftrightarrow':'↔','mapsto':'↦',
'odot':'⊙','oplus':'⊕','otimes':'⊗','langle':'⟨','rangle':'⟩','ldots':'…','cdots':'⋯','dots':'…',
'top':'⊤','perp':'⊥','angle':'∠','prime':'′','hbar':'ℏ','ell':'ℓ','Re':'ℜ','Im':'ℑ'}
_EQ_MCAL = {'A':'𝒜','B':'ℬ','C':'𝒞','D':'𝒟','E':'ℰ','F':'ℱ','G':'𝒢','H':'ℋ','I':'ℐ','J':'𝒥',
'K':'𝒦','L':'ℒ','M':'ℳ','N':'𝒩','O':'𝒪','P':'𝒫','Q':'𝒬','R':'ℛ','S':'𝒮','T':'𝒯','U':'𝒰',
'V':'𝒱','W':'𝒲','X':'𝒳','Y':'𝒴','Z':'𝒵'}
_EQ_BB = {'R':'ℝ','N':'ℕ','Z':'ℤ','Q':'ℚ','C':'ℂ','E':'𝔼','P':'ℙ'}
_EQ_ACC = {'hat':'̂','widehat':'̂','tilde':'̃','widetilde':'̃','bar':'̄',
'vec':'⃗','dot':'̇','ddot':'̈','check':'̌','breve':'̆','acute':'́',
'grave':'̀'}
_EQ_2D = {'frac','dfrac','tfrac','sqrt','begin','overline','underline','binom','matrix','pmatrix',
'bmatrix','vmatrix','overbrace','underbrace','substack'}

def _eq_read_group(s, i):
    depth = 0; j = i
    while j < len(s):
        if s[j] == '{': depth += 1
        elif s[j] == '}':
            depth -= 1
            if depth == 0: return s[i+1:j], j+1
        j += 1
    return s[i+1:], len(s)

def _eq_resolve(s):
    """Resolve a LaTeX chunk (no top-level _/^) → list of (display_char, italic_bool)."""
    out = []; i = 0
    while i < len(s):
        c = s[i]
        if c == '\\':
            j = i+1; name = ''
            while j < len(s) and s[j].isalpha(): name += s[j]; j += 1
            if name == '':                                   # escaped symbol: \|  \{  \}  \%
                sym = s[j] if j < len(s) else ''
                out.append(('‖' if sym == '|' else sym, False)); i = j+1; continue
            if name in _EQ_ACC:                              # accents: \hat{x} → x̂
                if j < len(s) and s[j] == '{': inner, j = _eq_read_group(s, j)
                else: inner = s[j] if j < len(s) else ''; j += 1
                u = _eq_resolve(inner)
                if u: u[0] = (u[0][0] + _EQ_ACC[name], u[0][1])
                out.extend(u); i = j; continue
            if name in ('mathcal','mathbf','mathrm','mathbb','mathit','text','operatorname','boldsymbol'):
                if j < len(s) and s[j] == '{': inner, j = _eq_read_group(s, j)
                else: inner = s[j] if j < len(s) else ''; j += 1
                if name == 'mathcal':
                    for ch in inner: out.append((_EQ_MCAL.get(ch, ch), ch not in _EQ_MCAL))
                elif name == 'mathbb':
                    for ch in inner: out.append((_EQ_BB.get(ch, ch), ch not in _EQ_BB))
                elif name in ('mathrm','text','operatorname'):
                    for ch in inner: out.append((ch, False))
                else:
                    for ch in inner: out.append((ch, ch.isalpha()))
                i = j; continue
            if name in _EQ_2D:
                raise NotImplementedError(f"\\{name} needs 2-D layout — use equation_png for this formula")
            if name in _EQ_GREEK: out.append((_EQ_GREEK[name], False)); i = j; continue
            if name in _EQ_SYM:   out.append((_EQ_SYM[name], False)); i = j; continue
            if name in ('left','right','displaystyle','textstyle','limits','nolimits','bigl','bigr','Bigl','Bigr'):
                i = j; continue
            if name in ('quad','qquad'): out.append(('  ', False)); i = j; continue
            for ch in name: out.append((ch, False))          # unknown command (\min,\arg,\log…) → upright
            i = j; continue
        if c in '{}': i += 1; continue
        if c == '|': out.append(('‖', False)); i += 1; continue
        if c == ' ': out.append((' ', False)); i += 1; continue
        out.append((c, c.isalpha())); i += 1
    return out

def latex_to_runs(latex):
    """LaTeX-subset string → list of (text, kind) tokens; kind ∈ n/i/sub/sup/isub/isup.
    Raises NotImplementedError on 2-D constructs (\\frac, matrices, …)."""
    s = latex.strip()
    for a, b in (('\\,',' '),('\\;',' '),('\\:',' '),('\\!',''),('~',' '),('\\ ',' ')):
        s = s.replace(a, b)
    out = []; i = 0
    def push(units, lvl):
        for ch, ital in units:
            kind = ('isup' if ital else 'sup') if lvl > 0 else \
                   ('isub' if ital else 'sub') if lvl < 0 else ('i' if ital else 'n')
            out.append((ch, kind))
    while i < len(s):
        c = s[i]
        if c in '_^':
            lvl = 1 if c == '^' else -1; j = i+1
            if j < len(s) and s[j] == '{':
                inner, j = _eq_read_group(s, j); push(_eq_resolve(inner), lvl)
            elif j < len(s) and s[j] == '\\':
                k = j+1
                while k < len(s) and s[k].isalpha(): k += 1
                push(_eq_resolve(s[j:k]), lvl); j = k
            else:
                push(_eq_resolve(s[j:j+1]) if j < len(s) else [], lvl); j = j+1
            i = j; continue
        if c == '\\':
            j = i+1
            while j < len(s) and s[j].isalpha(): j += 1
            if j < len(s) and s[j] == '{': _, j = _eq_read_group(s, j)
            push(_eq_resolve(s[i:j]), 0); i = j; continue
        push(_eq_resolve(c), 0); i += 1
    return out

EQ_MATHFONT = "STIX Two Math"   # math font with ℒ Σ ‖ … ; 'Cambria Math' is the Office-portable alt

def equation_native(slide, x, y, w, h, latex, *, size=20, color=DEEP, font=None,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    """EDITABLE native math — the PREFERRED way to put a formula the user may edit on a slide.
    Renders a LaTeX-subset (or a pre-tokenised list) as real, click-editable PowerPoint TEXT
    RUNS (italic variables · upright operators · true sub/superscripts · math glyphs) in a math
    `font`, so it renders identically in PowerPoint / Keynote / LibreOffice / PDF AND stays
    editable — unlike `equation_png` (a flat raster) and unlike an OMML equation object (which is
    invisible in the LibreOffice render & PDF export). For LINEAR formulas; 2-D math (fractions,
    matrices, stacked limits) raises NotImplementedError → use `equation_png` for those.

    `latex` e.g. r"\\mathcal{L} = \\sum_i \\|A x_i - y_i\\|_2^2 + \\lambda R(x_i)" (or a list of
    (text, kind) tokens). `font` defaults to a math font (`EQ_MATHFONT` = 'STIX Two Math'; set it
    to 'Cambria Math' for Office portability — flag the dependency at hand-off). `size` is the
    base point size; keep it ≈ the deck's body size, consistent across slides. Returns the textbox."""
    toks = latex if isinstance(latex, list) else latex_to_runs(latex)
    fnt = font or EQ_MATHFONT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    try: tf.vertical_anchor = anchor
    except Exception: pass
    for m in ('left','right','top','bottom'):
        setattr(tf, 'margin_' + m, Inches(0.02))
    p = tf.paragraphs[0]; p.alignment = align
    for txt, k in toks:
        r = p.add_run(); r.text = txt
        sz = size * (0.62 if k in ('sub','sup','isub','isup') else 1.0)
        set_font(r, sz, color, italic=(k in ('i','isub','isup')), font=fnt)
        if k in ('sup','isup'): r._r.get_or_add_rPr().set('baseline', '30000')
        if k in ('sub','isub'): r._r.get_or_add_rPr().set('baseline', '-22000')
    return tb


# ================================================================ template reuse
def open_template(path):
    """Open the user's deck and delete its slides while KEEPING masters/layouts.
    A template's branding (header band, logos, footer) lives on the layouts, so new
    slides added afterwards inherit all of it automatically. Dropping the slide
    relationships also prunes the old slides' heavy media (e.g. GIFs) on save."""
    prs = Presentation(path)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.get(qn('r:id')))
        sldIdLst.remove(sldId)
    return prs


def drop_placeholders(slide, keep_idx):
    """Remove inherited placeholders whose idx is not in keep_idx (set)."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx not in keep_idx:
            ph._element.getparent().remove(ph._element)


def content_slide(prs, layout_idx, title_txt, size=23, footer="", date="",
                  title_color=WHITE, body_idx=1):
    """Add a content slide on the template's 'Title and Content with Logo' layout.
    Sets a WHITE title on the band, removes the empty body placeholder, fills the
    footer/date placeholders. Returns the slide for you to draw on.

    NOTE: layout_idx and placeholder types are template-specific — confirm them by
    inspecting the template once (inspect_template.py) and save them to the template's
    profile.md in the active template registry, e.g. ~/.codex/slide-templates/<name>/
    for Codex or ~/.claude/slide-templates/<name>/ for Claude Code."""
    s = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    tf = s.shapes.title.text_frame; tf.text = title_txt
    for r in tf.paragraphs[0].runs:
        set_font(r, size, title_color, bold=True)
    # remove the empty body content placeholder so it doesn't show a prompt
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx == body_idx:
            ph._element.getparent().remove(ph._element)
    # fill footer (type 15) and date (type 16) placeholders if present
    for ph in list(s.placeholders):
        t = ph.placeholder_format.type
        if t == 15 and footer:
            ph.text = footer
            for r in ph.text_frame.paragraphs[0].runs: set_font(r, 9, MUTE)
        elif t == 16 and date:
            ph.text = date
            for r in ph.text_frame.paragraphs[0].runs: set_font(r, 9, MUTE)
    return s


# ===================================== no-template branch (build chrome yourself)
def blank_deck(w_in=10.0, h_in=5.625):
    """A fresh 16:9 deck when the user has NO template to match. Use add_slide() +
    title_bar()/footer() to give it simple, consistent branding you define from
    their brand colors (set the palette constants above) or a clean default."""
    prs = Presentation()
    prs.slide_width = Inches(w_in); prs.slide_height = Inches(h_in)
    return prs


def add_slide(prs):
    """Add a truly blank slide (layout 6) to draw on from scratch."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_size(slide):
    """Return the slide's real width/height in inches."""
    prs = slide.part.package.presentation_part.presentation
    return prs.slide_width / 914400, prs.slide_height / 914400


def title_bar(slide, title, kicker="", accent=MAGENTA, title_c=DEEP, w_in=None):
    """Lightweight slide chrome for the no-template branch: optional kicker, title,
    and a short accent rule. Pair with footer().

    By default this reads the actual deck width from ``slide`` so it works on
    widescreen templates, custom ``blank_deck(w_in, h_in)`` sizes, and posters.
    Pass ``w_in`` only when deliberately overriding that geometry."""
    if w_in is None:
        w_in, _ = _slide_size(slide)
    if kicker:
        text(slide, 0.55, 0.30, w_in - 1.1, 0.3, [[(kicker.upper(), 11, BLUE, True, False)]], space_after=0)
        ty = 0.54
    else:
        ty = 0.40
    tb = text(slide, 0.55, ty, w_in - 1.1, 0.7,
              [[(title, 26, title_c, True, False, DISPLAY or FONT)]], space_after=0)  # title gets the DISPLAY face
    if EADISPLAY:                                    # ...and a distinct CJK display face if set
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                _apply_ea(r, EADISPLAY)
    box(slide, 0.57, ty + 0.62, 1.1, 0.045, fill=accent)


def footer(slide, tag="", page=None, w_in=None, h_in=None):
    """Footer tag (left) + optional page number (right) for the no-template branch.

    Defaults to the actual deck size, so custom-sized decks don't end up with a
    page number stranded at the old 10x5.625 coordinate."""
    if w_in is None or h_in is None:
        sw, sh = _slide_size(slide)
        w_in = sw if w_in is None else w_in
        h_in = sh if h_in is None else h_in
    if tag:
        text(slide, 0.55, h_in - 0.35, 6.0, 0.3, [[(tag, 8, MUTE, False, False)]], space_after=0)
    if page is not None:
        text(slide, w_in - 1.0, h_in - 0.35, 0.6, 0.3,
             [[(str(page), 9, MUTE, True, False)]], align=PP_ALIGN.RIGHT, space_after=0)


def logo(slide, path, *, corner="tr", h=0.42, margin=0.3, w_in=None, h_in=None, alt=None):
    """Place a brand / institution / product logo as PERSISTENT chrome — the SAME mark in the
    SAME spot on every slide. For a deck that is *about* a company, institution, or product,
    the entity's real logo belongs on every content slide (top-right by convention) so the
    audience always knows whose deck this is and the brand stays present; in the no-template /
    generated branch there are no layouts to carry it, so call this once per slide (after the
    background, before content) — and on the hero/cover too. It's chrome, not content: keep `h`
    small (~0.35-0.5 in) so it never competes with the title, and keep `corner`/`margin`
    identical across slides so it doesn't jump. The logo holds its aspect ratio (a wide wordmark
    and a square mark both sit right). `corner` is "tr" (default), "tl", "br", or "bl".

    Use the REAL logo (see references/image-generation.md's real-asset hierarchy) — if it's
    missing, ask the user for it or drop in an honest "logo here" placeholder; NEVER an
    AI-imagined or recolored look-alike. On a busy/dark background, give the logo a small scrim
    or light plate behind it so it stays legible. Returns the picture shape."""
    from PIL import Image
    sw, sh = _slide_size(slide)
    w_in = sw if w_in is None else w_in
    h_in = sh if h_in is None else h_in
    with Image.open(path) as im:
        iw, ih = im.size
    w = h * (iw / ih) if ih else h
    top = "t" in corner.lower()
    left = "l" in corner.lower()
    x = margin if left else (w_in - margin - w)
    y = margin if top else (h_in - margin - h)
    return picture(slide, path, x, y, w, h, fit="contain", alt=("" if alt is None else alt))


# ===================================================================== notes
def speaker_notes(slide, notes):
    """Attach speaker notes — the SPOKEN script — to a slide, off the visible canvas.
    The slide should show the phrase; the full sentences the presenter says live here.
    Notes do NOT render to the slide or the PNG (the critic won't see them); they appear
    in PowerPoint/Keynote Presenter View and print on Notes Pages. This is the right
    place to move full sentences OFF a slide for a talk/defense/lecture the user will
    rehearse — keeping the slide a clean visual aid while preserving the narration."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes
    return slide.notes_slide


def alt_text(shape, description):
    """Set a shape's ACCESSIBILITY alt-text (the screen-reader description). Call it after
    add_picture() — and on any informative figure/diagram — so assistive tech can describe
    the visual. Alt-text does NOT render (it won't appear in the PNG or to the pixel
    critic); it lives in PowerPoint's accessibility metadata (the cNvPr 'descr' attribute)
    and shows in the Selection/Accessibility pane. Give a one-line factual description,
    e.g. alt_text(pic, "ROC curve: proposed method (pink) above baseline (grey)")."""
    cNvPr = shape._element.find(".//" + qn("p:cNvPr"))
    if cNvPr is not None:
        cNvPr.set("descr", description)
        if not cNvPr.get("title"):
            cNvPr.set("title", description[:120])
    return shape


# ════════════════════════════════════════════════════════════════════════════════════════
# Design-pattern components mined from professional sample decks (ppt-master gallery).
# See references/design-gallery.md + semantic-color-contract.md for when/why.
# ════════════════════════════════════════════════════════════════════════════════════════

def _blend(c, other, t):
    """Blend RGBColor c toward `other` by fraction t (0..1). For tints / faint watermarks."""
    if isinstance(c, str):
        c = RGBColor.from_string(c.lstrip("#"))
    if isinstance(other, str):
        other = RGBColor.from_string(other.lstrip("#"))
    return RGBColor(*(int(round(a + (b - a) * t)) for a, b in zip(c, other)))


def highlight(s, size, base_c, accent_c, *, key_bold=True, font=None):
    """Split a string with <k>…</k> tags into a text() RUN PARAGRAPH where tagged spans are
    recolored to `accent_c` (and bolded). Gives a sentence/headline a scannable second layer —
    recolor exactly ONE phrase per headline, a few per body line.
        text(s, x,y,w,h, [highlight("the <k>one phrase</k> that matters", 16, INK, ACCENT)])
    Returns a list of run tuples (a single paragraph)."""
    import re
    runs = []
    for i, seg in enumerate(re.split(r"<k>(.*?)</k>", s)):
        if not seg:
            continue
        key = (i % 2 == 1)
        runs.append((seg, size, accent_c if key else base_c, key and key_bold, False, font) if font
                    else (seg, size, accent_c if key else base_c, key and key_bold, False))
    return runs


def node(slide, x, y, w, h, label, *, shape="roundrect", fill=None, line=None, line_w=1.4,
         tcolor=None, sub="", dashed=False, hub=False, accent=None):
    """One diagram NODE (box/connector kit — the general architecture/flowchart builder).
    `shape`='roundrect'|'rect'|'pill'|'circle'. By default a thin-outline pale node; `hub=True`
    promotes it to a SOLID accent fill (the ONE focal node — keep every other node thin-outline,
    exactly one hub). `dashed=True` marks an optional/inferred node. Returns the node's CENTER
    (cx, cy) in inches so connector() can join it. Pair with connector() + the stroke-semantics
    convention (solid=required · dashed=optional · dotted=feedback)."""
    acc = accent if accent is not None else BLUE
    ln = line if line is not None else acc
    sh = {"pill": MSO_SHAPE.ROUNDED_RECTANGLE, "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
          "rect": MSO_SHAPE.RECTANGLE, "circle": MSO_SHAPE.OVAL}.get(shape, MSO_SHAPE.ROUNDED_RECTANGLE)
    o = slide.shapes.add_shape(sh, Inches(x), Inches(y), Inches(w), Inches(h))
    o.shadow.inherit = False
    if hub:
        o.fill.solid(); o.fill.fore_color.rgb = acc; o.line.fill.background()
        tc = tcolor if tcolor is not None else (WHITE if contrast_ratio(WHITE, acc) >= contrast_ratio(DEEP, acc) else DEEP)
    else:
        o.fill.solid(); o.fill.fore_color.rgb = WHITE
        o.line.color.rgb = ln; o.line.width = Pt(line_w)
        tc = tcolor if tcolor is not None else DEEP
    if dashed and not hub:
        el = o.line._get_or_add_ln(); el.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="dash"/>'))
    if shape == "pill":
        try: o.adjustments[0] = 0.5
        except Exception: pass
    runs = [[(label, 13, tc, True, False)]]
    if sub:
        runs.append([(sub, 9.5, _blend(tc, WHITE, 0.25) if hub else MUTE, False, False, MONO)])
    text(slide, x + 0.06, y, w - 0.12, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         space_after=1, line_spacing=0.96)
    return (x + w / 2, y + h / 2)


def connector(slide, p0, p1, *, style="solid", color=None, width=1.5, label="", arrow=True,
              label_c=None):
    """Join two node points (cx,cy tuples from node()) with a connector. `style`='solid'
    (required) | 'dashed' (optional) | 'dotted' (feedback/inferred) — the stroke SEMANTICS that
    make a technical diagram readable. `label` = an on-shaft mono edge label (centred at midpoint).
    `arrow=True` adds an arrowhead at p1."""
    col = color if color is not None else MUTE
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(p0[0]), Inches(p0[1]),
                                   Inches(p1[0]), Inches(p1[1]))
    c.line.color.rgb = col; c.line.width = Pt(width); c.shadow.inherit = False
    if style in ("dashed", "dotted"):
        ln = c.line._get_or_add_ln()
        ln.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="{"dash" if style == "dashed" else "sysDot"}"/>'))
    if arrow:
        ln = c.line._get_or_add_ln()
        ln.append(parse_xml(f'<a:tailEnd {nsdecls("a")} type="triangle" w="med" len="med"/>'))
    if label:
        mx, my = (p0[0] + p1[0]) / 2, (p0[1] + p1[1]) / 2
        text(slide, mx - 0.8, my - 0.16, 1.6, 0.3, [[(label, 9, label_c or col, False, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return c


def elbow_connector(slide, pts, *, style="solid", color=None, width=1.5, arrow=True, label="", label_c=None):
    """A multi-segment ELBOW / U-shaped connector through `pts` (list of (x,y) inch tuples) — the right
    arrow when a STRAIGHT line would be wrong or cross other shapes: a **feedback / repeat loop** that
    drops below a row and returns, a **return path**, or a link between **non-adjacent** nodes. Don't
    default every arrow to straight — straight is for direct adjacent flow; an elbow reads as
    'goes back / around'. Same stroke SEMANTICS as `connector` (solid=required · dashed=optional ·
    dotted=feedback). Arrowhead on the FINAL segment only. Helper `loop_path(...)` builds a common U.
    Example (a repeat-loop under a 4-node row at y≈3.1, dropping to 3.5):
        dk.elbow_connector(s, dk.loop_path(x_last, x_first, 3.1, 3.5), style="dotted", color=CYAN)"""
    col = color if color is not None else MUTE
    segs = []
    for i in range(len(pts) - 1):
        a, b = pts[i], pts[i + 1]
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1]))
        c.line.color.rgb = col; c.line.width = Pt(width); c.shadow.inherit = False
        if style in ("dashed", "dotted"):
            ln = c.line._get_or_add_ln()
            ln.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="{"dash" if style == "dashed" else "sysDot"}"/>'))
        segs.append(c)
    if arrow and segs:
        ln = segs[-1].line._get_or_add_ln()
        ln.append(parse_xml(f'<a:tailEnd {nsdecls("a")} type="triangle" w="med" len="med"/>'))
    if label and len(pts) >= 2:
        mid = pts[len(pts) // 2]
        text(slide, mid[0] - 0.9, mid[1] - 0.16, 1.8, 0.3, [[(label, 9, label_c or col, False, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return segs


def loop_path(x_from, x_to, y_row, y_drop):
    """Waypoints for a U-shaped feedback/repeat loop: from (x_from, y_row) DOWN to y_drop, across to
    x_to, and UP to (x_to, y_row). Feed to `elbow_connector`. y_drop should clear the row's content."""
    return [(x_from, y_row), (x_from, y_drop), (x_to, y_drop), (x_to, y_row)]


_ALGO_KW = {"input", "output", "require", "ensure", "initialize", "for", "while", "repeat",
            "until", "if", "else", "elif", "endfor", "endif", "endwhile", "end", "then", "do",
            "return", "break", "continue", "function", "procedure", "foreach"}

def _algo_runs(text_s, size, ink, kwc, font):
    """Split a pseudocode line into runs, bolding + colouring control-flow keywords."""
    import re
    out = []
    for tok in re.split(r"(\s+)", text_s):
        key = tok.strip().lower().rstrip(":")
        if tok.strip() and key in _ALGO_KW:
            out.append((tok, size, kwc, True, False, font))
        else:
            out.append((tok, size, ink, False, False, font))
    return out or [(text_s, size, ink, False, False, font)]

def algorithm_block(slide, x, y, w, lines, *, title="Algorithm 1", caption=None, size=10.5,
                    ink=None, kw_color=None, rule_c=None, font=None, indent=0.20, gutter=0.30,
                    pad=0.16, number=True, boxed=False, fill=None):
    """Render a pseudocode ALGORITHM block — LaTeX `algorithm`/`algorithmic`-environment style, for
    describing a CS/AI method, training loop, or optimization procedure as exact, skimmable steps.

    `lines` = list of entries, each either:
        "text"                      → indent level 0, auto-numbered
        (text, indent_level)        → indent_level * indent of left padding
        (text, indent_level, False) → not numbered (e.g. a Input:/Output: header or a sub-note)
    Control-flow KEYWORDS (Input, Output, for, while, if, else, then, do, return, end…) are auto-bolded
    and tinted `kw_color`. Lines auto-number 1..N (skip with the per-entry flag). Default look is academic
    **booktabs rules** (a thick top rule, a hairline under the title, a thick bottom rule, no side
    borders); pass `boxed=True` (or a `fill`) for a full rounded card instead. Use a MONO `font` for the
    classic pseudocode feel. Pair the block with one prose line of intuition — the block gives the exact
    procedure, the prose gives the why. Returns the block height H (so you can place a caption below).
    See references/form-selection.md ('an algorithm / procedure') and design-gallery.md."""
    ink = DEEP if ink is None else ink
    kwc = ink if kw_color is None else kw_color
    rc = ink if rule_c is None else rule_c
    f = font or MONO or FONT
    lh = size * 1.66 / 72.0
    th = size * 2.0 / 72.0
    norm = []
    for e in lines:
        if isinstance(e, (tuple, list)):
            t = e[0]; ind = e[1] if len(e) > 1 else 0; nm = e[2] if len(e) > 2 else True
        else:
            t = e; ind = 0; nm = True
        norm.append((t, ind, nm))
    # estimate wrapped visual-line count per entry (mono char-width ≈ 0.62·size) so long lines
    # don't collide with the next step; keep lines short (one display line) for the cleanest look.
    cw = 0.66 * size / 72.0
    counts = []
    for (t, ind, nm) in norm:
        lx0 = pad + gutter + 0.09 + ind * indent
        cpl = max(6, int((w - pad - lx0) / cw))
        counts.append(max(1, -(-len(t) // cpl)))
    H = pad + th + 0.06 + sum(counts) * lh + pad * 0.7
    framed = boxed or fill is not None
    if framed:
        box(slide, x, y, w, H, fill=(fill if fill is not None else WHITE), line=rc, line_w=1.1, round=True)
    else:
        box(slide, x, y, w, 0.020, fill=rc)            # thick top rule
        box(slide, x, y + H, w, 0.020, fill=rc)        # thick bottom rule
    # title row ("Algorithm N: caption")
    cap = (":  " + caption) if caption else ""
    text(slide, x + pad, y + pad * 0.7, w - 2 * pad, th,
         [[(title, size + 0.5, ink, True, False, f), (cap, size + 0.5, ink, False, False, f)]],
         space_after=0)
    ry_mid = y + pad * 0.7 + th + 0.02
    box(slide, x, ry_mid, w, 0.010, fill=rc)           # hairline under title
    # numbered, indented lines (advance by each entry's wrapped height)
    yy = ry_mid + 0.05
    i = 0
    for (t, ind, nm), c in zip(norm, counts):
        i += 1
        if number and nm:
            text(slide, x + pad, yy, gutter, lh, [[(str(i) + ":", size, ink, False, False, f)]],
                 align=PP_ALIGN.RIGHT, space_after=0)
        lx = x + pad + gutter + 0.09 + ind * indent
        text(slide, lx, yy, x + w - pad - lx, c * lh, [_algo_runs(t, size, ink, kwc, f)],
             space_after=0, line_spacing=1.32)
        yy += c * lh
    return H


def flow_chain(slide, x, y, w, h, labels, *, accent=None, gap=None, subs=None, hub_idx=None,
               vertical=False, hub_accent=None):
    """Convenience over node()+connector(): a CHAIN of nodes joined by arrows (a pipeline).
    `labels` = list of node titles; `subs` optional same-length sub-labels; `hub_idx` promotes
    one node to a solid fill. `hub_accent` colours that hub differently from the chain `accent`
    (e.g. a coral 'Generate' hub among cyan retrieval nodes — keeps a semantic-colour contract).
    Horizontal by default; `vertical=True` stacks + down-arrows. Returns the list of node centers."""
    acc = accent if accent is not None else BLUE
    hubacc = hub_accent if hub_accent is not None else acc
    n = len(labels); g = gap if gap is not None else 0.34
    centers = []
    if vertical:
        nh = (h - g * (n - 1)) / n
        for i, lab in enumerate(labels):
            ny = y + i * (nh + g)
            c = node(slide, x, ny, w, nh, lab, sub=(subs[i] if subs else ""),
                     hub=(i == hub_idx), accent=(hubacc if i == hub_idx else acc))
            if i: connector(slide, (x + w / 2, ny - g + 0.02), (x + w / 2, ny - 0.02), color=_blend(acc, WHITE, 0.3))
            centers.append(c)
    else:
        nw = (w - g * (n - 1)) / n
        for i, lab in enumerate(labels):
            nx = x + i * (nw + g)
            c = node(slide, nx, y, nw, h, lab, sub=(subs[i] if subs else ""),
                     hub=(i == hub_idx), accent=(hubacc if i == hub_idx else acc))
            if i: connector(slide, (nx - g + 0.02, y + h / 2), (nx - 0.02, y + h / 2), color=_blend(acc, WHITE, 0.3))
            centers.append(c)
    return centers


def step_list(slide, x, y, w, items, *, orientation="vertical", accent=None, ink=None,
              body_c=None, numeral_style="arabic", active_idx=None, gap=None):
    """A NUMBERED process / step list. items = [(title, body), …]. orientation='vertical'
    (numbered spine, title+body rows) or 'horizontal' (connected pill steps with arrows).
    `numeral_style`='arabic'|'pad2' (01) |'cjk'. `active_idx` accents one step (terminal/current).
    Returns the bottom y (vertical)."""
    acc = accent if accent is not None else BLUE
    ic = ink if ink is not None else DEEP
    bc = body_c if body_c is not None else SLATE
    def numr(i):
        if numeral_style == "cjk": return cjk_numeral(i + 1)
        if numeral_style == "pad2": return f"{i + 1:02d}"
        return str(i + 1)
    if orientation == "horizontal":
        n = len(items); g = gap if gap is not None else 0.3
        cw = (w - g * (n - 1)) / n
        for i, (title, body) in enumerate(items):
            cx = x + i * (cw + g); on = (active_idx == i)
            d = 0.5
            box(slide, cx + cw / 2 - d / 2, y, d, d, fill=acc if on else WHITE,
                line=None if on else acc, line_w=1.4, round=True, r=d / 2)
            tc = (WHITE if contrast_ratio(WHITE, acc) >= contrast_ratio(DEEP, acc) else DEEP) if on else acc
            text(slide, cx + cw / 2 - d / 2, y, d, d, [[(numr(i), 15, tc, True, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
            if i: connector(slide, (cx - g + 0.02, y + d / 2), (cx - 0.02, y + d / 2), color=_blend(acc, WHITE, 0.4))
            text(slide, cx, y + d + 0.12, cw, 0.34, [[(title, 13, ic, True, False)]], align=PP_ALIGN.CENTER, space_after=0)
            if body:
                text(slide, cx, y + d + 0.46, cw, 0.5, [[(body, 11, bc, False, False)]], align=PP_ALIGN.CENTER, space_after=0)
        return y + d + 1.0
    # vertical
    g = gap if gap is not None else 0.2
    cy = y
    for i, (title, body) in enumerate(items):
        d = 0.42; on = (active_idx == i)
        box(slide, x, cy, d, d, fill=acc if (on or active_idx is None) else _blend(acc, WHITE, 0.0),
            round=True, r=d / 2)
        tc = WHITE if contrast_ratio(WHITE, acc) >= contrast_ratio(DEEP, acc) else DEEP
        text(slide, x, cy, d, d, [[(numr(i), 14, tc, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        text(slide, x + d + 0.18, cy - 0.02, w - d - 0.18, 0.32, [[(title, 14.5, ic, True, False)]], space_after=0)
        bh = 0.3
        if body:
            text(slide, x + d + 0.18, cy + 0.3, w - d - 0.18, 0.6, [[(body, 11.5, bc, False, False)]], space_after=0, line_spacing=1.02)
            bh = 0.62
        cy += max(d, bh) + g
    return cy


def ghost_numeral(slide, x, y, w, h, text_str, *, color=None, bg=None, opacity=0.12, font=None,
                  align=PP_ALIGN.LEFT):
    """A giant FAINT index/ordinal/year numeral sitting BEHIND content as silent wayfinding +
    texture (8–18% strength). Draw it FIRST (behind the card/title). The **bg-aware successor to
    `big_numeral(mode='ghost')`** — it blends `color` toward `bg` by (1-opacity), so unlike that
    light-only watermark it works on a DARK deck too (pass the deck's `bg`). For a *foreground* hero
    figure use `big_numeral` / `stat_row` instead. No alpha needed."""
    c = _blend(color if color is not None else MUTE, bg if bg is not None else WHITE, 1 - opacity)
    sz = int(min(h * 72 * 1.05, 220))
    text(slide, x, y, w, h, [[(str(text_str), sz, c, True, False, font or DISPLAY or FONT)]],
         align=align, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.9)


def insight_banner(slide, x, y, w, body, *, label="INSIGHT", fill=None, accent=None,
                   tcolor=None, h=0.62):
    """The consulting 'so-what' BANNER — a full-width dark rounded bar under a slide's action
    title carrying the one-sentence implication (label caps in accent + the sentence). Returns
    bottom y."""
    f = fill if fill is not None else DEEP
    acc = accent if accent is not None else GOLD
    tc = tcolor if tcolor is not None else WHITE
    box(slide, x, y, w, h, fill=f, round=True)
    box(slide, x, y + 0.1, 0.06, h - 0.2, fill=acc)
    runs = [(label + "   ", 11, acc, True, False), (body, 13.5, tc, False, False)]
    text(slide, x + 0.28, y, w - 0.5, h, [runs], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return y + h


def bilingual_lockup(slide, x, y, w, zh, en, *, zh_size=30, en_size=11, ink=None, accent=None,
                     rule=True, zh_font=None, en_font=None, anchor_top=True):
    """A CJK (or any) heavy display headline auto-paired with a wide-tracked ALL-CAPS Latin/pinyin
    strap line beneath — the most universal 'instantly professional' lockup. Optional short accent
    rule between. Returns bottom y."""
    ic = ink if ink is not None else DEEP
    acc = accent if accent is not None else MAGENTA
    tb = text(slide, x, y, w, 0.7, [[(zh, zh_size, ic, True, False, zh_font or EADISPLAY or DISPLAY or FONT)]], space_after=0)
    eaface = zh_font or EADISPLAY      # CJK glyphs render from <a:ea> — set the DISPLAY face there too
    if eaface:
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                _apply_ea(r, eaface)
    yy = y + zh_size / 72.0 + 0.12
    if rule:
        box(slide, x + 0.02, yy, 0.9, 0.035, fill=acc); yy += 0.14
    text(slide, x + 0.01, yy, w, 0.3,
         [[(" ".join(en.upper()) if len(en) < 26 else en.upper(), en_size, MUTE, True, False, en_font or FONT)]],
         space_after=0)
    return yy + 0.3


def concept_equation(slide, x, y, w, h, terms, *, op="=", accent=None, ink=None, highlight_idx=None,
                     term_size=30, op_size=34, font=None, term_colors=None):
    """A word-EQUATION headline device (not LaTeX math): big display terms joined by oversized
    accent operators — 'ZINE = MAGAZINE', 'A ≠ B ≠ C', 'Answer = Retrieve + Generate'. `op` is the
    joiner: a single string ('='|'≠'|'×'|'+'|'→') for all gaps, OR a **list** of operators (one per
    gap, len == len(terms)-1) for a MIXED equation like Answer = Retrieve + Generate (op=['=','+']).
    `highlight_idx` recolors one term to the accent; `term_colors` (a list aligned to `terms`,
    `None` per slot = default) overrides per-term colour for a semantic palette. Centred row."""
    acc = accent if accent is not None else MAGENTA
    ic = ink if ink is not None else DEEP
    ops = list(op) if isinstance(op, (list, tuple)) else [op] * (len(terms) - 1)
    runs = []
    for i, t in enumerate(terms):
        col = (term_colors[i] if term_colors and i < len(term_colors) and term_colors[i] is not None
               else (acc if i == highlight_idx else ic))
        runs.append((t, term_size, col, True, False, font or DISPLAY or FONT))
        if i < len(terms) - 1:
            runs.append(("  " + ops[i] + "  ", op_size, acc, True, False, font or DISPLAY or FONT))
    tb = text(slide, x, y, w, h, [runs], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    eaface = font or EADISPLAY        # honor the CJK display face for any CJK terms
    if eaface:
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                _apply_ea(r, eaface)
    return y + h


def cta_button(slide, x, y, w, h, label, *, variant="primary", accent=None, tcolor=None,
               arrow=True, mono=False):
    """A call-to-action button: filled `variant='primary'` (accent fill) or `variant='secondary'`
    (outline). Optional trailing → arrow / mono command label. Use on closing/cover masters."""
    acc = accent if accent is not None else BLUE
    if variant == "primary":
        box(slide, x, y, w, h, fill=acc, round=True, r=min(h / 2, 0.16))
        tc = tcolor if tcolor is not None else (WHITE if contrast_ratio(WHITE, acc) >= contrast_ratio(DEEP, acc) else DEEP)
    else:
        box(slide, x, y, w, h, fill=None, line=acc, line_w=1.4, round=True, r=min(h / 2, 0.16))
        tc = tcolor if tcolor is not None else acc
    lab = label + ("  →" if arrow else "")
    text(slide, x, y, w, h, [[(lab, 13, tc, True, False, MONO if mono else None)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def cta_pair(slide, x, y, w, h, primary, secondary, *, accent=None, gap=0.2):
    """Primary (filled) + secondary (outline) CTA buttons side by side, equal height."""
    bw = (w - gap) / 2
    cta_button(slide, x, y, bw, h, primary, variant="primary", accent=accent)
    cta_button(slide, x + bw + gap, y, bw, h, secondary, variant="secondary", accent=accent)


def status_stamp(slide, x, y, text_str, *, color=None, size=0.95, rotation=-12):
    """A rotated state STAMP ('SOLD OUT', 'CONFIDENTIAL') — a bordered caps mark attached to a
    card/footer. Independent of the CJK `seal`."""
    c = color if color is not None else MAGENTA
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(size * 1.7), Inches(size * 0.5))
    sh.fill.background(); sh.line.color.rgb = c; sh.line.width = Pt(1.6); sh.rotation = rotation
    sh.shadow.inherit = False
    text(slide, x, y, size * 1.7, size * 0.5, [[(text_str.upper(), 13, c, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return sh


def corner_tab(slide, card_x, card_y, card_w, label, *, fill=None, tcolor=None, w=1.6, h=0.32):
    """A 'RECOMMENDED' / 'MOST POPULAR' tab sitting ON the top edge of a card (centred on it). Its
    bottom meets the card's top edge so it reads as attached WITHOUT overlapping into the card (no
    false overlap lint). Build the card first, then call this with the card's x/y/w."""
    f = fill if fill is not None else MAGENTA
    tc = tcolor if tcolor is not None else WHITE
    bx = card_x + card_w / 2 - w / 2
    ty = card_y - h            # bottom edge meets the card top — attached, not overlapping
    box(slide, bx, ty, w, h, fill=f, round=True, r=h / 2)
    text(slide, bx, ty, w, h, [[(label.upper(), 9.5, tc, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def concentric_rings(slide, cx, cy, layers, *, accent=None, ink=None, r0=1.6, ring=0.62):
    """A nested-containment / synthesis diagram (core → ring → ring) with leader labels — for a
    qualitative framework (e.g. CMT 色彩·材质·纹理). layers = [outer…inner] labels."""
    acc = accent if accent is not None else GOLD
    ic = ink if ink is not None else DEEP
    n = len(layers)
    for i in range(n):
        r = r0 - i * ring
        col = _blend(acc, WHITE, 0.75 - 0.22 * i)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
        o.fill.solid(); o.fill.fore_color.rgb = col; o.line.color.rgb = acc; o.line.width = Pt(1.2); o.shadow.inherit = False
    for i, lab in enumerate(layers):
        r = r0 - i * ring
        ly = cy - r + ring / 2 if i < n - 1 else cy
        text(slide, cx - 1.3, ly - 0.16, 2.6, 0.32, [[(lab, 13 if i < n - 1 else 14, ic, True, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def pull_quote(slide, x, y, w, quote, *, attribution="", accent=None, ink=None, serif=None, size=20):
    """An italic-serif PULL QUOTE with an oversized accent quote-mark and attribution — turns a
    statement into an argument. Returns bottom y."""
    acc = accent if accent is not None else MAGENTA
    ic = ink if ink is not None else DEEP
    f = serif or DISPLAY or "Georgia"
    text(slide, x, y - 0.1, 0.7, 0.7, [[("“", 54, acc, True, False, f)]], space_after=0)
    text(slide, x + 0.05, y + 0.5, w, 1.4, [[(quote, size, ic, False, True, f)]], space_after=2, line_spacing=1.1)
    yy = y + 0.5 + (size / 72.0) * 1.1 * (1 + len(quote) // max(1, int(w * 9))) + 0.2
    if attribution:
        text(slide, x + 0.05, yy, w, 0.3, [[("— " + attribution, 12, MUTE, False, False)]], space_after=0)
        yy += 0.3
    return yy


def standfirst(slide, x, y, w, text_str, *, ink=None, serif=None, size=15):
    """An italic-serif STANDFIRST / dekker — the one-line editorial gloss under a headline.
    Returns bottom y."""
    text(slide, x, y, w, 0.5, [[(text_str, size, ink if ink is not None else SLATE, False, True, serif or DISPLAY or "Georgia")]],
         space_after=0, line_spacing=1.05)
    return y + size / 72.0 * 1.05 + 0.1


def dot_meter(slide, x, y, n, total, *, accent=None, off=None, d=0.13, gap=0.07):
    """A ●●○ complexity/level meter — n filled of `total` dots."""
    acc = accent if accent is not None else BLUE
    of = off if off is not None else _blend(acc, WHITE, 0.72)
    for i in range(total):
        box(slide, x + i * (d + gap), y, d, d, fill=acc if i < n else of, round=True, r=d / 2)


def tradeoff_list(slide, x, y, w, plus, minus, *, pos=None, neg=None, recommended=False):
    """A +/− trade-off list: green '+' pros and red '−' cons. plus/minus = lists of strings."""
    pc = pos if pos is not None else GREEN
    nc = neg if neg is not None else RGBColor(0xD0, 0x3A, 0x2E)
    cy = y
    for sign, col, items in (("+", pc, plus), ("−", nc, minus)):
        for it in items:
            text(slide, x, cy, 0.3, 0.28, [[(sign, 14, col, True, False)]], space_after=0)
            text(slide, x + 0.32, cy, w - 0.32, 0.3, [[(it, 12, SLATE, False, False)]], space_after=0)
            cy += 0.3
    return cy


def segmented_bar(slide, x, y, w, h, parts, *, labels=None, accents=None, show_pct=True):
    """A cumulative 100% SEGMENTED bar. parts = list of values (auto-normalised). Distinct hues.
    Returns bottom y."""
    tot = sum(parts) or 1
    cols = accents or palette(len(parts), ACCENTS)
    cx = x
    for i, v in enumerate(parts):
        seg = w * v / tot
        box(slide, cx, y, seg, h, fill=cols[i], round=False)
        if show_pct and seg > 0.5:
            tc = WHITE if contrast_ratio(WHITE, cols[i]) >= contrast_ratio(DEEP, cols[i]) else DEEP
            lab = (labels[i] + " " if labels else "") + f"{round(100 * v / tot)}%"
            text(slide, cx, y, seg, h, [[(lab, 10.5, tc, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        cx += seg
    return y + h


def meter_bar(slide, x, y, w, frac, *, label=None, value=None, value_unit=None,
              accent=None, track=None, h=0.28, gap_above=0.34, value_pos="right",
              value_pad=0.2, value_w=2.6, label_size=12.5, value_size=18, unit_size=11,
              ink=None, label_c=None, font=None, value_font=None):
    """A horizontal METER / progress bar with a value label that is **always vertically
    centered on the bar** (never floating above or below it) — the correct, reusable form of
    the hand-built "track + fill + number" row. Use for percentile / share / progress / "want
    vs have" rows (e.g. ``第10百分位``, ``95%``).

    Draws: an optional caption ``label`` ABOVE the bar, a rounded ``track``, an ``accent`` fill
    to ``frac`` (0..1), and ``value`` (+ optional ``value_unit``) set on the bar's centerline.

    ``value_pos``:
      - ``"right"`` (default) — value sits just past the END OF THE TRACK, so a column of
        meter_bars shares one value column and reads aligned. Never overlaps the fill.
      - ``"end"`` — value sits just past the FILL (tied to the bar length); good for a single
        bar, but for a short fill it floats over the empty track, so prefer ``"right"`` in a
        stack.

    Colours default to deckkit's light palette; on a DARK deck pass ``track=`` your panel
    colour, ``ink=`` your body colour, ``label_c=`` your muted colour, and ``accent=``.
    Returns the bar's BOTTOM y (so a caller can stack rows)."""
    acc = accent if accent is not None else BLUE
    tr = track if track is not None else RGBColor(0xE6, 0xE9, 0xEE)
    vink = ink if ink is not None else DEEP
    lc = label_c if label_c is not None else MUTE
    f = max(0.0, min(1.0, frac))
    yt = y + (gap_above if label else 0.0)
    if label:
        text(slide, x, y, w, gap_above, [[(str(label), label_size, lc, False, False, font)]],
             anchor=MSO_ANCHOR.BOTTOM, space_after=0)
    box(slide, x, yt, w, h, fill=tr, round=True, r=h / 2)
    if f > 0:
        box(slide, x, yt, max(w * f, h), h, fill=acc, round=True, r=h / 2)
    if value is not None:
        vx = (x + w + value_pad) if value_pos == "right" else (x + w * f + value_pad)
        runs = [(str(value), value_size, vink, True, False, value_font or font)]
        if value_unit:
            runs.append((" " + str(value_unit), unit_size, lc, True, False, value_font or font))
        # value box shares the bar's y/h with MIDDLE anchor → text on the bar's centerline
        text(slide, vx, yt, value_w, h, [runs], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return yt + h


def year_badge(slide, x, y, text_str, *, fill=None, tcolor=None, w=0.95, h=0.4):
    """A small year/date PILL badge (anchors chronology on timelines/cards)."""
    f = fill if fill is not None else GOLD
    tc = tcolor if tcolor is not None else (WHITE if contrast_ratio(WHITE, f) >= contrast_ratio(DEEP, f) else DEEP)
    box(slide, x, y, w, h, fill=f, round=True, r=h / 2)
    text(slide, x, y, w, h, [[(str(text_str), 12, tc, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def spec_card(slide, x, y, w, h, rows, *, ink=None, accent=None, fill=None, title=""):
    """A mono key→value PLACARD (Rendering / Palette / Layout; a CAD title-block) that documents a
    figure/layout as a recipe. rows = [(key, value), …]."""
    ic = ink if ink is not None else DEEP
    box(slide, x, y, w, h, fill=fill if fill is not None else LIGHT, line=_blend(ic, WHITE, 0.82), line_w=1.0, round=True)
    cy = y + 0.16
    if title:
        text(slide, x + 0.18, cy, w - 0.36, 0.3, [[(title.upper(), 10.5, accent if accent is not None else MAGENTA, True, False, MONO)]], space_after=0)
        cy += 0.32
    for k, v in rows:
        text(slide, x + 0.18, cy, w * 0.42, 0.26, [[(str(k).upper(), 9.5, MUTE, False, False, MONO)]], space_after=0)
        text(slide, x + w * 0.44, cy, w * 0.54 - 0.18, 0.26, [[(str(v), 10.5, ic, True, False, MONO)]], space_after=0)
        cy += 0.28
    return y + h


def diagram_island(slide, x, y, w, h, *, caption="", bezel=None, fill=None, cap_c=None, pad=0.3):
    """A bright rounded device-bezel PANEL hosting a flowchart/figure on a DARK slide (the 'white
    island' move). Draw the island, then build your diagram inside the returned inner rect. Adds an
    optional 'Figure N' caption below."""
    bz = bezel if bezel is not None else WHITE
    f = fill if fill is not None else WHITE
    box(slide, x, y, w, h, fill=bz, round=True, r=0.14)
    inner = (x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    if f != bz:
        box(slide, *inner, fill=f, round=True, r=0.08)
    if caption:
        text(slide, x, y + h + 0.08, w, 0.3, [[(caption, 11, cap_c if cap_c is not None else MUTE, False, False)]],
             align=PP_ALIGN.CENTER, space_after=0)
    return inner


def gradient_rule(slide, x, y, w, c0, c1, *, h=0.05, angle=0):
    """A thin two-stop GRADIENT rule (navy→emerald, amber→blue) — a brand signature under a title
    or along an edge."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.fill.gradient()
    try:
        stops = sh.fill.gradient_stops
        stops[0].color.rgb = c0 if not isinstance(c0, str) else RGBColor.from_string(c0.lstrip("#"))
        stops[0].position = 0.0
        stops[1].color.rgb = c1 if not isinstance(c1, str) else RGBColor.from_string(c1.lstrip("#"))
        stops[1].position = 1.0
        sh.fill.gradient_angle = angle
    except Exception:
        sh.fill.solid(); sh.fill.fore_color.rgb = c0 if not isinstance(c0, str) else RGBColor.from_string(c0.lstrip("#"))
    return sh


def catalogue_frame(slide, *, inset=0.32, gap=0.06, color=None, line_w=1.0, slide_obj=None,
                    w_in=None, h_in=None):
    """A thin DOUBLE-LINE full-bleed frame inset from the slide edges — the printed-specimen /
    exhibition-catalogue look (pair with the museum_memorial / eastern presets). Draw it as a
    background element (before content). Reads the slide's real size when possible."""
    if w_in is None or h_in is None:
        s = slide_obj or slide
        try:
            prs = s.part.package.presentation_part.presentation
            w_in = prs.slide_width / 914400 if w_in is None else w_in
            h_in = prs.slide_height / 914400 if h_in is None else h_in
        except Exception:
            w_in = 10.0 if w_in is None else w_in
            h_in = 5.625 if h_in is None else h_in
    c = color if color is not None else GOLD
    for k in (0, gap):
        box(slide, inset + k, inset + k, w_in - 2 * (inset + k), h_in - 2 * (inset + k),
            fill=None, line=c, line_w=line_w)
