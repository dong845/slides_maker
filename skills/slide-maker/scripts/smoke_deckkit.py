#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""smoke_deckkit — exercise the crash-prone CORE of deckkit (+ the chart recipes) with canonical
examples on a blank deck and assert none raises. NOT exhaustive: it covers roughly half of the
public helpers (~68 of 138 — the positional/tuple colour contracts and geometry paths that have
actually broken), so a passing run is a regression guard, not proof every helper works. Run after
editing deckkit/designed_charts:

    python scripts/smoke_deckkit.py     # exits non-zero on any failure
"""
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deckkit as dk          # noqa: E402
import designed_charts as dc  # noqa: E402
from deckkit import RGBColor  # noqa: E402

def C(h):
    return RGBColor.from_string(h)

TMP = tempfile.gettempdir()
IMG = os.path.join(TMP, "_smoke.png")
from PIL import Image  # noqa: E402
Image.new("RGB", (800, 450), (120, 150, 200)).save(IMG)

W, H = 13.333, 7.5
prs = dk.blank_deck(W, H)
def S():
    return dk.add_slide(prs)
def last():
    return prs.slides[-1]

fails = []
def ok(name, fn):
    try:
        fn()
    except Exception as e:
        fails.append((name, repr(e))); print(f"  FAIL {name} -> {e!r}")
def raises(name, fn):                       # a guard that SHOULD reject bad input
    try:
        fn(); fails.append((name, "did not raise")); print(f"  FAIL {name} (expected an error)")
    except Exception:
        pass

s = S()
ok("box (hex colours)", lambda: dk.box(s, 0, 0, 3, 1, fill="C0362C", line="333333"))     # unified colour API
ok("text", lambda: dk.text(s, 0, 2, 4, 1, [[("hi", 18, C("222222"), True, False)]]))
ok("title_bar/footer", lambda: (dk.title_bar(s, "T", kicker="k"), dk.footer(s, "tag", 1)))
ok("columns/rows", lambda: (dk.columns(3, slide=s), dk.rows(2, slide=s)))
ok("callout", lambda: dk.callout(s, 0, 3, 4, 1, "head", "body"))
ok("chip", lambda: dk.chip(s, 6, 4, 2, 0.6, "A", "sub", C("007CC2")))
ok("bullet", lambda: dk.bullet(s, 5, 1, 5, [("Lead ", "body")]))
ok("hrule", lambda: dk.hrule(s, 0, 5, 5))
ok("table", lambda: dk.table(S(), 0.6, 1.5, 6, [["A", "B"], ["1", "2"]], highlight=1))
ok("scorecard (numeric value)", lambda: dk.scorecard(S(), 0.6, 1, 2.5, 1.8, "Users", 1234, delta="3.2pp"))
ok("scorecard glass", lambda: dk.scorecard(last(), 4, 1, 2.5, 1.8, "X", "9", glass_tint=C("5B8DEF")))
ok("leaderboard (mixed + short row)", lambda: dk.leaderboard(S(), 0.6, 1, 5,
   [(C("007CC2"), "a", 42), (C("C0362C"), "b", "18", "sub"), (C("1B7A3D"), "short-row")]))
ok("takeaway_rail", lambda: dk.takeaway_rail(last(), 7, 1, 3, "lab", "+5", "body"))
ok("change_stat", lambda: dk.change_stat(S(), 4, 3, 3, 0.7, "<10%", "≈40%"))
ok("stat_row", lambda: dk.stat_row(S(), 0.7, 2, 8, [("8", "x", "label"), ("99", "%", "two")]))
ok("stat_row (2-tuple, no unit)", lambda: dk.stat_row(last(), 0.7, 3.2, 8, [("8", "no-unit"), ("9", "also")]))
ok("quadrant", lambda: dk.quadrant(S(), 1.6, 1.6, 5, 4, x_labels=("lo", "hi"), y_labels=("a", "b")))
ok("hub_spoke", lambda: dk.hub_spoke(S(), 6, 4, 2.0, "Core", ["a", "b", "c", "d"]))
    # r=2.0 (was 1.8): at 1.8 two spoke chips sat a ~0.05in sliver apart, which the
    # SLIVER_GAP lint correctly flags — the fixture is nudged, the check is not weakened.
ok("timeline h", lambda: dk.timeline(S(), 0.7, 2, 11.4, [("1979", "first"), ("2026", "now", "cap")], highlight=1))
ok("timeline v", lambda: dk.timeline(S(), 0.7, 2, 5, [("x", "a"), ("y", "b")], orientation="v"))
ok("image_tab", lambda: dk.image_tab(S(), 0.5, 0.5, "BEFORE"))
ok("before_after", lambda: dk.before_after(last(), 0.7, 2, 7, 3, IMG, IMG))
ok("photo_triptych", lambda: dk.photo_triptych(S(), [IMG, IMG, IMG]))
ok("photo_card", lambda: dk.photo_card(S(), 0.5, 0.5, 3, 1, role="primary"))
ok("corner_frame", lambda: dk.corner_frame(S()))
ok("accent_one", lambda: dk.accent_one(["a", "b", "c"], 1, C("C0362C")))
ok("cover", lambda: dk.cover(S(), "Title", issue_label="No 1", subtitle="sub"))
ok("colophon (list credits)", lambda: dk.colophon(S(), "tag", credits=["a", "b"], tooling="x"))
ok("sources_page", lambda: dk.sources_page(S(), [f"ref {i}" for i in range(6)]))
ok("part_eyebrow/page_marker", lambda: (dk.part_eyebrow(S(), 0.7, 0.5, "x"), dk.page_marker(last(), 2, 8)))
ok("specimen_card", lambda: dk.specimen_card(S(), 1, 1, 2, 2.5, "Aa", "Sans"))
ok("specimen_card (small h)", lambda: dk.specimen_card(last(), 4, 1, 1, 0.5, "Aa", "tiny"))
ok("wireframe_grid/spec_list", lambda: (dk.wireframe_grid(S(), 0.7, 1.5, 7, 4, [("H", 0, 4, 0, 1)]),
                                        dk.spec_list(last(), 8, 1.5, ["base = 8 px"])))
ok("glass/glow/scrim/offset", lambda: (lambda sl: (dk.box(sl, 0, 0, W, H, fill="0A0E27"),
   dk.glow(sl, 3, 3, 5, 4, C("5B4BE0")), dk.glass_card(sl, 1, 1, 3, 2, C("5B8DEF")),
   dk.scrim_overlay(sl, 0, 5, W, 2), dk.offset_shadow(sl, 5, 5, 2, 0.6, C("F5B301"))))(S()))
ok("editorial_header", lambda: dk.editorial_header(S(), "eyebrow", "Title", serif="Georgia"))
ok("big_numeral", lambda: dk.big_numeral(S(), 0.5, 0.5, "04", mode="ghost"))
ok("picture", lambda: dk.picture(S(), IMG, 1, 1, 4, 3, round=True))
def _wordmark():
    p = os.path.join(TMP, "_wm.png")
    dk.wordmark("Acme Labs", p, rule=True, monogram=True)
    assert os.path.exists(p) and os.path.getsize(p) > 0, "wordmark PNG missing/empty"
    dk.logo(S(), p, h=0.4)                       # placed afterward, exactly like a real logo
ok("wordmark (typographic logo stand-in)", _wordmark)
ok("backdrop_motif", lambda: dk.backdrop_motif(S(), accent_disc=C("C0362C")))
ok("native_chart (editable)", lambda: dk.native_chart(S(), 0.6, 1, 6, 3.2, ["Q1", "Q2", "Q3"],
   [("新客", [10, 18, 26]), ("老客", [30, 33, 36])], kind="column", dark=True, highlight=0))
ok("native_dual_axis (editable)", lambda: dk.native_dual_axis(S(), 0.6, 1, 7, 3.2, ["m1", "m2", "m3"],
   [5, 24, 40], [100, 80, 62], left_name="占比（%）", right_name="成本（指数）", dark=True))
ok("native_donut (editable)", lambda: dk.native_donut(S(), 0.6, 1, 5, 4, [("私域", 40), ("公域", 35), ("其他", 25)], "40%", "私域占比", dark=True))
ok("native_pareto (editable)", lambda: dk.native_pareto(S(), 0.6, 1, 8, 4, [("华东", 45), ("华北", 28), ("华南", 18)],
   dark=True, count_name="数量", cum_name="累计 %"))    # CJK deck → pass the CJK series names explicitly
ok("native_bubble (editable)", lambda: dk.native_bubble(S(), 0.6, 1, 8, 4, [(1, 2, 30), (2, 3, 55), (3, 2.4, 20)], dark=True))
def _csv_chart():
    p = os.path.join(TMP, "_s.csv"); open(p, "w").write("m,a,b\nJ,1,\"2,000\"\nF,3,4\n")
    cats, series = dk.series_from_csv(p, "m", ["a", "b"])
    assert cats == ["J", "F"] and series[0] == ("a", [1.0, 3.0]) and series[1][1] == [2000.0, 4.0], "series_from_csv parse"
    dk.native_chart(S(), 0.6, 1, 6, 3, cats, series, kind="column")
ok("series_from_csv + native_chart", _csv_chart)
raises("series_from_csv rejects a missing column", lambda: dk.series_from_csv(os.path.join(TMP, "_s.csv"), "nope", ["a"]))

ok("dc.donut_kpi", lambda: dc.donut_kpi(os.path.join(TMP, "_d.png"), [("a", 3), ("b", 2)], "5", "x"))
ok("dc.dumbbell", lambda: dc.dumbbell(os.path.join(TMP, "_db.png"), [("a", 1, 2)], highlight=0))
ok("dc.slope", lambda: dc.slope(os.path.join(TMP, "_sl.png"), [("a", 1, 2)], highlight=0))
ok("dc.dual_axis", lambda: dc.dual_axis(os.path.join(TMP, "_da.png"), [1, 2], [1, 2], [2, 1]))
ok("dc.bubble_trend", lambda: dc.bubble_trend(os.path.join(TMP, "_bt.png"), [(1, 2, 3, "a"), (2, 3, 5, "b")]))
ok("dc.pareto", lambda: dc.pareto(os.path.join(TMP, "_pa.png"), [("a", 4), ("b", 2)]))

# --- the build-time geometry gate must actually catch a fault and (strict) block the save ---
def _offcanvas_deck():
    p = dk.blank_deck(10, 5.625); s = dk.add_slide(p)
    dk.text(s, 9.2, 2.0, 3.0, 0.5, [[("this text runs off the right slide edge entirely here", 16, dk.DEEP, False, False)]])
    return p
def _lint_layout_gate():
    crit = [f for f in dk.lint_layout(_offcanvas_deck(), verbose=False) if f[1] == "CRITICAL"]
    assert crit, "lint_layout missed an off-canvas CRITICAL"
    p = dk.blank_deck(10, 5.625); s = dk.add_slide(p)
    dk.text(s, 1.0, 1.0, 4.0, 0.5, [[("a tidy line well inside the slide", 16, dk.DEEP, False, False)]])
    assert not [f for f in dk.lint_layout(p, verbose=False) if f[1] == "CRITICAL"], "lint_layout false-flagged a clean slide"
ok("lint_layout catches off-canvas + passes clean", _lint_layout_gate)
raises("lint_layout(strict=True) raises on a CRITICAL", lambda: dk.lint_layout(_offcanvas_deck(), strict=True, verbose=False))

raises("donut_kpi([]) rejects empty", lambda: dc.donut_kpi(os.path.join(TMP, "_x.png"), [], "0", "n"))
raises("dual_axis rejects empty", lambda: dc.dual_axis(os.path.join(TMP, "_x.png"), [], [], []))
raises("wireframe_grid rejects cols=0", lambda: dk.wireframe_grid(S(), 0, 0, 5, 4, [], cols=0))
raises("picture rejects missing file", lambda: dk.picture(S(), "/no/such/img.png", 0, 0, 3, 2))

ok("cycle_diagram (4 nodes + feedback + center)", lambda: dk.cycle_diagram(
    S(), 6.5, 3.6, [("获客", "3 行业"), ("见效", "7 天"), ("使用", "3 指标"), ("续费", "达标增购")],
    rx=1.5, ry=1.0, center="飞轮", feedback=(3, 0), feedback_label="回流"))
    # rx=1.5 (was 1.4): at 1.4 the hub sat a ~0.05in sliver from two nodes, which the
    # SLIVER_GAP lint correctly flags — the fixture is nudged, the check is not weakened.
ok("dumbbell_board (hero + threshold)", lambda: dk.dumbbell_board(
    S(), 0.8, 1.6, 11.0, [("ARR", "+51%", 4980, 6350, 4300, 6800, "万"),
                          ("NRR", "首次>100%", 92, 108, 85, 118, "%")],
    hero=1, threshold=(1, 100, "100%")))

# --- columns()/rows() weights= : proportional split, symmetric outer margins, guarded input ---
def _weights_grid():
    s = S()
    rail, main = dk.columns(2, slide=s, weights=(1, 2))
    assert abs(main[2] / rail[2] - 2.0) < 1e-6, "columns weights=(1,2) should give a 2x width ratio"
    assert abs(rail[0] - (W - (main[0] + main[2]))) < 1e-6, "outer margins must stay symmetric"
    r_top, r_bot = dk.rows(2, slide=s, weights=(3, 1))
    assert abs(r_top[3] / r_bot[3] - 3.0) < 1e-6, "rows weights=(3,1) should give a 3x height ratio"
    for a, b in zip(dk.columns(3, slide=s), dk.columns(3, slide=s, weights=(1, 1, 1))):
        assert all(abs(u - v) < 1e-9 for u, v in zip(a, b)), "equal weights must reproduce the default grid"
ok("columns/rows weights= (1:2 split, symmetric margins)", _weights_grid)
raises("columns rejects a length-mismatched weights tuple", lambda: dk.columns(2, slide=last(), weights=(1, 2, 3)))
raises("rows rejects a non-positive weight", lambda: dk.rows(2, slide=last(), weights=(1, 0)))

# --- dumbbell_board: leftward (lower-is-better) rows flip labels OUTWARD; optional v_mid dot.
#     Asserted on INK GEOMETRY directly (dk._ink_rect + dk._overlap_area), because lint_layout's
#     overlap_tol provably passes the old broken output — plus a byte-identity guard that a
#     rightward row's label frames still sit at the legacy coordinates. ---
def _dumbbell_directional():
    s = S()
    x, y0, w = 0.8, 1.6, 11.0
    dk.dumbbell_board(s, x, y0, w, [
        ("Cost / query", "", 0.62, 0.55, 0, 0.8, "$"),      # leftward: the proven-colliding row
        ("Latency", "p95", 1240, 620, 0, 1500, "ms"),       # leftward, wide span
        ("ARR", "+51%", 4980, 6350, 4300, 6800, "万"),      # rightward control (legacy geometry)
        ("NRR", "", 92, 108, 85, 118, "%", 100),            # rightward + v_mid (8-element row)
    ], label_w=0.42 * w)   # explicit label_w exercises the byte-identity frame guard below
    labels = {"0.62", "0.55 $", "1240", "620 ms", "4980", "6350 万", "92", "108 %", "100"}
    inks = []
    for sh in s.shapes:
        bb = dk._bbox_in(sh)
        if bb is None or not dk._is_text(sh):
            continue
        t = sh.text_frame.text.strip()
        if t in labels:
            r = dk._ink_rect(sh, bb)
            if r:
                inks.append((t, r[0], bb))
    assert len(inks) == 9, f"expected 9 value labels, found {len(inks)}"
    for i in range(len(inks)):
        for j in range(i + 1, len(inks)):
            ov = dk._overlap_area(inks[i][1], inks[j][1])
            assert ov < 1e-6, f"value-label ink overlap {ov:.3f}in^2: {inks[i][0]} x {inks[j][0]}"
    # byte-identity guard: the rightward row's label frames at the pre-change coordinates
    lw = 0.42 * w
    bx0, bx1 = x + lw + 0.15, x + w - 1.02
    span = 6800.0 - 4300.0
    x0 = bx0 + (4980 - 4300) / span * (bx1 - bx0)
    x1 = bx0 + (6350 - 4300) / span * (bx1 - bx0)
    ry = y0 + 2 * 0.52
    fb = next(b for t, _i, b in inks if t == "4980")
    fa = next(b for t, _i, b in inks if t == "6350 万")
    assert abs(fb[0] - (x0 - 0.62)) < 1e-3 and abs(fb[1] - (ry - 0.185)) < 1e-3, \
        "rightward before-label moved from its legacy frame"
    assert abs(fa[0] - (x1 + 0.10)) < 1e-3 and abs(fa[1] - (ry - 0.21)) < 1e-3, \
        "rightward after-label moved from its legacy frame"
ok("dumbbell_board leftward + v_mid rows (outward labels, zero ink overlap)", _dumbbell_directional)

# --- lint_layout SLIVER_GAP: warns on a near-touching seam, silent on a rule-compliant gap ---
def _sliver_gap():
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.box(s2, 1, 1.00, 4, 1.02, fill="E3E6EC")
    dk.box(s2, 1, 2.04, 4, 1.02, fill="E3E6EC")      # the documented 0.02in pitch-seam bug
    f = dk.lint_layout(p, verbose=False)
    assert any(c == "SLIVER_GAP" for _n, _sv, c, _m in f), "SLIVER_GAP missed a 0.02in seam"
    p2 = dk.blank_deck(10, 5.625); s3 = dk.add_slide(p2)
    dk.box(s3, 1, 1.0, 4, 1.0, fill="E3E6EC")
    dk.box(s3, 1, 2.2, 4, 1.0, fill="E3E6EC")        # a clear 0.2in gap — must stay silent
    assert not any(c == "SLIVER_GAP" for _n, _sv, c, _m in dk.lint_layout(p2, verbose=False)), \
        "SLIVER_GAP false-positive on a rule-compliant 0.2in gap"
ok("lint_layout SLIVER_GAP (0.02in seam warns, 0.2in gap silent)", _sliver_gap)

ok("tint mixes toward white", lambda: dk.tint("1B7F5C", 0.14))
ok("kpi_card (delta + strip, tall enough)", lambda: dk.kpi_card(
    S(), 0.8, 0.8, 3.4, 2.3, "净收入留存 NRR", "108", unit="%", delta="+16pt",
    delta_color=dk.RGBColor(0x1B, 0x7F, 0x5C), sub="从 92%", strip="首次超过流失"))
ok("flow_compare (old/new + highlight + note)", lambda: dk.flow_compare(
    S(), 0.8, 1.4, 11.5, ["签约", "排期", "对接", "上线"], ["达标签约", "复用模板", "首次转化"],
    old_label="旧流程", new_label="新流程",             # CJK deck → pass the CJK row labels explicitly
    old_result="27 天", new_result="7.5 天", highlight_old=2, highlight_new=2,
    note="40% 卡在此", transition_label="模板化"))

# --- axis_scale: the shared value→x mapper (linear, degenerate-safe) ---
def _axis_scale():
    X, draw = dk.axis_scale(1.0, 8.0, 0.0, 100.0)
    assert abs(X(0) - 1.0) < 1e-6 and abs(X(100) - 9.0) < 1e-6 and abs(X(50) - 5.0) < 1e-6, \
        "axis_scale linear mapping wrong"
    Xd, _ = dk.axis_scale(1.0, 8.0, 5.0, 5.0)          # hi==lo must not divide by zero
    assert Xd(5.0) == 1.0, "axis_scale degenerate span not guarded"
    draw(S(), 3.0)                                       # draw_axis must not raise
ok("axis_scale (linear map + degenerate span + draw_axis)", _axis_scale)

# --- dot_strip: value-mapped dots with anti-collision labels, lint-clean ---
def _dot_strip():
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.dot_strip(s2, 0.6, 2.0, 8.0, [("博后", 70), ("学术", 100), ("工业", 180)],
                 60, 190, highlight=2, unit="k")
    dk.lint_layout(p, strict=True)                       # dots + labels stay in-canvas, no overlap
    # dense cluster near one end must not overflow the frame or collide (anti-collision path)
    p2 = dk.blank_deck(10, 5.625); s3 = dk.add_slide(p2)
    dk.dot_strip(s3, 0.6, 2.0, 8.0, [("a", 61), ("b", 63), ("c", 65), ("d", 188)], 60, 190, unit="k")
    dk.lint_layout(p2, strict=True)
ok("dot_strip (value-mapped + dense-cluster anti-collision, lint-clean)", _dot_strip)

# --- pangu: opt-in 盘古之白 normalizer — default OFF is byte-identical, modes are idempotent ---
def _pangu():
    s = "用K99机制占58.3%的博后"
    assert dk.pangu(s) == s, "pangu(default None) must be a no-op (byte-identical guarantee)"
    sp = dk.pangu(s, "spaced")
    assert "用 K99 机制" in sp and "58.3% 的" in sp, "pangu spaced did not insert boundary spaces"
    assert dk.pangu(sp, "spaced") == sp, "pangu spaced not idempotent"
    assert dk.pangu("用 K99 机制", "unspaced") == "用K99机制", "pangu unspaced did not strip"
    try:
        dk.pangu(s, "bogus"); raise AssertionError("pangu accepted an invalid mode")
    except ValueError:
        pass
ok("pangu (default no-op + spaced/unspaced idempotent + bad-mode guard)", _pangu)

# --- new geometry components: tier_stack / gantt / harvey_ball+eval_matrix / heat_matrix / device_frame ---
def _tier_stack():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.pyramid(s2, 1.0, 1.2, 4.4, 3.8, ["Vision", "Strategy", "Delivery", "Foundation"])
    dk.funnel(s2, 7.2, 1.2, 5.0, 3.8, ["Visits", "Signups", "Trials", "Paid"],
              values=[100, 54, 30, 12], labels="side")
    dk.lint_layout(p, strict=True)                       # tiers taper + labels stay collision-free
ok("tier_stack / pyramid + funnel (value-proportional, lint-clean)", _tier_stack)
raises("tier_stack rejects an empty tier list", lambda: dk.tier_stack(last(), 1, 1, 4, 3, []))

def _gantt():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.gantt(s2, 0.6, 1.3, 12.1,
             [("Discovery", 0, 1, 0), ("Design", 1, 2.5, 0),
              ("Build", 2, 5, 1), ("QA", 4.5, 6, 1), ("Launch", 6, 6.5, 2)],
             axis_min=0, axis_max=7, ticks=[0, 2, 4, 6], tick_labels=["Q1", "Q2", "Q3", "Q4"],
             lanes=["Plan", "Engineering", "GTM"], today=3.2, highlight=2)
    dk.lint_layout(p, strict=True)                       # bars keyed to axis_scale, swimlanes clean
ok("gantt (swimlanes + ticks + today + highlight, lint-clean)", _gantt)
raises("gantt raises on an off-axis bar", lambda: dk.gantt(last(), 0.6, 1.3, 11.0,
       [("X", 0, 5)], axis_min=0, axis_max=3))          # end 5 > axis_max 3 → off-canvas

def _harvey_ball():
    s2 = S()
    from pptx.oxml.ns import qn as _qn
    wsh = dk.harvey_ball(s2, 3.0, 3.0, 2, d=0.4)
    g = wsh._element.spPr.find(_qn("a:prstGeom"))
    assert g is not None and g.get("prst") == "pie", "harvey_ball(level=2) must build a PIE wedge"
    assert len(wsh.adjustments) == 2, "PIE must expose 2 adjustments (start/end angle)"
    dk.harvey_ball(s2, 4.0, 3.0, 0)                      # empty ring (no wedge)
    dk.harvey_ball(s2, 5.0, 3.0, 4)                      # full disc
ok("harvey_ball (PIE wedge geometry builds)", _harvey_ball)

def _eval_matrix():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.eval_matrix(s2, 0.8, 1.6, 11.0, ["Option A", "Option B", "Option C"],
                   ["Cost", "Speed", "Risk", "Support"],
                   [[4, 2, 3], [3, 4, 2], [2, 3, 4], [4, 1, 3]], recommend=0)
    dk.eval_matrix(dk.add_slide(p), 0.8, 1.6, 11.0, ["A", "B", "C"], ["c1", "c2", "c3"],
                   [["yes", "no", "partial"], ["partial", "yes", "no"], ["yes", "yes", "partial"]],
                   mark="mark", legend=False)
    dk.lint_layout(p, strict=True)
ok("eval_matrix (harvey balls + marks + recommend column, lint-clean)", _eval_matrix)

def _heat_matrix():
    p = dk.blank_deck(W, H); s2 = dk.add_slide(p)
    dk.heat_matrix(s2, 1.2, 1.4, 6.2, 4.2,
                   [[1, 2, 3, 4, 5], [2, 4, 6, 8, 10], [3, 6, 9, 12, 15],
                    [4, 8, 12, 16, 20], [5, 10, 15, 20, 25]],
                   ["Rare", "Unlikely", "Possible", "Likely", "Certain"],
                   ["Trivial", "Minor", "Moderate", "Major", "Severe"],
                   scale="risk", cell_labels=True)
    dk.heat_matrix(dk.add_slide(p), 1.2, 1.4, 6.0, 3.0,
                   [[10, 20, 30], [40, 50, 60]], ["r1", "r2"], ["c1", "c2", "c3"],
                   scale="div", cell_labels=True)
    dk.lint_layout(p, strict=True)                       # cells + contrast text stay collision-free
ok("heat_matrix (risk + div scales, contrast-aware text, lint-clean)", _heat_matrix)

def _device_frame():
    r = dk.device_frame(S(), IMG, 0.8, 1.3, 6.2, 3.6, chrome="browser",
                        url="app.example.com/dashboard")
    assert len(r) == 4 and r[2] > 0 and r[3] > 0, "device_frame must return the inner picture rect"
    dk.device_frame(S(), IMG, 5.2, 1.0, 2.2, 4.6, chrome="phone")
ok("device_frame (browser chrome + phone bezel)", _device_frame)

def _waterfall():
    q = os.path.join(TMP, "_wf.png")
    dc.waterfall(q, [("Start", None), ("Q1", 25), ("Q2", -12), ("Q3", 18), ("End", None)])
    assert os.path.exists(q) and os.path.getsize(q) > 1500, "waterfall PNG missing/trivial"
ok("dc.waterfall (floating step bars + dashed connectors)", _waterfall)
raises("waterfall rejects an empty item list", lambda: dc.waterfall(os.path.join(TMP, "_x.png"), []))

# --- box/connector kit: every node() shape + all three arrowhead variants ---
def _node_kit():
    s2 = S()
    centers = []
    for i, shp in enumerate(["roundrect", "rect", "pill", "circle", "diamond",
                             "parallelogram", "cylinder"]):
        centers.append(dk.node(s2, 0.4 + i * 1.8, 1.2, 1.6, 0.8, shp, shape=shp, sub="sub"))
    dk.node(s2, 0.4, 2.6, 1.6, 0.8, "hub", hub=True)
    dk.node(s2, 2.4, 2.6, 1.6, 0.8, "opt", dashed=True)
    dk.connector(s2, centers[0], centers[1], head="triangle", label="req")
    dk.connector(s2, (0.6, 4.0), (3.0, 4.0), style="dashed", head="open")
    dk.connector(s2, (0.6, 4.4), (3.0, 4.4), style="dotted", head="none")
    dk.elbow_connector(s2, dk.loop_path(8.0, 5.0, 4.0, 4.6), style="dotted", head="open", label="retry")
    dk.elbow_connector(s2, dk.loop_path(12.0, 9.0, 4.0, 4.6), head="none")
    dk.elbow_connector(s2, [(9.0, 5.2), (9.0, 5.6), (12.0, 5.6)], head="triangle")
ok("node (all 7 shapes + hub/dashed) + connector/elbow heads (triangle/open/none)", _node_kit)

raises("vstack overflow raises at build time", lambda: dk.vstack(
    S(), 0.6, 1.0, 5.0, [(2.0, lambda x, y, w: None), (2.0, lambda x, y, w: None)], bottom=3.0))

# --- content_band + bottom_callout: the footer-safe return geometry must be exact ---
def _band_and_callout():
    s2 = S()
    bx, by, bw, bh = dk.content_band(s2)
    assert abs(bx - dk.GUTTER) < 1e-6 and abs(bw - (W - 2 * dk.GUTTER)) < 1e-6, "content_band x/w wrong"
    assert abs((by + bh) - (H - dk.FOOTER_BAND - 0.15)) < 1e-6, "content_band bottom must clear the footer band"
    body = "the takeaway body — long enough that it could wrap on a narrower frame"
    top = dk.bottom_callout(s2, 0.6, W - 1.2, "TAKEAWAY", body)
    ch = dk.measure_callout("TAKEAWAY", body, W - 1.2)
    assert abs(top - (H - dk.FOOTER_BAND - 0.15 - ch)) < 1e-3, "bottom_callout bottom must anchor above the footer band"
    assert by < top < H, "bottom_callout top must fall inside the content band"
ok("content_band + bottom_callout (footer-safe return geometry)", _band_and_callout)

ok("equation_native (linear LaTeX -> editable runs)", lambda: dk.equation_native(
    S(), 0.8, 1.2, 9.0, 0.8, r"\mathcal{L} = \sum_i \|A x_i - y_i\|_2^2 + \lambda R(x_i)"))
raises("equation_native rejects 2-D \\frac (use equation_png)", lambda: dk.equation_native(
    S(), 0.8, 2.4, 6.0, 0.8, r"\frac{a}{b}"))

# --- step_list active_idx: the active disc's fill must DIFFER from an inactive one
#     (regression: the vertical branch once used _blend(acc, WHITE, 0.0) == acc, a no-op) ---
def _step_list_active():
    p = dk.blank_deck(10, 5.625); s2 = dk.add_slide(p)
    dk.step_list(s2, 0.8, 1.0, 8.0, [("Collect", "gather the inputs"), ("Train", "fit the model")],
                 active_idx=0)
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _T
    discs = [sh for sh in s2.shapes if sh.shape_type == _T.AUTO_SHAPE
             and abs(sh.width - dk.Inches(0.42)) < 2000 and abs(sh.height - dk.Inches(0.42)) < 2000]
    assert len(discs) == 2, f"expected 2 number discs, found {len(discs)}"
    fills = [d.fill.fore_color.rgb for d in discs]
    assert fills[0] != fills[1], "active_idx must render the active disc's fill differently from inactive"
    dk.step_list(s2, 0.8, 3.0, 8.0, [("A", ""), ("B", "")], orientation="horizontal", active_idx=1)
ok("step_list active vs inactive disc fills differ (v + h)", _step_list_active)

ok("cycle_diagram (3 nodes, defaults)", lambda: dk.cycle_diagram(
    S(), 6.5, 3.6, [("Plan", ""), ("Do", "daily"), ("Check", "")]))
ok("dumbbell_board (basic row, no hero/threshold)", lambda: dk.dumbbell_board(
    S(), 0.8, 1.6, 11.0, [("Latency", "", 120, 80, 0, 150, "ms")]))
ok("sources_page rule=True / rule=False", lambda: (dk.sources_page(S(), ["a", "b", "c"], rule=True),
                                                   dk.sources_page(S(), ["d", "e"], rule=False)))
ok("icon_tile (glass=True frosted tile)", lambda: dk.icon_tile(S(), 1.0, 1.0, 0.8, IMG, glass=True, fill=C("5B8DEF")))

prs.save(os.path.join(TMP, "_smoke_deck.pptx"))
print(f"\nsmoke_deckkit: {len(fails)} failure(s)" + ("" if not fails else " — " + "; ".join(n for n, _ in fails)))
sys.exit(1 if fails else 0)
