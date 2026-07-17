#!/usr/bin/env python3
"""Pull the CONTENT out of an existing .pptx so a redesign can reuse it.

When the user hands you their own deck to *improve*, you need its real material —
the actual text per slide, the tables, and the embedded figures — not a guess at
what's on it. This dumps all three so the rebuild is faithful to what they wrote and
reuses their figures whole (see design-principles: "use the source's own figures").

Pairs with render_deck.sh (which shows you what each slide LOOKS like) — run both:
render to see the design problems, extract to recover the content to carry forward.

Usage:  python3 extract_deck.py /path/to/their_deck.pptx [out_dir]
Output: <out_dir>/content.md         — per-slide text + tables + image filenames
        <out_dir>/assets/slideNN_*   — every embedded picture, saved whole
(default out_dir: ./extracted)
"""
import sys, os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

if len(sys.argv) < 2:
    print("usage: python3 extract_deck.py /path/to/their_deck.pptx [out_dir]")
    raise SystemExit(2)
path = sys.argv[1]
out = sys.argv[2] if len(sys.argv) > 2 else "./extracted"
assets = os.path.join(out, "assets")
os.makedirs(assets, exist_ok=True)
prs = Presentation(path)

W = Emu(prs.slide_width).inches
H = Emu(prs.slide_height).inches
lines = [f"# Extracted content — {os.path.basename(path)}",
         f"\nSlide size: {W:.2f} x {H:.2f} in · {len(prs.slides._sldIdLst)} slides\n"]
img_n = 0


def walk(shapes, slide_idx, texts, tables, imgs):
    """Recurse shapes (groups included), collecting text, tables, and pictures."""
    global img_n
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, slide_idx, texts, tables, imgs)
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = sh.image
                img_n += 1
                fn = f"slide{slide_idx:02d}_img{img_n:02d}.{img.ext}"
                with open(os.path.join(assets, fn), "wb") as f:
                    f.write(img.blob)
                w = Emu(sh.width).inches; h = Emu(sh.height).inches
                imgs.append(f"{fn}  ({w:.1f}x{h:.1f} in)")
            except Exception as e:
                imgs.append(f"<unreadable picture: {e}>")
            continue
        if sh.has_table:
            t = sh.table
            rows = [" | ".join(c.text.strip() for c in r.cells) for r in t.rows]
            tables.append("\n".join(rows))
            continue
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                texts.append(txt)


for i, slide in enumerate(prs.slides, 1):
    texts, tables, imgs = [], [], []
    walk(slide.shapes, i, texts, tables, imgs)
    lines.append(f"\n## Slide {i}")
    if texts:
        lines.append("**Text:**")
        for t in texts:
            lines.append("- " + t.replace("\n", "  \n  "))
    if tables:
        lines.append("\n**Table(s):**")
        for tb in tables:
            lines.append("```\n" + tb + "\n```")
    if imgs:
        lines.append("\n**Figures (saved to assets/, reuse these whole):**")
        for im in imgs:
            lines.append("- " + im)
    if not (texts or tables or imgs):
        lines.append("_(empty / decorative only)_")

with open(os.path.join(out, "content.md"), "w") as f:
    f.write("\n".join(lines) + "\n")
print(f"extracted {len(prs.slides._sldIdLst)} slides -> {out}/content.md  "
      f"({img_n} images -> {assets}/)")
