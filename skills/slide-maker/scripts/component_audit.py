#!/usr/bin/env python3
"""Component audit — did this deck hand-roll a form the library already implements?

WHY THIS EXISTS, measured: across three delivered decks (45 slides) the build scripts reached for
**3 of deckkit's 59 form components**. Everything else was composed from raw `box` + `text`. That is
not a style complaint — a hand-rolled form re-introduces the exact geometry bugs the component was
written to fix (a baseline hardcoded short of the last bar; a value label floating off the bar's
centerline; a segmented band whose parts do not sum; a timeline whose end caption drifts off its dot).
SKILL.md has said "when a COMPONENT exists for the form, BUILD that component" since forever — as
PROSE, with no gate. It was violated dozens of times and detected zero times.

WHY IT IS ADVISORY, NOT A BLOCKER — and this is a deliberate design decision, not timidity:
geometry alone CANNOT distinguish a lazy hand-roll from a deliberate bespoke composition, and the
deliberate one is the *signature move* — the single thing that makes a deck not-a-template. A gate
that blocked on "you hand-rolled a bar row" would delete the best page of the deck that motivated
it. So this tool reports and directs attention; it never vetoes. What it CAN state as fact is the
usage ratio and the specific component whose geometry the cluster matches.

CLI:  python component_audit.py <build_script.py> [<deck.pptx>] [--json]
Exit: 0 = nothing to report · 2 = at least one cluster matches an unused component (advisory).
"""
import argparse
import json
import os
import re
import sys

# form components + the concrete guarantee each one gives that a hand-roll does not
FORM_GUARANTEE = {
    "native_chart": "a real editable chart; axis derived from the data, non-Latin labels render",
    "segmented_bar": "parts are normalised to sum, and each segment's label is fitted to its width",
    "meter_bar": "the value label is centered ON the bar's centerline by construction",
    "stat_row": "even column pitch and one baseline for every figure",
    "scorecard": "equal tiles with one hero, from the region — not a hand-picked pitch",
    "timeline": "end markers are inset so the first/last caption stays co-centred with its dot",
    "step_list": "numbered steps with a derived pitch; no stranded gap under the header",
    "tier_stack": "one taper (funnel/pyramid) sized from the values, not clamped by a min-size floor",
    "waterfall": "increments or a total — never both as peer bars (the double-count bug)",
    "dot_strip": "even spacing with the end dots inset",
    "eval_matrix": "options x criteria with fitted glyph cells",
    "heat_matrix": "one colour scale bound to the value range",
    "leaderboard": "ranked rows keyed to a chart, with one highlight",
    "table": "header/rule/highlight styling and column widths that respect the content",
    "org_tree": "centroid parents and a horizontal bus; raises when it cannot fit legibly",
    "position_map": "labelled 2-D positions with anti-collision labels",
    "small_multiples": "every panel pinned to ONE shared value axis",
}


def _script_calls(path):
    src = open(path, encoding="utf-8", errors="ignore").read()
    return {m.group(1) for m in re.finditer(r"\bdk\.([a-z][a-z0-9_]*)\s*\(", src)}


def _shapes(pptx_path):
    from pptx import Presentation
    from pptx.util import Emu
    out = []
    prs = Presentation(pptx_path)
    for n, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            try:
                x, y = Emu(sh.left).inches, Emu(sh.top).inches
                w, h = Emu(sh.width).inches, Emu(sh.height).inches
            except Exception:
                continue
            txt = ""
            try:
                if sh.has_text_frame:
                    txt = sh.text_frame.text.strip()
            except Exception:
                pass
            kind = "chart" if getattr(sh, "has_chart", False) else (
                "table" if getattr(sh, "has_table", False) else (
                    "pic" if "PICTURE" in str(getattr(sh, "shape_type", "")) else
                    ("text" if txt else "rect")))
            out.append({"slide": n, "x": x, "y": y, "w": w, "h": h, "kind": kind, "txt": txt})
    return out


def _clusters(shapes):
    """Geometry signatures that a deckkit component already implements."""
    hits = []
    by_slide = {}
    for s in shapes:
        by_slide.setdefault(s["slide"], []).append(s)

    for n, shs in by_slide.items():
        rects = [s for s in shs if s["kind"] == "rect" and s["w"] > 0.12 and s["h"] > 0.06]
        texts = [s for s in shs if s["kind"] == "text"]

        def near_text(r, pad=0.55):
            return [t for t in texts
                    if abs(t["y"] - r["y"]) < pad and t["x"] > r["x"] - 2.6 and t["x"] < r["x"] + r["w"] + 2.6]

        # BAR ROW — same x0 and h, different w, stacked vertically, with digit-bearing labels
        rows = {}
        for r in rects:
            rows.setdefault((round(r["x"], 2), round(r["h"], 2)), []).append(r)
        for (x0, h), group in rows.items():
            if len(group) < 3:
                continue
            if len({round(g["w"], 2) for g in group}) < 2:
                continue                       # all identical width -> that is a card stack, below
            if any(re.search(r"\d", t["txt"]) for g in group for t in near_text(g)):
                hits.append({"slide": n, "pattern": "bar row",
                             "detail": f"{len(group)} rects sharing x={x0}in and h={h}in with varying "
                                       f"width and numeric labels",
                             "components": ["native_chart", "meter_bar", "segmented_bar",
                                            "bullet_graph", "range_bars", "dumbbell_board",
                                            "dot_strip", "waterfall", "leaderboard"]})

        # 100% BAND — >=3 rects on one y, abutting (gap ~0), spanning a wide run
        bands = {}
        for r in rects:
            bands.setdefault((round(r["y"], 2), round(r["h"], 2)), []).append(r)
        for (y0, h), group in bands.items():
            if len(group) < 3:
                continue
            g = sorted(group, key=lambda r: r["x"])
            gaps = [g[i + 1]["x"] - (g[i]["x"] + g[i]["w"]) for i in range(len(g) - 1)]
            span = (g[-1]["x"] + g[-1]["w"]) - g[0]["x"]
            if span > 3.0 and all(abs(v) < 0.06 for v in gaps):
                hits.append({"slide": n, "pattern": "abutting 100% band",
                             "detail": f"{len(group)} rects abutting across {span:.1f}in at y={y0}in",
                             "components": ["segmented_bar", "tier_stack", "native_chart",
                                            "org_tree", "heat_matrix", "eval_matrix"]})

        # CARD/TILE ROW — >=3 identical rects evenly spaced on one axis
        for (y0, h), group in bands.items():
            if len(group) < 3:
                continue
            g = sorted(group, key=lambda r: r["x"])
            if len({round(r["w"], 2) for r in g}) != 1:
                continue
            pitches = [g[i + 1]["x"] - g[i]["x"] for i in range(len(g) - 1)]
            if pitches and max(pitches) - min(pitches) < 0.05 and min(pitches) > g[0]["w"]:
                hits.append({"slide": n, "pattern": "tile row",
                             "detail": f"{len(group)} identical rects on an even pitch at y={y0}in",
                             "components": ["scorecard", "stat_row", "columns", "org_tree",
                                            "step_list", "timeline", "tier_stack",
                                            "position_map", "eval_matrix", "heat_matrix"]})

        # MARKER ROW — >=3 small same-size shapes evenly spaced with captions under them
        dots = [r for r in rects if r["w"] < 0.4 and r["h"] < 0.4]
        if len(dots) >= 3:
            g = sorted(dots, key=lambda r: r["x"])
            ys = {round(r["y"], 1) for r in g}
            pitches = [g[i + 1]["x"] - g[i]["x"] for i in range(len(g) - 1)]
            if len(ys) == 1 and pitches and max(pitches) - min(pitches) < 0.08:
                caps = [t for t in texts if any(abs(t["x"] + t["w"] / 2 - (r["x"] + r["w"] / 2)) < 1.2
                                                and t["y"] > r["y"] for r in g)]
                if len(caps) >= 3:
                    hits.append({"slide": n, "pattern": "marker row with captions",
                                 "detail": f"{len(g)} evenly spaced markers with captions below",
                                 "components": ["timeline", "spaced_centers", "step_list",
                                                "dot_strip", "position_map", "org_tree"]})
    return hits


def audit(script_path, pptx_path=None):
    called = _script_calls(script_path)
    used_forms = sorted(f for f in FORM_GUARANTEE if f in called)
    hits = []
    if pptx_path and os.path.isfile(pptx_path):
        try:
            hits = _clusters(_shapes(pptx_path))
        except Exception as e:                                   # noqa: BLE001
            hits = [{"slide": 0, "pattern": "could not read the deck", "detail": str(e),
                     "components": []}]
    # Only report a cluster when NO component that could have EMITTED that geometry was called
    # anywhere in the script. This is what stops a component's own output being reported as a
    # hand-roll (org_tree draws a row of identical node rects; in the finished pptx that is
    # indistinguishable from three hand-placed boxes). Keep every emitting component in each
    # pattern's list — a short list here means false positives, not sharper detection.
    # a deck that uses meter_bar elsewhere and hand-rolls one bar has clearly made a choice.
    actionable = [h for h in hits if h["components"] and not (set(h["components"]) & called)]
    return {"used_forms": used_forms, "form_total": len(FORM_GUARANTEE),
            "clusters": hits, "actionable": actionable}


def main():
    ap = argparse.ArgumentParser(description="did this deck hand-roll a form the library implements?")
    ap.add_argument("script")
    ap.add_argument("pptx", nargs="?")
    ap.add_argument("--json", action="store_true", dest="as_json")
    a = ap.parse_args()
    pptx = a.pptx
    if pptx is None:
        d = os.path.dirname(os.path.abspath(a.script))
        cand = [f for f in os.listdir(d) if f.endswith(".pptx") and not f.startswith("~$")]
        pptx = os.path.join(d, cand[0]) if len(cand) == 1 else None
    r = audit(a.script, pptx)
    if a.as_json:
        print(json.dumps(r, indent=1)); sys.exit(2 if r["actionable"] else 0)
    print("[components] this deck calls {} of {} form components: {}".format(
        len(r["used_forms"]), r["form_total"], ", ".join(r["used_forms"]) or "none"))
    if not r["actionable"]:
        print("[components] no hand-rolled cluster matches an unused component.")
        sys.exit(0)
    print("[components] {} cluster(s) look like a form the library already implements. "
          "This is ADVISORY:".format(len(r["actionable"])))
    print("             a bespoke composition IS the signature move — but a hand-rolled COMMON form")
    print("             re-inherits the geometry bugs the component fixed. Confirm each is deliberate.")
    for h in r["actionable"]:
        print(f"  slide {h['slide']:>2}  {h['pattern']}: {h['detail']}")
        for c in h["components"]:
            g = FORM_GUARANTEE.get(c)
            print(f"            deckkit.{c}() — {g}" if g else f"            deckkit.{c}()")
    sys.exit(2)


if __name__ == "__main__":
    main()
